"""
Style analysis and profiling tools for PowerPoint MCP Server.

Provides MCP tools to analyze presentation styles, create reusable
style profiles, and apply them to new presentations.
"""

from typing import Any, Dict, List, Optional
from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations

# Style profiles stored in memory (keyed by profile name)
_style_profiles: Dict[str, Any] = {}


def _apply_legacy_flat_profile(pres, profile, profile_name: str,
                               pres_id: str) -> Dict:
    """Apply a legacy flat profile (create_style_profile output).

    Historical behavior, unchanged: every run gets the profile's
    primary font; runs in title-positioned shapes get the title size,
    all others the body size.
    """
    from pptx.util import Pt

    modified_runs = 0
    for slide in pres.slides:
        for shape in slide.shapes:
            if not (hasattr(shape, "text_frame") and shape.text_frame):
                continue

            top = shape.top.inches if shape.top else 5
            is_title = top < 2 and len(shape.text_frame.text) < 100

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = profile.primary_font
                    if is_title:
                        run.font.size = Pt(profile.title_font_size)
                    else:
                        run.font.size = Pt(profile.body_font_size)
                    modified_runs += 1

    return {
        "message": f"Applied profile '{profile_name}' to presentation '{pres_id}'",
        "font_applied": profile.primary_font,
        "title_size": profile.title_font_size,
        "body_size": profile.body_font_size,
        "runs_modified": modified_runs,
    }


def register_style_tools(app: FastMCP, presentations: Dict, get_current_presentation_id):
    """Register style analysis and profiling tools with the FastMCP app."""

    @app.tool(
        annotations=ToolAnnotations(
            title="Analyze Presentation Style",
        ),
    )
    def analyze_presentation_style(
        file_path: str,
        resolved: bool = True,
        slide_range: Optional[str] = None,
        detail: str = "summary",
        groups: Optional[List[str]] = None,
    ) -> Dict:
        """Analyze a PowerPoint file to extract fonts, colors, layouts, and text hierarchy.

        Returns a comprehensive style analysis including:
        - Primary font and font usage patterns
        - Color palette with usage contexts
        - Layout positioning patterns
        - Text hierarchy (title/subtitle/body styles)
        - Consistency score (0-1)
        - A "resolved" section (default on) with inheritance-resolved
          EFFECTIVE values -- what PowerPoint actually displays -- for
          fonts, colors, paragraph spacing, bullets, theme schemes, and
          shape line/fill, resolved through the placeholder -> layout ->
          master -> theme chain even where the XML stores nothing

        The resolved section is hard-capped (~40KB); when truncated it
        carries {"truncated": true, "hint": ...} -- narrow with
        slide_range / detail / groups to see more.

        Args:
            file_path: Path to the .pptx file to analyze
            resolved: Include the inheritance-resolved section (default
                True; set False for the raw python-pptx view only)
            slide_range: 1-based slides for the resolved section, e.g.
                "1-3,5" (default: all slides)
            detail: "summary" (deck-level rollup, default) or "full"
                (per-run dump) for the resolved section
            groups: Property-group filter for the resolved section, any
                subset of ["fonts", "colors", "paragraphs", "shapes",
                "theme"] (default: all)
        """
        from pathlib import Path

        from pptx import Presentation

        from utils.resolve_analysis import build_resolved_analysis
        from utils.style_utils import analyze_open_presentation

        try:
            if not Path(file_path).exists():
                raise FileNotFoundError(
                    f"Presentation not found: {file_path}"
                )
            # Parse the .pptx once; the raw analysis and the resolved
            # section share the same Presentation object. The resolver
            # runs FIRST: it is strictly read-only, while python-pptx
            # property reads in the legacy analyzer (e.g.
            # ``run.font.color``) mutate the in-memory tree.
            prs = Presentation(file_path)
            resolved_report = (
                build_resolved_analysis(
                    prs,
                    slide_range=slide_range,
                    detail=detail,
                    groups=groups,
                )
                if resolved else None
            )
            analysis = analyze_open_presentation(prs, file_path)
            result = {
                "message": f"Analyzed {analysis['slide_count']} slides from {file_path}",
                "primary_font": analysis["fonts"]["primary_font"],
                "font_count": len(analysis["fonts"]["font_usage"]),
                "color_count": analysis["colors"]["total_unique_colors"],
                "top_colors": analysis["colors"]["primary_palette"][:5],
                "common_sizes": analysis["fonts"]["common_sizes"][:5],
                "text_hierarchy": analysis.get("text_hierarchy", {}),
                "consistency_score": analysis["consistency_score"],
                "slide_dimensions": f"{analysis['slide_width']}x{analysis['slide_height']} inches",
                "total_shapes": analysis["shapes"]["total_shapes"],
                "full_analysis": analysis,
            }
            if resolved_report is not None:
                result["resolved"] = resolved_report
            return result
        except FileNotFoundError as e:
            return {"error": str(e)}
        except ValueError as e:
            return {"error": f"Invalid analysis options: {str(e)}"}
        except Exception as e:
            return {"error": f"Analysis failed: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Create Style Profile",
        ),
    )
    def create_style_profile(file_path: str, profile_name: str) -> Dict:
        """Analyze a presentation and create a named, reusable style profile.

        The profile captures fonts, colors, layout patterns, and text hierarchy
        that can later be applied to new presentations or referenced for
        consistent formatting.

        Args:
            file_path: Path to the .pptx file to profile
            profile_name: Name for the style profile (e.g., 'firm_standard', 'client_acme')
        """
        from utils.style_utils import analyze_presentation, create_profile

        try:
            analysis = analyze_presentation(file_path)
            profile = create_profile(analysis, profile_name)
            _style_profiles[profile_name] = profile

            return {
                "message": f"Created style profile '{profile_name}' from {file_path}",
                "profile_name": profile_name,
                "primary_font": profile.primary_font,
                "title_font_size": profile.title_font_size,
                "body_font_size": profile.body_font_size,
                "color_count": len(profile.color_palette),
                "consistency_score": profile.consistency_score,
            }
        except FileNotFoundError as e:
            return {"error": str(e)}
        except Exception as e:
            return {"error": f"Failed to create profile: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Create House Style Profile",
        ),
    )
    def create_house_profile(paths: List[str], profile_name: str) -> Dict:
        """Learn a compact house-style profile from multiple reference decks.

        Runs the inheritance-resolved analyzer across 5-10 reference
        presentations and aggregates the EFFECTIVE values into a
        prescriptive "house-profile/1" rules JSON (hard 8KB budget):

        - typography/paragraph/palette/shape_defaults: modal house rules
          (fonts, sizes, colors, spacing, bullets, borders) -- the
          deterministic subset apply_style_profile consumes
        - grid: the inferred alignment-column grid (left/right/center
          edge positions in inches + tolerance) -- consult when placing
          shapes
        - archetypes: learned slide-layout types with title_band /
          body_region boxes -- consult when composing slides
        - images/distributions: placement zones and house scales
          (font-size scale, spacing quanta, palette shares) for
          analysis and lint

        The profile is stored in memory under profile_name; persist it
        with save_style_profile and reload with load_style_profile.

        Args:
            paths: Paths of the reference .pptx decks (all must share
                one slide size)
            profile_name: Name to store the profile under
        """
        from utils.profile_extract import create_house_profile as build
        from utils.profile_schema import enforce_size_budget

        try:
            profile = build(paths, profile_name)
            _style_profiles[profile_name] = profile
            payload = enforce_size_budget(profile)
            return {
                "message": (
                    f"Created house profile '{profile_name}' from "
                    f"{len(paths)} deck(s)"
                ),
                "profile_name": profile_name,
                "schema_version": profile["schema_version"],
                "source_decks": profile["source_decks"],
                "profile_bytes": len(payload),
                "archetypes": {
                    name: spec["count"]
                    for name, spec in profile["archetypes"].items()
                },
                "grid_edges": profile["grid"]["edges"],
                "profile": profile,
            }
        except FileNotFoundError as e:
            return {"error": str(e)}
        except ValueError as e:
            return {"error": f"Invalid house-profile input: {str(e)}"}
        except Exception as e:
            return {"error": f"Failed to create house profile: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Save Style Profile",
        ),
    )
    def save_style_profile(profile_name: str, output_path: str) -> Dict:
        """Save a style profile to a JSON file for reuse across sessions.

        Args:
            profile_name: Name of a previously created profile
            output_path: Path where the JSON profile will be saved
        """
        from utils.style_utils import save_profile

        if profile_name not in _style_profiles:
            available = list(_style_profiles.keys()) or ["(none)"]
            return {
                "error": f"Profile '{profile_name}' not found. Available: {', '.join(available)}"
            }

        try:
            profile = _style_profiles[profile_name]
            save_profile(profile, output_path)
            return {
                "message": f"Saved profile '{profile_name}' to {output_path}",
                "profile_name": profile_name,
                "output_path": output_path,
            }
        except Exception as e:
            return {"error": f"Failed to save profile: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Load Style Profile",
        ),
    )
    def load_style_profile(file_path: str) -> Dict:
        """Load a previously saved style profile from a JSON file.

        Accepts both legacy flat profiles (created by
        create_style_profile) and house-profile/1 multi-deck profiles
        (schema_version "house-profile/1"). The loaded profile is
        available for reference and can be applied with
        apply_style_profile.

        Args:
            file_path: Path to the JSON profile file
        """
        import json
        from pathlib import Path

        from utils.style_apply import is_house_profile
        from utils.style_utils import load_profile

        try:
            data = json.loads(Path(file_path).read_text(encoding="utf-8"))
            if is_house_profile(data):
                name = data.get("name")
                if not name or not isinstance(name, str):
                    return {
                        "error": "house-profile/1 file is missing its "
                                 "'name' key"
                    }
                _style_profiles[name] = data
                return {
                    "message": f"Loaded house profile '{name}' from "
                               f"{file_path}",
                    "profile_name": name,
                    "schema_version": data["schema_version"],
                    "source_decks": data.get("source_decks", []),
                    "sections": sorted(
                        key for key in data
                        if key not in ("schema_version", "name",
                                       "source_decks")
                    ),
                }
            profile = load_profile(file_path)
            _style_profiles[profile.name] = profile
            return {
                "message": f"Loaded profile '{profile.name}' from {file_path}",
                "profile_name": profile.name,
                "primary_font": profile.primary_font,
                "title_font_size": profile.title_font_size,
                "body_font_size": profile.body_font_size,
                "color_palette_count": len(profile.color_palette),
                "consistency_score": profile.consistency_score,
            }
        except FileNotFoundError:
            return {"error": f"Profile file not found: {file_path}"}
        except Exception as e:
            return {"error": f"Failed to load profile: {str(e)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Get Style Profile",
        ),
    )
    def get_style_profile(profile_name: str) -> Dict:
        """Get the full details of a loaded style profile.

        Returns all extracted style information including font hierarchy,
        color palette, and layout patterns.

        Args:
            profile_name: Name of a loaded profile
        """
        from dataclasses import asdict

        if profile_name not in _style_profiles:
            available = list(_style_profiles.keys()) or ["(none)"]
            return {
                "error": f"Profile '{profile_name}' not found. Available: {', '.join(available)}"
            }

        profile = _style_profiles[profile_name]
        return {
            "profile": asdict(profile) if hasattr(profile, "__dataclass_fields__") else profile,
            "profile_name": profile_name,
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="List Style Profiles",
        ),
    )
    def list_style_profiles() -> Dict:
        """List all style profiles currently loaded in memory."""
        profiles = []
        for name, profile in _style_profiles.items():
            if hasattr(profile, "primary_font"):
                profiles.append({
                    "name": name,
                    "primary_font": profile.primary_font,
                    "source": profile.source_file,
                    "consistency_score": profile.consistency_score,
                })
            else:
                profiles.append({"name": name, "type": "dict"})

        return {
            "profiles": profiles,
            "total": len(profiles),
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Apply Style to Presentation",
        ),
    )
    def apply_style_profile(
        profile_name: str,
        presentation_id: Optional[str] = None,
    ) -> Dict:
        """Apply a style profile's formatting rules to the current presentation.

        Two profile kinds are supported:

        - house-profile/1 (multi-deck house profiles): applies the
          DETERMINISTIC sections — typography (title/body/footer font,
          size, bold, color, mapped via placeholder type), paragraph
          spacing and per-level bullets, palette re-linking (hardcoded
          srgbClr equal to a scheme color becomes the schemeClr token),
          and shape_defaults (border weight/color/dash, hardcoded
          off-palette srgbClr fills — theme-linked/tinted fills are
          never repainted — and corner radius). Writes are
          MINIMAL-DIFF: a value is
          rewritten only where the inheritance-resolved EFFECTIVE value
          deviates from the rule, so conformant content gets no new
          explicit overrides. Slide masters, layouts, themes and shape
          GEOMETRY are never touched — grid/archetypes/distributions/
          images sections are consumed at generation time and enforced
          by lint, not applied here.
        - legacy flat profiles (create_style_profile): updates all text
          to the profile's primary font, title size, and body size.

        Does not change content — only formatting.

        Args:
            profile_name: Name of a loaded style profile
            presentation_id: ID of presentation to style (default: current)
        """
        from utils.style_apply import apply_house_profile, is_house_profile

        if profile_name not in _style_profiles:
            return {"error": f"Profile '{profile_name}' not found"}

        profile = _style_profiles[profile_name]
        pres_id = presentation_id or get_current_presentation_id()
        if not pres_id or pres_id not in presentations:
            return {"error": "No presentation loaded. Create or open one first."}

        pres = presentations[pres_id]

        if is_house_profile(profile):
            try:
                summary = apply_house_profile(pres, profile)
            except ValueError as e:
                return {"error": f"House-profile apply rejected invalid "
                                 f"input: {str(e)}"}
            except Exception as e:
                return {"error": f"House-profile apply failed: {str(e)}"}
            return {
                "message": (
                    f"Applied house profile '{profile_name}' to "
                    f"presentation '{pres_id}': {summary['writes']} "
                    f"deviation(s) corrected on slides "
                    f"{summary['slides_touched']}"
                ),
                **summary,
            }

        return _apply_legacy_flat_profile(pres, profile, profile_name,
                                          pres_id)
