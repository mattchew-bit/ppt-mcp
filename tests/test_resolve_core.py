"""Tests for utils.resolve_core -- cascade walker + placeholder matching.

Two layers:
  * synthetic-XML unit tests for the walker mechanics and the ported
    SCPShapeTree placeholder-matching semantics (type+idx, fallback
    type-only), and
  * fixture-driven integration tests that walk the real inheritance
    chains of the PowerPoint-authored decks and assert EXACT values from
    the COM-extracted ``expected_values`` JSON (never just non-None).
"""

import pytest
from lxml import etree
from pptx.oxml.ns import qn

from tests.conftest import load_expected, skip_if_fixture_missing
from utils.resolve_core import (
    KIND_LIST_STYLE,
    KIND_PARA_PROPS,
    KIND_RUN_PROPS,
    StyleSource,
    build_text_cascade_sources,
    indent_level_of,
    master_text_style_for_ph_type,
    match_placeholder_in_tree,
    placeholder_of,
    resolve_first,
    resolve_paragraph_property,
    resolve_run_property,
)
from utils.resolve_theme import ThemeContext

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def a_element(fragment: str) -> etree._Element:
    return etree.fromstring(f'<root xmlns:a="{_A_NS}">{fragment}</root>')[0]


def p_element(fragment: str) -> etree._Element:
    return etree.fromstring(
        f'<root xmlns:a="{_A_NS}" xmlns:p="{_P_NS}">{fragment}</root>'
    )[0]


def sp_with_ph(ph_attrs: str) -> str:
    """A minimal placeholder ``p:sp`` fragment for matching tests."""
    return (
        "<p:sp><p:nvSpPr><p:cNvPr id='2' name='x'/><p:cNvSpPr/>"
        f"<p:nvPr><p:ph {ph_attrs}/></p:nvPr></p:nvSpPr></p:sp>"
    )


# -- shared extractors -------------------------------------------------------

def extract_sz(rpr_element):
    value = rpr_element.get("sz")
    return None if value is None else int(value)


def extract_bold(rpr_element):
    return rpr_element.get("b")


def extract_italic(rpr_element):
    return rpr_element.get("i")


def extract_latin(rpr_element):
    latin = rpr_element.find(qn("a:latin"))
    return None if latin is None else latin.get("typeface")


def extract_solid_fill(rpr_element):
    return rpr_element.find(qn("a:solidFill"))


# ---------------------------------------------------------------------------
# StyleSource mechanics
# ---------------------------------------------------------------------------

class TestStyleSource:
    def test_list_style_level_lookup(self):
        style = a_element(
            '<a:lstStyle>'
            '<a:lvl1pPr><a:defRPr sz="2400"/></a:lvl1pPr>'
            '<a:lvl2pPr><a:defRPr sz="2000"/></a:lvl2pPr>'
            '</a:lstStyle>'
        )
        source = StyleSource("s", KIND_LIST_STYLE, style)
        assert source.run_props(1).get("sz") == "2400"
        assert source.run_props(2).get("sz") == "2000"
        assert source.run_props(3) is None

    def test_list_style_defppr_fallback(self):
        style = a_element(
            "<a:lstStyle><a:defPPr><a:defRPr sz='1800'/></a:defPPr></a:lstStyle>"
        )
        source = StyleSource("s", KIND_LIST_STYLE, style)
        assert source.run_props(4).get("sz") == "1800"

    def test_run_props_kind_ignores_level(self):
        rpr = a_element('<a:rPr sz="1200"/>')
        source = StyleSource("run", KIND_RUN_PROPS, rpr)
        assert source.run_props(1) is rpr
        assert source.run_props(9) is rpr
        assert source.paragraph_props(1) is None

    def test_para_props_kind(self):
        ppr = a_element('<a:pPr algn="ctr"><a:defRPr sz="900"/></a:pPr>')
        source = StyleSource("para", KIND_PARA_PROPS, ppr)
        assert source.paragraph_props(3) is ppr
        assert source.run_props(3).get("sz") == "900"

    def test_invalid_kind_rejected(self):
        with pytest.raises(ValueError):
            StyleSource("bad", "notAKind", a_element("<a:lstStyle/>"))

    def test_missing_element_rejected(self):
        with pytest.raises(ValueError):
            StyleSource("bad", KIND_LIST_STYLE, None)

    @pytest.mark.parametrize("level", [0, 10, "1"])
    def test_out_of_range_level_rejected(self, level):
        source = StyleSource("s", KIND_LIST_STYLE, a_element("<a:lstStyle/>"))
        with pytest.raises(ValueError):
            source.run_props(level)


# ---------------------------------------------------------------------------
# The generic walker
# ---------------------------------------------------------------------------

class TestWalker:
    def _sources(self):
        first = a_element("<a:lstStyle><a:lvl1pPr><a:defRPr b='1'/></a:lvl1pPr></a:lstStyle>")
        second = a_element(
            "<a:lstStyle><a:lvl1pPr><a:defRPr sz='2000' b='0'/></a:lvl1pPr></a:lstStyle>"
        )
        return (
            StyleSource("first", KIND_LIST_STYLE, first),
            StyleSource("second", KIND_LIST_STYLE, second),
        )

    def test_first_explicit_value_wins(self):
        assert resolve_run_property(self._sources(), 1, extract_bold) == "1"

    def test_walk_continues_past_silent_sources(self):
        assert resolve_run_property(self._sources(), 1, extract_sz) == 2000

    def test_returns_none_when_no_source_defines(self):
        assert resolve_run_property(self._sources(), 1, extract_latin) is None

    def test_generic_resolve_first(self):
        assert resolve_first(self._sources(), lambda source: source.name) == "first"

    def test_paragraph_property_resolution(self):
        style = a_element(
            "<a:lstStyle><a:lvl2pPr algn='r'/></a:lstStyle>"
        )
        sources = (StyleSource("s", KIND_LIST_STYLE, style),)
        assert resolve_paragraph_property(
            sources, 2, lambda ppr: ppr.get("algn")
        ) == "r"
        assert resolve_paragraph_property(
            sources, 1, lambda ppr: ppr.get("algn")
        ) is None

    def test_extractor_required(self):
        with pytest.raises(ValueError):
            resolve_first((), None)


# ---------------------------------------------------------------------------
# Placeholder matching (SCPShapeTree semantics)
# ---------------------------------------------------------------------------

class TestPlaceholderMatching:
    def test_type_and_idx_beats_type_only(self):
        tree = p_element(
            "<p:spTree>"
            + sp_with_ph("type='body' idx='2'")
            + sp_with_ph("type='body' idx='1'")
            + "</p:spTree>"
        )
        target = p_element("<p:ph type='body' idx='1'/>")
        matched = match_placeholder_in_tree(target, tree)
        assert placeholder_of(matched).get("idx") == "1"

    def test_fallback_to_type_only(self):
        tree = p_element(
            "<p:spTree>" + sp_with_ph("type='body' idx='1'") + "</p:spTree>"
        )
        target = p_element("<p:ph type='body' idx='5'/>")
        matched = match_placeholder_in_tree(target, tree)
        assert placeholder_of(matched).get("type") == "body"

    def test_absent_attributes_match_absent(self):
        tree = p_element("<p:spTree>" + sp_with_ph("type='title'") + "</p:spTree>")
        target = p_element("<p:ph type='title'/>")
        assert match_placeholder_in_tree(target, tree) is not None

    def test_untyped_target_gets_no_fallback(self):
        tree = p_element(
            "<p:spTree>" + sp_with_ph("type='body' idx='1'") + "</p:spTree>"
        )
        target = p_element("<p:ph idx='3'/>")  # absent type
        assert match_placeholder_in_tree(target, tree) is None

    def test_no_placeholders_returns_none(self):
        tree = p_element(
            "<p:spTree><p:sp><p:nvSpPr><p:cNvPr id='2' name='x'/>"
            "<p:cNvSpPr/><p:nvPr/></p:nvSpPr></p:sp></p:spTree>"
        )
        target = p_element("<p:ph type='body'/>")
        assert match_placeholder_in_tree(target, tree) is None

    def test_ctr_title_does_not_alias_to_title(self):
        # Documented SCPShapeTree semantics: no title/ctrTitle aliasing
        # at the matching layer -- master txStyles covers it instead.
        tree = p_element("<p:spTree>" + sp_with_ph("type='title'") + "</p:spTree>")
        target = p_element("<p:ph type='ctrTitle'/>")
        assert match_placeholder_in_tree(target, tree) is None

    def test_validates_inputs(self):
        tree = p_element("<p:spTree/>")
        with pytest.raises(ValueError):
            match_placeholder_in_tree(None, tree)
        with pytest.raises(ValueError):
            match_placeholder_in_tree(p_element("<p:ph/>"), None)
        with pytest.raises(ValueError):
            placeholder_of(None)


# ---------------------------------------------------------------------------
# Master txStyles selection
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("theme_only.pptx")
class TestMasterTextStyleSelection:
    @pytest.fixture
    def master_element(self, open_deck):
        prs = open_deck("theme_only.pptx")
        return prs.slides[0].slide_layout.slide_master.element

    @pytest.mark.parametrize("ph_type,expected_tag", [
        ("title", "p:titleStyle"),
        ("ctrTitle", "p:titleStyle"),
        ("body", "p:bodyStyle"),
        ("subTitle", "p:bodyStyle"),
        ("obj", "p:bodyStyle"),
        (None, "p:bodyStyle"),
        ("ftr", "p:otherStyle"),
        ("sldNum", "p:otherStyle"),
    ])
    def test_mapping(self, master_element, ph_type, expected_tag):
        style = master_text_style_for_ph_type(master_element, ph_type)
        assert style is not None
        assert style.tag == qn(expected_tag)

    def test_none_master_rejected(self):
        with pytest.raises(ValueError):
            master_text_style_for_ph_type(None, "title")


# ---------------------------------------------------------------------------
# Fixture integration: the walker reproduces COM-recorded values exactly
# ---------------------------------------------------------------------------

def _run_font_expectation(fixture_name, slide_number, shape_index):
    expected = load_expected(fixture_name)
    slide_entry = next(
        entry for entry in expected["slides"] if entry["index"] == slide_number
    )
    return slide_entry["shapes"][shape_index]["paragraphs"][0]["runs"][0]["font"]


def _resolved_run_font(prs, slide_index, shape_name):
    """Resolve size/font/bold/italic/color for a shape's first run."""
    slide = prs.slides[slide_index]
    shape = next(s for s in slide.shapes if s.name == shape_name)
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    level = indent_level_of(paragraph)
    sources = build_text_cascade_sources(slide, shape, paragraph, run)
    context = ThemeContext.for_slide(slide)

    hundredths = resolve_run_property(sources, level, extract_sz)
    typeface = resolve_run_property(sources, level, extract_latin)
    fill = resolve_run_property(sources, level, extract_solid_fill)
    return {
        "size_pt": None if hundredths is None else hundredths / 100.0,
        "name": None if typeface is None else context.resolve_typeface(typeface),
        "bold": resolve_run_property(sources, level, extract_bold) == "1",
        "italic": resolve_run_property(sources, level, extract_italic) == "1",
        "color_rgb": None if fill is None else context.resolve_solid_fill(fill),
    }


@skip_if_fixture_missing("theme_only.pptx")
class TestThemeOnlyIntegration:
    def test_title_run_matches_powerpoint(self, open_deck):
        prs = open_deck("theme_only.pptx")
        resolved = _resolved_run_font(prs, 0, "Title 1")
        expected = _run_font_expectation("theme_only", 1, 0)
        # 60pt: layout placeholder; bold + accent1: master titleStyle;
        # Georgia: +mj-lt through the theme's major font.
        assert resolved["size_pt"] == expected["size_pt"]
        assert resolved["name"] == expected["name"]
        assert resolved["bold"] == expected["bold"]
        assert resolved["italic"] == expected["italic"]
        assert resolved["color_rgb"] == expected["color_rgb"]

    def test_subtitle_run_matches_powerpoint(self, open_deck):
        prs = open_deck("theme_only.pptx")
        resolved = _resolved_run_font(prs, 0, "Subtitle 2")
        expected = _run_font_expectation("theme_only", 1, 1)
        # 24pt: layout subTitle lstStyle; 1F4E79: master bodyStyle dk2;
        # Arial: +mn-lt. The subTitle -> bodyStyle routing is the
        # SCPShapeTree fallback miss working as designed.
        assert resolved["size_pt"] == expected["size_pt"]
        assert resolved["name"] == expected["name"]
        assert resolved["color_rgb"] == expected["color_rgb"]

    def test_cascade_source_order_for_placeholder(self, open_deck):
        prs = open_deck("theme_only.pptx")
        slide = prs.slides[0]
        shape = slide.shapes[0]
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        names = [
            source.name
            for source in build_text_cascade_sources(slide, shape, paragraph, run)
        ]
        # No master-placeholder step: ctrTitle has no match on the
        # master tree (documented -- txStyles covers titles instead).
        # No theme-txDef step: fixture themes carry no txDef.
        assert names == [
            "run",
            "shape-lstStyle",
            "layout-placeholder",
            "master-txStyles",
            "presentation-defaultTextStyle",
        ]


@skip_if_fixture_missing("multi_master.pptx")
class TestMultiMasterIntegration:
    def test_alt_master_title(self, open_deck):
        prs = open_deck("multi_master.pptx")
        resolved = _resolved_run_font(prs, 1, "Title 1")
        expected = _run_font_expectation("multi_master", 2, 0)
        # 34pt italic 0B6E4F Times New Roman -- every value from the
        # SECOND master's txStyles + theme, per-slide resolution.
        assert resolved["size_pt"] == expected["size_pt"]
        assert resolved["name"] == expected["name"]
        assert resolved["bold"] == expected["bold"]
        assert resolved["italic"] == expected["italic"]
        assert resolved["color_rgb"] == expected["color_rgb"]

    def test_base_master_title_same_deck(self, open_deck):
        prs = open_deck("multi_master.pptx")
        resolved = _resolved_run_font(prs, 2, "Title 1")
        expected = _run_font_expectation("multi_master", 3, 0)
        assert resolved["size_pt"] == expected["size_pt"]
        assert resolved["name"] == expected["name"]
        assert resolved["color_rgb"] == expected["color_rgb"]

    def test_floating_text_box_uses_default_text_style(self, open_deck):
        prs = open_deck("multi_master.pptx")
        resolved = _resolved_run_font(prs, 3, "FloatNoteCutover")
        expected = _run_font_expectation("multi_master", 4, 1)
        # 18pt Georgia 20232A: presentation defaultTextStyle resolved
        # against the ALT master's theme (tx1 -> dk1 via clrMap).
        assert resolved["size_pt"] == expected["size_pt"]
        assert resolved["name"] == expected["name"]
        assert resolved["color_rgb"] == expected["color_rgb"]

    def test_float_cascade_skips_placeholder_steps(self, open_deck):
        prs = open_deck("multi_master.pptx")
        slide = prs.slides[3]
        shape = next(s for s in slide.shapes if s.name == "FloatNoteCutover")
        names = [
            source.name for source in build_text_cascade_sources(slide, shape)
        ]
        assert "layout-placeholder" not in names
        assert "master-placeholder" not in names
        assert "master-txStyles" not in names
        assert "presentation-defaultTextStyle" in names


@skip_if_fixture_missing("theme_only.pptx")
class TestBuilderValidation:
    def test_requires_slide_and_shape(self, open_deck):
        prs = open_deck("theme_only.pptx")
        with pytest.raises(ValueError):
            build_text_cascade_sources(None, prs.slides[0].shapes[0])
        with pytest.raises(ValueError):
            build_text_cascade_sources(prs.slides[0], None)

    def test_indent_level(self, open_deck):
        prs = open_deck("theme_only.pptx")
        paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        assert indent_level_of(paragraph) == 1
        with pytest.raises(ValueError):
            indent_level_of(None)
