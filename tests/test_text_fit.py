"""Tests for ``utils.text_fit`` -- standalone text-fit prediction
(Step 5 gate d: synthetic overflow + comfortable boxes + the borderline
band).

Font-dependent tests skip when no font file can be resolved on the
machine (bare CI containers); the resolver itself is tested to degrade
to ``(None, None)`` rather than crash.
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from utils.text_fit import (
    BORDERLINE_BAND,
    assess_frame_record,
    predict_text_fit,
    resolve_font_file,
)

_NO_FONTS = resolve_font_file("Arial")[0] is None

needs_fonts = pytest.mark.skipif(
    _NO_FONTS, reason="no resolvable font files on this machine")


def _normalize_bodypr(box, autofit=None):
    """python-pptx textboxes default to wrap="none" + spAutoFit; real
    fixed frames wrap and do not auto-grow, so normalize, then add the
    autofit the test asks for."""
    from lxml import etree

    bodypr = box.text_frame._txBody.find(qn("a:bodyPr"))
    bodypr.set("wrap", "square")
    for tag in ("a:spAutoFit", "a:normAutofit", "a:noAutofit"):
        node = bodypr.find(qn(tag))
        if node is not None:
            bodypr.remove(node)
    if autofit:
        etree.SubElement(bodypr, qn(f"a:{autofit}"))


def _one_frame_deck(tmp_path, text, width_in, height_in, size_pt,
                    autofit=None):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1),
                                   Inches(width_in), Inches(height_in))
    box.name = "Frame"
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = "Arial"
    _normalize_bodypr(box, autofit)
    path = tmp_path / "fit.pptx"
    prs.save(str(path))
    return str(path)


def _frame(result):
    frames = [f for f in result["frames"] if f["shape"] == "Frame"]
    assert len(frames) == 1
    return frames[0]


# --------------------------------------------------------------- resolver

def test_resolve_font_file_never_raises():
    path, family = resolve_font_file("No Such Font Family 123",
                                     bold=True, italic=True)
    assert (path is None) == (family is None)


@needs_fonts
def test_resolve_font_file_reports_substitution_family():
    path, family = resolve_font_file("No Such Font Family 123")
    assert path is not None
    assert family != "No Such Font Family 123"


# --------------------------------------------------------------- verdicts

@needs_fonts
def test_comfortable_box_fits(tmp_path):
    deck = _one_frame_deck(tmp_path, "One short line.", 8.0, 4.0, 12.0)
    frame = _frame(predict_text_fit(deck))
    assert frame["verdict"] == "fits"
    assert frame["ratio"] < 1.0 - BORDERLINE_BAND


@needs_fonts
def test_overflowing_box_predicted(tmp_path):
    text = ("This block of text is far too long for a tiny box. " * 8)
    deck = _one_frame_deck(tmp_path, text, 2.0, 0.5, 24.0)
    frame = _frame(predict_text_fit(deck))
    assert frame["verdict"] == "overflow"
    assert frame["required_pt"] > frame["available_pt"]


@needs_fonts
def test_borderline_box_says_confirm_by_render(tmp_path):
    """A frame sized to ~its own required height lands in the band."""
    probe = _one_frame_deck(tmp_path, "Calibrated single line", 4.0,
                            2.0, 18.0)
    required = _frame(predict_text_fit(probe))["required_pt"]
    # Rebuild the frame with height = required + insets (0.1in total).
    height_in = (required / 72.0) + 0.1
    deck = _one_frame_deck(tmp_path, "Calibrated single line", 4.0,
                           height_in, 18.0)
    frame = _frame(predict_text_fit(deck))
    assert frame["verdict"] == "borderline"


@needs_fonts
def test_sp_autofit_frames_report_fits(tmp_path):
    text = "Way too much text for the declared height " * 10
    deck = _one_frame_deck(tmp_path, text, 2.0, 0.3, 24.0,
                           autofit="spAutoFit")
    frame = _frame(predict_text_fit(deck))
    assert frame["verdict"] == "fits"
    assert frame["autofit"] == "spAutoFit"


@needs_fonts
def test_norm_autofit_frames_are_still_predicted(tmp_path):
    """normAutofit is the stale-fontScale hole -- never skipped."""
    text = "Shrink-on-overflow text that will not fit as declared " * 6
    deck = _one_frame_deck(tmp_path, text, 2.0, 0.5, 24.0,
                           autofit="normAutofit")
    frame = _frame(predict_text_fit(deck))
    assert frame["verdict"] == "overflow"
    assert frame["autofit"] == "normAutofit"


@needs_fonts
def test_line_spacing_override_raises_required_height(tmp_path):
    single = _one_frame_deck(tmp_path, "Two lines of wrapped text "
                             "in a narrowish box", 2.5, 3.0, 14.0)
    base = _frame(predict_text_fit(single))["required_pt"]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2.5),
                                   Inches(3.0))
    box.name = "Frame"
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "Two lines of wrapped text in a narrowish box"
    run.font.size = Pt(14)
    run.font.name = "Arial"
    para.line_spacing = 2.0
    _normalize_bodypr(box)
    path = tmp_path / "spaced.pptx"
    prs.save(str(path))
    doubled = _frame(predict_text_fit(str(path)))["required_pt"]
    assert doubled == pytest.approx(base * 2.0, rel=0.01)


# --------------------------------------------------------------- API

def test_predict_rejects_bad_slide_number(tmp_path):
    deck = _one_frame_deck(tmp_path, "hello", 4.0, 1.0, 12.0)
    with pytest.raises(ValueError, match="slide_number"):
        predict_text_fit(deck, slide_number=9)


def test_predict_filters_by_slide_and_shape(tmp_path):
    deck = _one_frame_deck(tmp_path, "hello", 4.0, 1.0, 12.0)
    result = predict_text_fit(deck, slide_number=1, shape_name="Frame")
    assert [f["shape"] for f in result["frames"]] == ["Frame"]
    assert sum(result["summary"].values()) == 1


def test_predict_is_read_only(tmp_path):
    import hashlib
    from pathlib import Path

    deck = _one_frame_deck(tmp_path, "hello", 4.0, 1.0, 12.0)
    before = hashlib.sha256(Path(deck).read_bytes()).hexdigest()
    predict_text_fit(deck)
    after = hashlib.sha256(Path(deck).read_bytes()).hexdigest()
    assert before == after


@needs_fonts
def test_predict_on_house_corpus_finds_no_overflow():
    """Verification #5 alignment: the COM-rendered house corpus shows
    no overflowing frames, so prediction must not claim one (outside
    the borderline band)."""
    from tests.conftest import HOUSE_CORPUS_DIR, house_corpus_missing

    if house_corpus_missing():
        pytest.skip("house corpus not present")
    result = predict_text_fit(str(HOUSE_CORPUS_DIR / "house_01.pptx"))
    overflowing = [f for f in result["frames"]
                   if f["verdict"] == "overflow"]
    assert overflowing == []
