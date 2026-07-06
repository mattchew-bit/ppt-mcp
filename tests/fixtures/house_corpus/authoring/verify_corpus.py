"""Self-check the authored house corpus (pure python-pptx, no COM).

Checks:
    1. every corpus deck opens via python-pptx with the expected slide
       count (cross-referenced against corpus_truth.json / deviations.json)
    2. corpus_truth.json is internally consistent: >= 24 labeled slides,
       all six archetypes present, archetype slide refs resolve
    3. deviations.json cross-references cleanly: >= 8 violations, every
       referenced slide/shape/paragraph exists in deviant_01.pptx
    4. grid audit: every house-deck shape edge snaps to the seeded
       3-column grid (within tolerance); the deviant deck's off-grid
       shapes are EXACTLY the seeded ones
    5. images: sidebar exhibits on >= 6 slides, inside the recorded zone
    6. tinted elements resolve to real shapes
    7. committed corpus stays under the 8 MB budget

Exit code 0 and "CORPUS VERIFY: PASS" when everything holds.
"""

from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.util import Emu

import _bootstrap
import house_style as hs

MSO_PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE
SIZE_BUDGET_BYTES = 8 * 1024 * 1024
SNAP_EDGES_PT = sorted(hs.GRID_LEFT_EDGES_PT + hs.GRID_RIGHT_EDGES_PT
                       + hs.GRID_CENTER_EDGES_PT)
EXPECTED_OFF_GRID = {("deviant_01", 3, "OffGridPanel"),
                     ("deviant_01", 5, "StragglerNote")}


def _load(filename: str) -> dict:
    with open(_bootstrap.corpus_path(filename), encoding="utf-8") as handle:
        return json.load(handle)


def _pt(emu_value) -> float:
    return Emu(emu_value).pt


def _snapped(value_pt: float) -> bool:
    return any(abs(value_pt - edge) <= hs.GRID_TOLERANCE_PT
               for edge in SNAP_EDGES_PT)


def open_decks(truth: dict, deviations: dict, problems: list[str]) -> dict:
    decks: dict[str, object] = {}
    expected = {name: record["slide_count"]
                for name, record in truth["decks"].items()}
    expected[deviations["deck"]] = len(deviations["slides"])
    for filename, count in expected.items():
        path = _bootstrap.corpus_path(filename)
        if not os.path.exists(path):
            problems.append(f"missing deck: {filename}")
            continue
        presentation = Presentation(path)
        decks[filename] = presentation
        if len(presentation.slides) != count:
            problems.append(f"{filename}: {len(presentation.slides)} "
                            f"slides, metadata says {count}")
    return decks


def check_truth(truth: dict, problems: list[str]) -> None:
    labels = [label for record in truth["decks"].values()
              for label in record["slides"]]
    if len(labels) < 24:
        problems.append(f"only {len(labels)} labeled slides (need >= 24)")
    missing = set(hs.ARCHETYPES) - set(labels)
    if missing:
        problems.append(f"archetypes never used: {sorted(missing)}")
    if truth["labeled_slide_count"] != len(labels):
        problems.append("labeled_slide_count disagrees with deck labels")
    for name, record in truth["archetypes"].items():
        for ref in record["slides"]:
            deck, index = ref.split(":")
            slides = truth["decks"][f"{deck}.pptx"]["slides"]
            if slides[int(index) - 1] != name:
                problems.append(f"archetype ref {ref} is not {name!r}")


def _shape_map(presentation, slide_index: int) -> dict:
    slide = presentation.slides[slide_index - 1]
    return {shape.name: shape for shape in slide.shapes}


def check_deviations(deviations: dict, decks: dict,
                     problems: list[str]) -> None:
    if deviations["violation_count"] < 8:
        problems.append("fewer than 8 seeded violations")
    presentation = decks.get(deviations["deck"])
    if presentation is None:
        return
    for violation in deviations["violations"]:
        where = f"deviation {violation['id']}"
        index = violation["slide"]
        if not 1 <= index <= len(presentation.slides):
            problems.append(f"{where}: slide {index} out of range")
            continue
        shapes = _shape_map(presentation, index)
        shape = shapes.get(violation["shape"])
        if shape is None:
            problems.append(f"{where}: shape {violation['shape']!r} "
                            f"not on slide {index}")
            continue
        paragraph_index = violation.get("paragraph")
        if paragraph_index is not None:
            count = len(shape.text_frame.paragraphs)
            if not 1 <= paragraph_index <= count:
                problems.append(f"{where}: paragraph {paragraph_index} "
                                f"out of range (shape has {count})")


def check_grid(decks: dict, problems: list[str]) -> None:
    off_grid: set[tuple] = set()
    for filename, presentation in decks.items():
        deck = filename.removesuffix(".pptx")
        for index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                if shape.left is None:
                    problems.append(f"{deck}:{index} {shape.name}: "
                                    "no explicit position")
                    continue
                left = _pt(shape.left)
                right = _pt(shape.left + shape.width)
                if not (_snapped(left) and _snapped(right)):
                    off_grid.add((deck, index, shape.name))
    unexpected = off_grid - EXPECTED_OFF_GRID
    missing = EXPECTED_OFF_GRID - off_grid
    for deck, index, name in sorted(unexpected):
        problems.append(f"off-grid shape {deck}:{index} {name!r}")
    for deck, index, name in sorted(missing):
        problems.append(f"seeded off-grid shape {deck}:{index} {name!r} "
                        "unexpectedly snaps to the grid")


def check_images(truth: dict, decks: dict, problems: list[str]) -> None:
    zone = truth["images"]["zones"]["sidebar"]
    found: list[str] = []
    for filename, presentation in decks.items():
        deck = filename.removesuffix(".pptx")
        for index, slide in enumerate(presentation.slides, start=1):
            pictures = [shape for shape in slide.shapes
                        if shape.shape_type == MSO_PICTURE]
            if not pictures:
                continue
            found.append(f"{deck}:{index}")
            for picture in pictures:
                if deck.startswith("house") and (
                        abs(_pt(picture.left) - zone["x_pt"]) > 1.0
                        or abs(_pt(picture.top) - zone["y_pt"]) > 1.0):
                    problems.append(f"{deck}:{index} image off-zone at "
                                    f"({_pt(picture.left):.1f}, "
                                    f"{_pt(picture.top):.1f})pt")
    house_found = [ref for ref in found if ref.startswith("house")]
    if len(house_found) < 6:
        problems.append(f"images on only {len(house_found)} house slides "
                        "(need >= 6)")
    expected_refs = set(truth["images"]["slides_with_images"])
    if expected_refs != set(house_found):
        problems.append("slides_with_images disagrees with actual decks: "
                        f"{sorted(expected_refs ^ set(house_found))}")


def check_tints(truth: dict, decks: dict, problems: list[str]) -> None:
    for tint in truth["tinted_elements"]:
        presentation = decks.get(tint["deck"])
        if presentation is None:
            continue
        shapes = _shape_map(presentation, tint["slide"])
        if tint["shape"] not in shapes:
            problems.append(f"tinted shape {tint['deck']} slide "
                            f"{tint['slide']} {tint['shape']!r} missing")


def check_size(problems: list[str]) -> int:
    total = 0
    for root, _dirs, files in os.walk(_bootstrap.CORPUS_DIR):
        if "__pycache__" in root:
            continue
        total += sum(os.path.getsize(os.path.join(root, name))
                     for name in files)
    if total > SIZE_BUDGET_BYTES:
        problems.append(f"corpus is {total / 1e6:.1f} MB "
                        "(budget 8 MB)")
    return total


def main() -> bool:
    problems: list[str] = []
    truth = _load("corpus_truth.json")
    deviations = _load("deviations.json")
    decks = open_decks(truth, deviations, problems)
    check_truth(truth, problems)
    check_deviations(deviations, decks, problems)
    check_grid(decks, problems)
    check_images(truth, decks, problems)
    check_tints(truth, decks, problems)
    total = check_size(problems)
    print(f"Corpus size: {total / 1e6:.2f} MB across "
          f"{len(decks)} decks")
    if problems:
        print("CORPUS VERIFY: FAIL")
        for problem in problems:
            print(f"  {problem}")
        return False
    print("CORPUS VERIFY: PASS")
    return True


if __name__ == "__main__":
    sys.exit(0 if main() else 1)
