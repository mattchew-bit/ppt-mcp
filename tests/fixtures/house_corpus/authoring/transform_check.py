"""Fixture-arbitrate the resolver's color-transform math on real tints.

Step 2 shipped lumMod/lumOff math that was unit-tested but never
fixture-arbitrated (no Step 0 fixture used color transforms). The house
corpus closes that gap: house_01 carries theme-color brightness variants
(PowerPoint's "Lighter/Darker N%" picker), and this script checks, for
every tinted element, that THREE independent values agree exactly:

    1. COM truth      -- PowerPoint-reported effective RGB, recorded via
                         the EXISTING Step 0 extractor
                         (tests/fixtures/authoring/extract_expected.py)
                         into expected_values/house_01.json
    2. resolver       -- utils.resolve_utils on the corpus .pptx
    3. prediction     -- utils.resolve_colors.apply_color_transforms on
                         the recorded base hex + transforms
                         (already serialized in corpus_truth.json)

It also guards that the deck XML really contains ``a:schemeClr`` with
``a:lumMod`` children (i.e. PowerPoint wrote live theme tints, not baked
literals) so the check can never silently test the wrong thing.

Exit code 0 and "TRANSFORM CHECK: PASS" on full agreement.
"""

from __future__ import annotations

import json
import os
import sys

from pptx import Presentation
from pptx.oxml.ns import qn

import _bootstrap
import extract_expected
from com_helpers import powerpoint_app
from utils.resolve_theme import ThemeContext
from utils.resolve_utils import (
    TextContext,
    resolve_run_font,
    resolve_shape,
)

CHECK_DECK = "house_01"


def load_truth() -> dict:
    with open(_bootstrap.corpus_path("corpus_truth.json"),
              encoding="utf-8") as handle:
        return json.load(handle)


def extract_com_truth() -> dict:
    """Record COM effective values for the tinted deck (Step 0 extractor)."""
    os.makedirs(_bootstrap.EXPECTED_DIR, exist_ok=True)
    deck_path = _bootstrap.corpus_path(f"{CHECK_DECK}.pptx")
    with powerpoint_app() as app:
        document = extract_expected.extract_deck(app, deck_path, CHECK_DECK)
    out_path = os.path.join(_bootstrap.EXPECTED_DIR, f"{CHECK_DECK}.json")
    with open(out_path, "w", encoding="utf-8", newline="\n") as handle:
        json.dump(document, handle, indent=2, ensure_ascii=False)
        handle.write("\n")
    print(f"Wrote {out_path} ({len(document['anomalies'])} anomalies)")
    return document


def _com_shape(document: dict, slide_index: int, shape_name: str) -> dict:
    for record in document["slides"][slide_index - 1]["shapes"]:
        if record["name"] == shape_name:
            return record
    raise LookupError(
        f"shape {shape_name!r} not in COM extract of slide {slide_index}")


def _com_color(record: dict, target: str) -> str:
    if target == "fill":
        return record["fill"]["color_rgb"]
    return record["paragraphs"][0]["runs"][0]["font"]["color_rgb"]


def _pptx_shape(presentation, slide_index: int, shape_name: str):
    slide = presentation.slides[slide_index - 1]
    for shape in slide.shapes:
        if shape.name == shape_name:
            return slide, shape
    raise LookupError(
        f"shape {shape_name!r} not on corpus slide {slide_index}")


def _resolver_color(slide, shape, target: str) -> str:
    if target == "fill":
        return resolve_shape(shape, slide)["fill"]["color_hex"]
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    context = TextContext(slide=slide, shape=shape, paragraph=paragraph,
                          run=run, theme=ThemeContext.for_slide(slide))
    return resolve_run_font(context)["color_hex"]


def _assert_live_tint(shape, target: str, transforms: dict) -> list[str]:
    """The XML must carry schemeClr + the expected transform children."""
    problems = []
    parent = qn("p:spPr") if target == "fill" else qn("a:rPr")
    holder = shape._element.find(f".//{parent}/{qn('a:solidFill')}")
    scheme = None if holder is None else holder.find(qn("a:schemeClr"))
    if scheme is None:
        return [f"{shape.name}: no a:schemeClr under {parent} solidFill "
                "(tint was baked to a literal?)"]
    for name, val in transforms.items():
        child = scheme.find(qn(f"a:{name}"))
        if child is None or child.get("val") != str(val):
            problems.append(
                f"{shape.name}: expected a:{name} val={val}, "
                f"got {'missing' if child is None else child.get('val')}")
    return problems


def check(document: dict, truth: dict) -> list[str]:
    """Compare COM truth vs resolver vs prediction for every tint."""
    tints = [record for record in truth["tinted_elements"]
             if record["deck"] == f"{CHECK_DECK}.pptx"]
    if not tints:
        return [f"corpus_truth.json lists no tints for {CHECK_DECK}"]
    presentation = Presentation(_bootstrap.corpus_path(f"{CHECK_DECK}.pptx"))
    failures: list[str] = []
    for tint in tints:
        where = f"{tint['deck']} slide {tint['slide']} {tint['shape']}"
        com_hex = _com_color(
            _com_shape(document, tint["slide"], tint["shape"]),
            tint["target"])
        slide, shape = _pptx_shape(presentation, tint["slide"],
                                   tint["shape"])
        failures.extend(_assert_live_tint(shape, tint["target"],
                                          tint["xml_transforms"]))
        resolver_hex = _resolver_color(slide, shape, tint["target"])
        predicted = tint["expected_hex"]
        if not (com_hex == resolver_hex == predicted):
            failures.append(
                f"{where} ({tint['target']}, {tint['scheme_token']} "
                f"B{tint['brightness']:+.2f}): COM={com_hex} "
                f"resolver={resolver_hex} predicted={predicted}")
        else:
            print(f"  OK {where}: {com_hex} "
                  f"({tint['scheme_token']} B{tint['brightness']:+.2f})")
    return failures


def main() -> bool:
    truth = load_truth()
    document = extract_com_truth()
    failures = check(document, truth)
    if failures:
        print("TRANSFORM CHECK: FAIL")
        for failure in failures:
            print(f"  {failure}")
        return False
    print("TRANSFORM CHECK: PASS")
    return True


if __name__ == "__main__":
    sys.exit(0 if main() else 1)
