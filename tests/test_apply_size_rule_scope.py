"""Regression tests: typography size-rule level scoping (repair pass).

``utils.style_apply`` documents that the level-1-only restriction on
the size rule applies to the BODY role only (deeper body levels follow
the template's own level scale), but the code gated the size rule on
``level == 1`` for ALL roles -- a title/footer run at indent level > 1
silently missed the size rule. These tests pin the documented contract:

* a TITLE run at indent level 2 with a deviant size IS resized;
* a BODY run at indent level 2 keeps its size (template level scale).
"""

import shutil

import pytest
from pptx import Presentation
from pptx.util import Pt

from tests.conftest import fixture_path
from tests.test_apply_house_profile import (
    DEVIANT,
    HOUSE_PROFILE,
    _shape_named,
)
from utils.style_apply import apply_house_profile

TITLE_RULE_PT = 30.0  # HOUSE_PROFILE typography.title.size
DEVIANT_TITLE_PT = 40.0
DEVIANT_BODY_L2_PT = 13.0  # off the house scale, must survive apply


@pytest.fixture
def deviant_prs(tmp_path):
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    copy_path = tmp_path / "deviant.pptx"
    shutil.copyfile(source, copy_path)
    return Presentation(str(copy_path))


def test_title_run_at_indent_level_2_gets_size_rule(deviant_prs):
    title = _shape_named(deviant_prs.slides[1], "Title 1")
    paragraph = title.text_frame.paragraphs[0]
    paragraph.level = 1  # python-pptx 0-based -> resolver indent level 2
    run = paragraph.runs[0]
    run.font.size = Pt(DEVIANT_TITLE_PT)

    summary = apply_house_profile(deviant_prs, HOUSE_PROFILE)

    assert run.font.size.pt == TITLE_RULE_PT
    assert any(change["property"] == "font.size_pt"
               and change["shape"] == "Title 1"
               for change in summary["changes"])


def test_body_run_at_indent_level_2_size_stays_on_template_scale(
        deviant_prs):
    body = _shape_named(deviant_prs.slides[1], "BodyContent")
    paragraph = body.text_frame.paragraphs[0]
    paragraph.level = 1  # resolver indent level 2: size rule out of scope
    run = paragraph.runs[0]
    run.font.size = Pt(DEVIANT_BODY_L2_PT)

    apply_house_profile(deviant_prs, HOUSE_PROFILE)

    assert run.font.size.pt == DEVIANT_BODY_L2_PT
