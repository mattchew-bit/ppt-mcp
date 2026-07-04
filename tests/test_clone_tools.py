"""Tests for ``tools/clone_tools.py`` -- the MCP-facing duplicate/copy tools.

Verifies registration against a real FastMCP app plus the response-envelope
conventions the other tool modules follow (``{"error": ...}`` on failure,
``{"message": ...}`` on success, presentation registry + current-id default).
"""

import asyncio
import shutil

import pytest
from pptx import Presentation

from tests.conftest import FIXTURE_DECKS, fixture_path, skip_if_fixture_missing
from tests.support_clone import make_chart_deck, rename_all_layouts
from tools.clone_tools import register_clone_tools

FIRST_FIXTURE = FIXTURE_DECKS[0]


class _RecorderApp:
    """Minimal stand-in for FastMCP capturing the registered tool functions."""

    def __init__(self):
        self.tools = {}

    def tool(self, *args, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn

        return decorator


def _registered_tools(presentations, current_id):
    app = _RecorderApp()
    register_clone_tools(app, presentations, lambda: current_id)
    return app.tools


def test_tools_register_on_real_fastmcp_app():
    from mcp.server.fastmcp import FastMCP

    app = FastMCP(name="clone-tools-test")
    register_clone_tools(app, {}, lambda: None)
    tool_names = {tool.name for tool in asyncio.run(app.list_tools())}
    assert {"duplicate_slide", "copy_slide"} <= tool_names


# ------------------------------------------------------------- duplicate_slide


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_duplicate_slide_tool_success():
    prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    original_count = len(prs.slides)
    tools = _registered_tools({"deck": prs}, "deck")

    result = tools["duplicate_slide"](slide_index=0)

    assert "error" not in result
    assert "message" in result
    assert result["new_slide_index"] == original_count
    assert len(prs.slides) == original_count + 1


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_duplicate_slide_tool_explicit_presentation_id():
    prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    tools = _registered_tools({"other": prs}, None)

    result = tools["duplicate_slide"](slide_index=1, presentation_id="other")

    assert "error" not in result
    assert result["new_slide_index"] == len(prs.slides) - 1


def test_duplicate_slide_tool_no_presentation_loaded():
    tools = _registered_tools({}, None)
    result = tools["duplicate_slide"](slide_index=0)
    assert "error" in result


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_duplicate_slide_tool_invalid_index():
    prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    tools = _registered_tools({"deck": prs}, "deck")
    result = tools["duplicate_slide"](slide_index=999)
    assert "error" in result
    assert "slide index" in result["error"].lower()


def test_duplicate_slide_tool_chart_rejected_as_error():
    prs = make_chart_deck()
    tools = _registered_tools({"charts": prs}, "charts")
    result = tools["duplicate_slide"](slide_index=0)
    assert "error" in result
    assert "not supported" in result["error"].lower()


# ------------------------------------------------------------------ copy_slide


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_copy_slide_tool_success(tmp_path):
    src_prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    dst_file = tmp_path / FIRST_FIXTURE
    shutil.copyfile(fixture_path(FIRST_FIXTURE), dst_file)
    dst_prs = Presentation(str(dst_file))
    original_count = len(dst_prs.slides)

    tools = _registered_tools({"src": src_prs, "dst": dst_prs}, "dst")
    result = tools["copy_slide"](source_presentation_id="src", slide_index=0)

    assert "error" not in result
    assert "message" in result
    assert result["new_slide_index"] == original_count
    assert len(dst_prs.slides) == original_count + 1


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_copy_slide_tool_explicit_destination(tmp_path):
    src_prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    dst_file = tmp_path / FIRST_FIXTURE
    shutil.copyfile(fixture_path(FIRST_FIXTURE), dst_file)
    dst_prs = Presentation(str(dst_file))

    tools = _registered_tools({"src": src_prs, "dst": dst_prs}, None)
    result = tools["copy_slide"](
        source_presentation_id="src",
        slide_index=1,
        destination_presentation_id="dst",
    )
    assert "error" not in result


def test_copy_slide_tool_unknown_source():
    tools = _registered_tools({}, None)
    result = tools["copy_slide"](source_presentation_id="ghost", slide_index=0)
    assert "error" in result


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_copy_slide_tool_no_destination_loaded():
    src_prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    tools = _registered_tools({"src": src_prs}, None)
    result = tools["copy_slide"](source_presentation_id="src", slide_index=0)
    assert "error" in result


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_copy_slide_tool_layout_mismatch_surfaces_value_error():
    src_prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    dst_prs = Presentation()
    rename_all_layouts(dst_prs)
    tools = _registered_tools({"src": src_prs, "dst": dst_prs}, "dst")
    result = tools["copy_slide"](source_presentation_id="src", slide_index=0)
    assert "error" in result
    assert "layout" in result["error"].lower()


def test_copy_slide_tool_chart_rejected_as_error():
    src_prs = make_chart_deck()
    dst_prs = Presentation()
    tools = _registered_tools({"src": src_prs, "dst": dst_prs}, "dst")
    result = tools["copy_slide"](source_presentation_id="src", slide_index=0)
    assert "error" in result
    assert "not supported" in result["error"].lower()


@skip_if_fixture_missing(FIRST_FIXTURE)
def test_copy_slide_tool_invalid_index(tmp_path):
    src_prs = Presentation(str(fixture_path(FIRST_FIXTURE)))
    dst_file = tmp_path / FIRST_FIXTURE
    shutil.copyfile(fixture_path(FIRST_FIXTURE), dst_file)
    dst_prs = Presentation(str(dst_file))
    tools = _registered_tools({"src": src_prs, "dst": dst_prs}, "dst")
    result = tools["copy_slide"](source_presentation_id="src", slide_index=42)
    assert "error" in result
    assert "slide index" in result["error"].lower()
