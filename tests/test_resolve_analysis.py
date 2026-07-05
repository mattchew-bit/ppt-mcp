"""Tests for resolved-mode analysis + its ``analyze_presentation_style``
integration (``utils.resolve_analysis`` / ``tools.style_tools``).

Covers: the unchanged legacy tool contract (existing callers pass only
``file_path``), the default-on resolved section, slide_range / detail /
groups controls, the hard response cap with explicit truncation markers,
and input validation errors surfacing as the tool's ``{"error": ...}``
envelope.
"""

import json

import pytest
from pptx import Presentation

from tests.conftest import fixture_path, skip_if_fixture_missing
from tools.style_tools import register_style_tools
from utils.resolve_analysis import (
    MAX_RESPONSE_BYTES,
    build_resolved_analysis,
    parse_slide_range,
)

pytestmark = skip_if_fixture_missing("theme_only.pptx")

FIXTURE = str(fixture_path("theme_only.pptx"))


class _RecorderApp:
    """Minimal FastMCP stand-in capturing registered tool functions."""

    def __init__(self):
        self.tools = {}

    def tool(self, *args, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn

        return decorator


def _analyze_tool():
    app = _RecorderApp()
    register_style_tools(app, {}, lambda: None)
    return app.tools["analyze_presentation_style"]


# ---------------------------------------------------------------------------
# parse_slide_range
# ---------------------------------------------------------------------------

class TestParseSlideRange:
    def test_none_selects_all(self):
        assert parse_slide_range(None, 4) == [0, 1, 2, 3]

    def test_ranges_and_singles(self):
        assert parse_slide_range("1-2,4", 4) == [0, 1, 3]

    def test_deduplicates(self):
        assert parse_slide_range("2,1-3", 4) == [1, 0, 2]

    @pytest.mark.parametrize("bad", ["0", "5", "3-2", "1-9", ",", "x"])
    def test_invalid_specs_raise(self, bad):
        with pytest.raises(ValueError):
            parse_slide_range(bad, 4)


# ---------------------------------------------------------------------------
# build_resolved_analysis
# ---------------------------------------------------------------------------

class TestBuildResolvedAnalysis:
    def test_summary_contains_resolved_effective_values(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs)
        assert report["resolved"] is True
        assert report["truncated"] is False
        # Theme scheme is present (ISC: theme scheme in resolved output).
        assert report["themes"][0]["color_scheme"]["accent1"] == "C0504D"
        summary = report["summary"]
        fonts = {item["value"] for item in summary["fonts"]}
        assert {"Georgia", "Arial"} <= fonts  # inherited, not explicit
        colors = {item["value"] for item in summary["text_colors"]}
        assert "C0504D" in colors  # title accent1 via titleStyle
        # Bullets and paragraph spacing appear (ISC criterion).
        bullets = {item["value"] for item in summary["bullets"]}
        assert "L1:–" in bullets
        assert summary["paragraph_spacing"]

    def test_full_detail_dumps_runs(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs, slide_range="1",
                                         detail="full")
        shape = report["slides"][0]["shapes"][0]
        run_font = shape["paragraphs"][0]["runs"][0]["font"]
        assert run_font == {
            "name": "Georgia", "size_pt": 60.0, "bold": True,
            "italic": False, "color_hex": "C0504D",
        }
        assert shape["fill"]["type"] == "none"

    def test_slide_range_scopes_output(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs, slide_range="2-3",
                                         detail="full")
        assert report["slide_numbers"] == [2, 3]
        assert len(report["slides"]) == 2

    def test_groups_filter(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs, groups=["fonts"])
        assert report["groups"] == ["fonts"]
        assert "themes" not in report
        assert "text_colors" not in report["summary"]
        assert "fonts" in report["summary"]

    def test_cap_truncates_with_marker(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs, detail="full",
                                         max_bytes=4000)
        assert report["truncated"] is True
        assert "hint" in report
        assert len(json.dumps(report)) <= 4000

    def test_default_cap_respected(self):
        prs = Presentation(FIXTURE)
        report = build_resolved_analysis(prs, detail="full")
        assert len(json.dumps(report)) <= MAX_RESPONSE_BYTES

    def test_json_serializable(self):
        prs = Presentation(FIXTURE)
        json.dumps(build_resolved_analysis(prs, detail="full"))

    @pytest.mark.parametrize("kwargs", [
        {"detail": "everything"},
        {"groups": ["fonts", "nope"]},
        {"max_bytes": 0},
    ])
    def test_invalid_controls_raise(self, kwargs):
        prs = Presentation(FIXTURE)
        with pytest.raises(ValueError):
            build_resolved_analysis(prs, **kwargs)


# ---------------------------------------------------------------------------
# Performance regressions
# ---------------------------------------------------------------------------

class TestThemeParseEfficiency:
    """Regression: the theme part must be parsed once per master, not
    once per run/paragraph (the resolver used to re-parse it on every
    cascade build -- 36 parses for 32 runs on theme_only.pptx)."""

    def _count_theme_parses(self, monkeypatch):
        from utils import resolve_theme

        calls = []
        real = resolve_theme.theme_element_for_master

        def counting(master):
            calls.append(str(master.part.partname))
            return real(master)

        monkeypatch.setattr(resolve_theme, "theme_element_for_master",
                            counting)
        return calls

    def test_theme_parsed_once_for_full_analysis(self, monkeypatch):
        calls = self._count_theme_parses(monkeypatch)
        prs = Presentation(FIXTURE)
        build_resolved_analysis(prs, detail="full")
        assert len(prs.slide_masters) == 1
        assert len(calls) == 1

    @skip_if_fixture_missing("multi_master.pptx")
    def test_multi_master_parses_each_theme_once(self, monkeypatch):
        calls = self._count_theme_parses(monkeypatch)
        prs = Presentation(str(fixture_path("multi_master.pptx")))
        build_resolved_analysis(prs, detail="full")
        assert len(calls) == len(prs.slide_masters)


# ---------------------------------------------------------------------------
# House style regressions
# ---------------------------------------------------------------------------

class TestHouseStyle:
    def test_analysis_functions_stay_under_50_lines(self):
        """Regression: ``_summarize_slides`` grew past the <50-line
        house guideline once; keep every function in the module under."""
        import ast
        from pathlib import Path

        import utils.resolve_analysis as module

        source = Path(module.__file__).read_text(encoding="utf-8")
        offenders = [
            (node.name, node.end_lineno - node.lineno + 1)
            for node in ast.walk(ast.parse(source))
            if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef))
            and node.end_lineno - node.lineno + 1 > 50
        ]
        assert offenders == []


# ---------------------------------------------------------------------------
# analyze_presentation_style tool integration
# ---------------------------------------------------------------------------

class TestAnalyzeToolIntegration:
    def test_legacy_call_shape_is_preserved_and_resolved_added(self):
        result = _analyze_tool()(FIXTURE)  # existing callers: path only
        for legacy_key in ("message", "primary_font", "font_count",
                           "color_count", "top_colors", "common_sizes",
                           "text_hierarchy", "consistency_score",
                           "slide_dimensions", "total_shapes",
                           "full_analysis"):
            assert legacy_key in result
        assert result["resolved"]["resolved"] is True
        assert result["resolved"]["detail"] == "summary"

    def test_raw_mode_retained(self):
        result = _analyze_tool()(FIXTURE, resolved=False)
        assert "resolved" not in result
        assert "primary_font" in result

    def test_tool_passes_output_controls(self):
        result = _analyze_tool()(FIXTURE, slide_range="1",
                                 detail="full", groups=["shapes"])
        resolved = result["resolved"]
        assert resolved["slide_numbers"] == [1]
        assert resolved["groups"] == ["shapes"]
        shapes = resolved["slides"][0]["shapes"]
        assert all("paragraphs" not in shape for shape in shapes)
        assert shapes[0]["geometry"]["preset"] == "rect"

    def test_invalid_options_return_error_envelope(self):
        result = _analyze_tool()(FIXTURE, slide_range="99")
        assert "error" in result

    def test_missing_file_error(self):
        result = _analyze_tool()("no/such/deck.pptx")
        assert "error" in result

    def test_file_is_parsed_only_once(self, monkeypatch):
        """Regression: the tool used to open the .pptx twice per call
        (once for the raw analysis, once for the resolved section)."""
        import pptx

        real = pptx.Presentation
        calls = []

        def counting(path):
            calls.append(path)
            return real(path)

        monkeypatch.setattr(pptx, "Presentation", counting)
        result = _analyze_tool()(FIXTURE)
        assert "error" not in result
        assert result["resolved"]["resolved"] is True
        assert len(calls) == 1

    def test_registers_on_real_fastmcp_app(self):
        import asyncio

        from mcp.server.fastmcp import FastMCP

        app = FastMCP(name="style-tools-test")
        register_style_tools(app, {}, lambda: None)
        tools = {tool.name for tool in asyncio.run(app.list_tools())}
        assert "analyze_presentation_style" in tools
