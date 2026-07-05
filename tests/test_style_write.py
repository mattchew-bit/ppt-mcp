"""Tests for ``utils/style_write.py`` -- slide-level write helpers.

The helpers are the write-side counterpart of the read-only Step 2
resolver. Two contracts are pinned here:

1. Every write lands as correct, correctly-ORDERED OOXML at run /
   paragraph / shape level.
2. The never-touch-masters hard rule: every helper refuses targets on
   slideLayout / slideMaster parts (the theme has no shapes to target;
   the byte-identity of all three part families is asserted end-to-end
   in test_apply_house_profile.py).
"""

import shutil

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import fixture_path
from utils.style_write import (
    set_paragraph_bullet,
    set_paragraph_spacing,
    set_run_color,
    set_run_font,
    set_shape_border,
    set_shape_corner_radius,
    set_shape_fill,
)

DEVIANT = "house_corpus/deviant_01.pptx"

EM_DASH = "—"


@pytest.fixture
def deck(tmp_path):
    """A throwaway copy of the deviant corpus deck (never the original)."""
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    copy = tmp_path / "deck.pptx"
    shutil.copyfile(source, copy)
    return Presentation(str(copy))


def shape_named(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise AssertionError(f"no shape named {name!r} on slide")


def r_pr_of(run):
    return run._r.find(qn("a:rPr"))


def p_pr_of(paragraph):
    return paragraph._p.find(qn("a:pPr"))


# ---------------------------------------------------------------------------
# Run-level
# ---------------------------------------------------------------------------

def test_set_run_font_writes_explicit_rpr(deck):
    run = shape_named(deck.slides[1], "Title 1").text_frame.paragraphs[0].runs[0]
    set_run_font(run, name="Georgia", size_pt=20.0, bold=True, italic=False)
    r_pr = r_pr_of(run)
    assert r_pr.get("sz") == "2000"
    assert r_pr.get("b") == "1"
    assert r_pr.get("i") == "0"
    assert r_pr.find(qn("a:latin")).get("typeface") == "Georgia"


def test_set_run_font_writes_only_what_is_given(deck):
    run = shape_named(deck.slides[1],
                      "BodyContent").text_frame.paragraphs[1].runs[0]
    set_run_font(run, size_pt=14.0)
    r_pr = r_pr_of(run)
    assert r_pr.get("sz") == "1400"
    assert r_pr.get("b") is None
    assert r_pr.find(qn("a:latin")) is None


def test_set_run_font_rejects_empty_call_and_bad_values(deck):
    run = shape_named(deck.slides[1], "Title 1").text_frame.paragraphs[0].runs[0]
    with pytest.raises(ValueError):
        set_run_font(run)
    with pytest.raises(ValueError):
        set_run_font(run, size_pt=-3)
    with pytest.raises(ValueError):
        set_run_font(run, name="   ")


def test_set_run_color_srgb(deck):
    run = shape_named(deck.slides[1], "Title 1").text_frame.paragraphs[0].runs[0]
    set_run_color(run, hex_value="#1b7f79")
    srgb = r_pr_of(run).find(qn("a:solidFill")).find(qn("a:srgbClr"))
    assert srgb.get("val") == "1B7F79"


def test_set_run_color_scheme_replaces_hardcoded_srgb(deck):
    run = shape_named(deck.slides[3],
                      "ColumnPanelLeft").text_frame.paragraphs[0].runs[0]
    fill = r_pr_of(run).find(qn("a:solidFill"))
    assert fill.find(qn("a:srgbClr")).get("val") == "20262B"
    set_run_color(run, scheme="dk1")
    fill = r_pr_of(run).find(qn("a:solidFill"))
    assert fill.find(qn("a:srgbClr")) is None
    assert fill.find(qn("a:schemeClr")).get("val") == "dk1"


def test_set_run_color_argument_validation(deck):
    run = shape_named(deck.slides[1], "Title 1").text_frame.paragraphs[0].runs[0]
    with pytest.raises(ValueError):
        set_run_color(run)
    with pytest.raises(ValueError):
        set_run_color(run, hex_value="123456", scheme="dk1")
    with pytest.raises(ValueError):
        set_run_color(run, hex_value="not-a-color")
    with pytest.raises(ValueError):
        set_run_color(run, scheme="accent99")


# ---------------------------------------------------------------------------
# Paragraph spacing
# ---------------------------------------------------------------------------

def test_set_paragraph_spacing_writes_ppr(deck):
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[0]
    set_paragraph_spacing(paragraph, space_before_pt=2.0, space_after_pt=8.0,
                          line_spacing=1.2)
    p_pr = p_pr_of(paragraph)
    assert p_pr.find(qn("a:spcBef"))[0].get("val") == "200"
    assert p_pr.find(qn("a:spcAft"))[0].get("val") == "800"
    line = p_pr.find(qn("a:lnSpc")).find(qn("a:spcPct"))
    assert line.get("val") == "120000"


def test_set_paragraph_spacing_validation(deck):
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[0]
    with pytest.raises(ValueError):
        set_paragraph_spacing(paragraph)
    with pytest.raises(ValueError):
        set_paragraph_spacing(paragraph, space_after_pt=-1)


# ---------------------------------------------------------------------------
# Bullets
# ---------------------------------------------------------------------------

def test_set_paragraph_bullet_creates_ordered_children(deck):
    # Paragraph 3 of slide 2's body has NO a:pPr at all.
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[2]
    assert p_pr_of(paragraph) is None
    set_paragraph_bullet(paragraph, char=EM_DASH, color_scheme="dk2",
                         size_pct=95, font="Arial")
    p_pr = p_pr_of(paragraph)
    tags = [child.tag for child in p_pr]
    order = [tags.index(qn(tag))
             for tag in ("a:buClr", "a:buSzPct", "a:buFont", "a:buChar")]
    assert order == sorted(order), "bullet children out of ECMA-376 order"
    assert p_pr.find(qn("a:buChar")).get("char") == EM_DASH
    assert p_pr.find(qn("a:buSzPct")).get("val") == "95000"
    assert p_pr.find(qn("a:buFont")).get("typeface") == "Arial"
    assert p_pr.find(qn("a:buClr")).find(
        qn("a:schemeClr")).get("val") == "dk2"


def test_set_paragraph_bullet_replaces_existing_char(deck):
    # Paragraph 4 of slide 2's body carries the deviant bullet.
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[3]
    set_paragraph_bullet(paragraph, char=EM_DASH)
    p_pr = p_pr_of(paragraph)
    chars = p_pr.findall(qn("a:buChar"))
    assert len(chars) == 1
    assert chars[0].get("char") == EM_DASH
    # The pre-existing (conformant) buFont must survive untouched.
    assert p_pr.find(qn("a:buFont")).get("typeface") == "Arial"


def test_set_paragraph_bullet_removes_conflicting_bunone(deck):
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[2]
    p_pr = paragraph._p.get_or_add_pPr()
    p_pr.append(p_pr.makeelement(qn("a:buNone"), {}))
    set_paragraph_bullet(paragraph, char=EM_DASH)
    assert p_pr.find(qn("a:buNone")) is None
    assert p_pr.find(qn("a:buChar")).get("char") == EM_DASH


def test_set_paragraph_bullet_validation(deck):
    paragraph = shape_named(deck.slides[1],
                            "BodyContent").text_frame.paragraphs[2]
    with pytest.raises(ValueError):
        set_paragraph_bullet(paragraph)
    with pytest.raises(ValueError):
        set_paragraph_bullet(paragraph, char="ab")
    with pytest.raises(ValueError):
        set_paragraph_bullet(paragraph, size_pct=1000)
    with pytest.raises(ValueError):
        set_paragraph_bullet(paragraph, color_hex="123456",
                             color_scheme="dk1")


# ---------------------------------------------------------------------------
# Shapes
# ---------------------------------------------------------------------------

def test_set_shape_border_weight_only(deck):
    panel = shape_named(deck.slides[2], "TakeawayPanel")
    set_shape_border(panel, weight_pt=1.25)
    ln = panel._element.find(qn("p:spPr")).find(qn("a:ln"))
    assert ln.get("w") == "15875"
    # Existing color / dash untouched.
    assert ln.find(qn("a:solidFill")).find(
        qn("a:srgbClr")).get("val") == "14324F"
    assert ln.find(qn("a:prstDash")).get("val") == "dash"


def test_set_shape_border_color_and_dash(deck):
    panel = shape_named(deck.slides[2], "TakeawayPanel")
    set_shape_border(panel, color_scheme="dk2", dash="sysDash")
    ln = panel._element.find(qn("p:spPr")).find(qn("a:ln"))
    assert ln.find(qn("a:solidFill")).find(
        qn("a:schemeClr")).get("val") == "dk2"
    assert ln.find(qn("a:prstDash")).get("val") == "sysDash"


def test_set_shape_border_validation(deck):
    panel = shape_named(deck.slides[2], "TakeawayPanel")
    with pytest.raises(ValueError):
        set_shape_border(panel)
    with pytest.raises(ValueError):
        set_shape_border(panel, weight_pt=0)
    with pytest.raises(ValueError):
        set_shape_border(panel, dash="squiggle")
    with pytest.raises(ValueError):
        set_shape_border(panel, color_hex="14324F", color_scheme="dk2")


def test_set_shape_fill_hex(deck):
    panel = shape_named(deck.slides[3], "ColumnPanelRight")
    set_shape_fill(panel, hex_value="DCE3E8")
    fill = panel._element.find(qn("p:spPr")).find(qn("a:solidFill"))
    assert fill.find(qn("a:srgbClr")).get("val") == "DCE3E8"


def test_set_shape_corner_radius(deck):
    panel = shape_named(deck.slides[2], "TakeawayPanel")
    set_shape_corner_radius(panel, 0.2)
    geom = panel._element.find(qn("p:spPr")).find(qn("a:prstGeom"))
    gds = geom.find(qn("a:avLst")).findall(qn("a:gd"))
    assert len(gds) == 1
    assert gds[0].get("fmla") == "val 20000"


def test_set_shape_corner_radius_rejects_non_roundrect(deck):
    footer = shape_named(deck.slides[1], "FooterNote")
    with pytest.raises(ValueError):
        set_shape_corner_radius(footer, 0.12)
    panel = shape_named(deck.slides[2], "TakeawayPanel")
    with pytest.raises(ValueError):
        set_shape_corner_radius(panel, 0.9)


# ---------------------------------------------------------------------------
# The never-touch-masters hard rule
# ---------------------------------------------------------------------------

def _first_master_run(deck):
    for shape in deck.slide_masters[0].shapes:
        if shape._element.find(qn("p:txBody")) is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                return paragraph.runs[0]
    raise AssertionError("master carries no text runs")


def test_helpers_refuse_master_targets(deck):
    run = _first_master_run(deck)
    with pytest.raises(ValueError, match="refuses"):
        set_run_font(run, size_pt=14.0)
    with pytest.raises(ValueError, match="refuses"):
        set_run_color(run, scheme="dk1")
    master_shape = deck.slide_masters[0].shapes[0]
    master_paragraph = master_shape.text_frame.paragraphs[0]
    with pytest.raises(ValueError, match="refuses"):
        set_paragraph_spacing(master_paragraph, space_after_pt=8.0)
    with pytest.raises(ValueError, match="refuses"):
        set_paragraph_bullet(master_paragraph, char=EM_DASH)
    with pytest.raises(ValueError, match="refuses"):
        set_shape_border(master_shape, weight_pt=1.25)
    with pytest.raises(ValueError, match="refuses"):
        set_shape_fill(master_shape, hex_value="DCE3E8")
    with pytest.raises(ValueError, match="refuses"):
        set_shape_corner_radius(master_shape, 0.12)


def test_helpers_refuse_layout_targets(deck):
    layout_shape = deck.slides[1].slide_layout.shapes[0]
    with pytest.raises(ValueError, match="refuses"):
        set_shape_fill(layout_shape, hex_value="DCE3E8")
    layout_paragraph = layout_shape.text_frame.paragraphs[0]
    with pytest.raises(ValueError, match="refuses"):
        set_paragraph_spacing(layout_paragraph, space_after_pt=8.0)


def test_helpers_reject_none_target():
    with pytest.raises(ValueError):
        set_shape_fill(None, hex_value="DCE3E8")
