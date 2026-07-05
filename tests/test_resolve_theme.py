"""Tests for utils.resolve_theme -- the per-slide theme engine.

Ground truth: the PowerPoint-authored fixture decks. Their theme seeds
are deliberately non-default (Georgia/Arial majors/minors, C0504D
accent1, ...) so a resolver that silently fell back to Office defaults
would fail loudly here. ``multi_master.pptx`` carries two masters with
different themes; the per-slide tests prove the engine follows the
slide -> layout -> master -> theme relationship chain instead of
assuming ``theme1.xml``.
"""

import pytest
from lxml import etree
from pptx.oxml.ns import qn

from tests.conftest import skip_if_fixture_missing
from utils.resolve_theme import (
    COLOR_SCHEME_SLOTS,
    ThemeContext,
    master_for_slide,
    presentation_default_text_style,
    theme_element_for_master,
)

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Seed values authored into the fixtures (see tests/fixtures/authoring/).
THEME_BASE = {
    "dk1": "1A1A2E", "lt1": "FDFBF7", "dk2": "1F4E79", "lt2": "E8E4D8",
    "accent1": "C0504D", "accent2": "6B9F59", "accent3": "7C5CA6",
    "accent4": "E2A33D", "accent5": "3E8FB0", "accent6": "A63D57",
    "hlink": "2E86AB", "folHlink": "8C5E93",
}
THEME_ALT = {
    "dk1": "20232A", "lt1": "FBF7F0", "dk2": "4A1942",
    "accent1": "0B6E4F",
}


def a_element(fragment: str) -> etree._Element:
    return etree.fromstring(f'<root xmlns:a="{_A_NS}">{fragment}</root>')[0]


@pytest.fixture
def theme_only_context(open_deck):
    prs = open_deck("theme_only.pptx")
    return ThemeContext.for_slide(prs.slides[0])


# ---------------------------------------------------------------------------
# Part navigation
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("theme_only.pptx")
class TestPartNavigation:
    def test_theme_element_root_tag(self, open_deck):
        prs = open_deck("theme_only.pptx")
        master = master_for_slide(prs.slides[0])
        theme = theme_element_for_master(master)
        assert theme.tag == qn("a:theme")

    def test_default_text_style_found(self, open_deck):
        prs = open_deck("theme_only.pptx")
        style = presentation_default_text_style(prs.slides[0])
        assert style is not None
        # Known seed: deck-level default body size is 18pt (sz=1800).
        def_rpr = style.find(qn("a:lvl1pPr") + "/" + qn("a:defRPr"))
        assert def_rpr.get("sz") == "1800"

    def test_none_slide_rejected(self):
        with pytest.raises(ValueError):
            master_for_slide(None)
        with pytest.raises(ValueError):
            theme_element_for_master(None)
        with pytest.raises(ValueError):
            presentation_default_text_style(None)


# ---------------------------------------------------------------------------
# Scheme parsing (theme_only seeds)
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("theme_only.pptx")
class TestThemeOnlySchemes:
    def test_all_twelve_color_slots(self, theme_only_context):
        assert theme_only_context.color_scheme == THEME_BASE
        assert set(COLOR_SCHEME_SLOTS) == set(THEME_BASE)

    def test_font_scheme_seeds(self, theme_only_context):
        assert theme_only_context.major_latin_font == "Georgia"
        assert theme_only_context.minor_latin_font == "Arial"

    def test_theme_font_tokens(self, theme_only_context):
        context = theme_only_context
        assert context.theme_font("+mj-lt") == "Georgia"
        assert context.theme_font("+mn-lt") == "Arial"
        # A concrete typeface is not a token.
        assert context.theme_font("Georgia") is None

    def test_resolve_typeface(self, theme_only_context):
        assert theme_only_context.resolve_typeface("+mn-lt") == "Arial"
        assert theme_only_context.resolve_typeface("Calibri") == "Calibri"
        with pytest.raises(ValueError):
            theme_only_context.resolve_typeface("")

    def test_format_scheme_style_lists(self, theme_only_context):
        # Standard Office fmtScheme carries three entries per list.
        assert len(theme_only_context.fill_styles()) == 3
        assert len(theme_only_context.line_styles()) == 3
        assert len(theme_only_context.effect_styles()) == 3
        assert len(theme_only_context.background_fill_styles()) == 3

    def test_text_default_absent_in_fixture(self, theme_only_context):
        # Fixture themes carry objectDefaults/lnDef but no txDef.
        assert theme_only_context.text_default_list_style() is None


# ---------------------------------------------------------------------------
# clrMap indirection
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("theme_only.pptx")
class TestClrMapIndirection:
    def test_mapped_slots(self, theme_only_context):
        context = theme_only_context
        assert context.scheme_color_hex("tx1") == THEME_BASE["dk1"]
        assert context.scheme_color_hex("bg1") == THEME_BASE["lt1"]
        assert context.scheme_color_hex("tx2") == THEME_BASE["dk2"]
        assert context.scheme_color_hex("bg2") == THEME_BASE["lt2"]

    def test_direct_slots_bypass_map(self, theme_only_context):
        assert theme_only_context.scheme_color_hex("dk1") == THEME_BASE["dk1"]
        assert theme_only_context.scheme_color_hex("accent3") == THEME_BASE["accent3"]

    def test_clr_map_contents(self, theme_only_context):
        clr_map = theme_only_context.clr_map
        assert clr_map["tx1"] == "dk1"
        assert clr_map["bg1"] == "lt1"

    def test_phclr_raises(self, theme_only_context):
        with pytest.raises(ValueError):
            theme_only_context.scheme_color_hex("phClr")

    def test_unknown_name_raises(self, theme_only_context):
        with pytest.raises(ValueError):
            theme_only_context.scheme_color_hex("accent7")


# ---------------------------------------------------------------------------
# Color element resolution through the theme
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("theme_only.pptx")
class TestResolveColor:
    def test_scheme_color_element(self, theme_only_context):
        el = a_element('<a:schemeClr val="accent1"/>')
        assert theme_only_context.resolve_color(el) == "C0504D"

    def test_scheme_color_with_shade(self, theme_only_context):
        # accent2 6B9F59 shaded 50%: 107->54(36), 159->80(50), 89->45(2D)
        el = a_element('<a:schemeClr val="accent2"><a:shade val="50000"/></a:schemeClr>')
        assert theme_only_context.resolve_color(el) == "36502D"

    def test_srgb_passes_through(self, theme_only_context):
        el = a_element('<a:srgbClr val="123abc"/>')
        assert theme_only_context.resolve_color(el) == "123ABC"

    def test_solid_fill(self, theme_only_context):
        fill = a_element('<a:solidFill><a:schemeClr val="tx2"/></a:solidFill>')
        assert theme_only_context.resolve_solid_fill(fill) == THEME_BASE["dk2"]

    def test_empty_solid_fill_raises(self, theme_only_context):
        with pytest.raises(ValueError):
            theme_only_context.resolve_solid_fill(a_element("<a:solidFill/>"))

    def test_schemeclr_without_val_raises(self, theme_only_context):
        with pytest.raises(ValueError):
            theme_only_context.resolve_color(a_element("<a:schemeClr/>"))


# ---------------------------------------------------------------------------
# Multi-master: per-slide theme resolution (never assume theme1.xml)
# ---------------------------------------------------------------------------

@skip_if_fixture_missing("multi_master.pptx")
class TestMultiMasterResolution:
    """Slides 1/3/5 are FixtureBase (theme1); slides 2/4 FixtureAlt (theme2)."""

    def test_base_master_slide(self, open_deck):
        prs = open_deck("multi_master.pptx")
        context = ThemeContext.for_slide(prs.slides[0])
        assert context.color_scheme["accent1"] == THEME_BASE["accent1"]
        assert context.major_latin_font == "Georgia"
        assert context.minor_latin_font == "Arial"

    def test_alt_master_slide(self, open_deck):
        prs = open_deck("multi_master.pptx")
        context = ThemeContext.for_slide(prs.slides[1])
        assert context.color_scheme["accent1"] == THEME_ALT["accent1"]
        assert context.color_scheme["dk1"] == THEME_ALT["dk1"]
        assert context.color_scheme["lt1"] == THEME_ALT["lt1"]
        assert context.color_scheme["dk2"] == THEME_ALT["dk2"]
        assert context.major_latin_font == "Times New Roman"
        assert context.minor_latin_font == "Georgia"

    def test_same_deck_two_themes(self, open_deck):
        prs = open_deck("multi_master.pptx")
        base = ThemeContext.for_slide(prs.slides[0])
        alt = ThemeContext.for_slide(prs.slides[3])
        assert base.color_scheme != alt.color_scheme
        assert base.major_latin_font != alt.major_latin_font

    def test_clr_map_follows_master(self, open_deck):
        prs = open_deck("multi_master.pptx")
        alt = ThemeContext.for_slide(prs.slides[3])
        # tx1 -> dk1 of THAT master's theme (floats on slide 4 render
        # 20232A per the COM-extracted expected values).
        assert alt.scheme_color_hex("tx1") == THEME_ALT["dk1"]

    def test_expected_values_agree(self, open_deck):
        """The engine reproduces the COM-recorded title colors per master."""
        from tests.conftest import load_expected

        expected = load_expected("multi_master")
        prs = open_deck("multi_master.pptx")
        for slide_entry in expected["slides"]:
            slide = prs.slides[slide_entry["index"] - 1]
            context = ThemeContext.for_slide(slide)
            title_font = slide_entry["shapes"][0]["paragraphs"][0]["runs"][0]["font"]
            # Fixture titles are styled schemeClr accent1 on both masters.
            assert context.color_scheme["accent1"] == title_font["color_rgb"]


# ---------------------------------------------------------------------------
# Constructor validation
# ---------------------------------------------------------------------------

class TestConstructorValidation:
    def test_requires_elements(self):
        with pytest.raises(ValueError):
            ThemeContext(None, None)

    def test_requires_theme_elements_child(self):
        bare_theme = etree.fromstring(f'<a:theme xmlns:a="{_A_NS}"/>')
        bare_master = etree.fromstring(
            '<p:sldMaster xmlns:p="http://schemas.openxmlformats.org/'
            'presentationml/2006/main"/>'
        )
        with pytest.raises(ValueError):
            ThemeContext(bare_theme, bare_master)
