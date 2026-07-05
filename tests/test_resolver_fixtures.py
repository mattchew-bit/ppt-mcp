"""Exact-value resolver tests: EVERY COM-recorded effective value, all decks.

Walks the complete ``tests/fixtures/expected_values/<fixture>.json`` for
every fixture deck and asserts the Stage 2 resolver
(``utils.resolve_utils``) reproduces every recorded value: every run's
font name / size / bold / italic / color, every paragraph's alignment /
spacing / bullet, every shape's line / fill / geometry. No sampling, no
skips, no xfail -- a mismatch is a resolver bug.

Comparison adapter (representation differences between the COM object
model and the resolver's OOXML-derived output; each case verified
against the raw XML during Stage 2 development):

    * COM geometry is in points; the resolver converts EMU / 12700.
      Expected values were rounded to 2 dp at extraction, so floats
      compare with tolerance 0.01.
    * COM ``IndentLevel`` is 1-based; the resolver's ``indent_level`` is
      ``a:pPr/@lvl + 1`` -- both 1-based, compared directly.
    * COM ``SpaceWithin`` with rule "multiple" (e.g. 1.15) corresponds
      to ``a:spcPct val="115000"``; rule "points" to ``a:spcPts``.
      ``space_before_pt``/``space_before_lines`` mirror the same split
      for ``a:spcBef`` (and after).
    * Colors compare as uppercase ``RRGGBB`` hex on both sides.
    * Enum translations (``_COM_*`` tables below): MsoLineDashStyle /
      MsoFillType / PpParagraphAlignment / PpBulletType ints -> OOXML
      tokens. ``AutoShapeType`` 138 is what COM reports for text
      placeholders, whose XML carries no ``a:prstGeom`` at any level of
      the chain -- ECMA's default geometry is ``rect``, so 138 maps to
      ``rect`` (empirical, placeholder-only).
    * A bullet with ``color_source: "follow_text"`` (COM
      ``Bullet.UseTextColor`` = msoTrue) paints with the paragraph's
      first-run color; the resolver reports ``color_follows_text`` and
      the adapter substitutes the resolved first-run color. Font and
      relative size fall back the same way (``bu*Tx`` semantics,
      relative size default 1.0).
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import fixture_missing, fixture_path, load_expected
from utils.resolve_theme import ThemeContext
from utils.resolve_utils import (
    TextContext,
    resolve_paragraph,
    resolve_run_font,
    resolve_shape,
)

FIXTURES = ("theme_only", "layout_override", "explicit_override",
            "multi_master")

#: Float comparison tolerance in points / fraction units.
TOLERANCE = 0.01

# -- COM enum -> resolver token adapters -------------------------------------

_COM_ALIGNMENT = {
    1: "left", 2: "center", 3: "right", 4: "justify",
    5: "distribute", 6: "thai_distribute", 7: "justify_low",
}
_COM_BULLET_TYPE = {0: "none", 1: "char", 2: "autonum"}
_COM_DASH = {
    1: "solid", 2: "sysDash", 3: "sysDot", 4: "dash", 5: "dashDot",
    6: "sysDashDotDot", 7: "lgDash", 8: "lgDashDot",
}
_COM_FILL = {1: "solid", 2: "pattern", 3: "gradient", 6: "picture"}
_COM_AUTO_SHAPE = {
    1: "rect", 5: "roundRect", 9: "ellipse", 52: "chevron",
    138: "rect",  # placeholder without prstGeom (see module docstring)
}

#: Grand totals across every parametrized case (reported by the final
#: volume test so silent under-walking cannot pass unnoticed).
_TOTALS = {"checks": 0, "runs": 0, "paragraphs": 0, "shapes": 0}


class Comparison:
    """Collects check counts + mismatch descriptions for one slide."""

    def __init__(self):
        self.checks = 0
        self.mismatches = []

    def check(self, path, expected, actual):
        self.checks += 1
        if expected != actual:
            self.mismatches.append(
                f"{path}: expected {expected!r}, got {actual!r}"
            )

    def check_close(self, path, expected, actual):
        self.checks += 1
        if actual is None or abs(actual - expected) > TOLERANCE:
            self.mismatches.append(
                f"{path}: expected {expected!r}, got {actual!r}"
            )


# -- adapters ----------------------------------------------------------------

def _expected_font_name(font: dict) -> str:
    """COM-recorded font name (theme tokens carry name_resolved)."""
    name = font["name"]
    if name.startswith("+"):
        return font["name_resolved"]
    return name


def _compare_run(cmp, path, expected_run, resolved_font, run):
    font = expected_run["font"]
    cmp.check(f"{path}.text", expected_run["text"], run.text)
    cmp.check(f"{path}.font.name", _expected_font_name(font),
              resolved_font["name"])
    cmp.check_close(f"{path}.font.size_pt", font["size_pt"],
                    resolved_font["size_pt"])
    cmp.check(f"{path}.font.bold", font["bold"], resolved_font["bold"])
    cmp.check(f"{path}.font.italic", font["italic"],
              resolved_font["italic"])
    cmp.check(f"{path}.font.color", font["color_rgb"],
              resolved_font["color_hex"])


def _compare_spacing(cmp, path, expected_para, key, resolved_value):
    """space_before / space_after: *_pt vs {"points"}, *_lines vs {"lines"}."""
    pt_key, lines_key = f"{key}_pt", f"{key}_lines"
    if pt_key in expected_para:
        cmp.check(f"{path}.{key}.tag", "points",
                  next(iter(resolved_value)))
        cmp.check_close(f"{path}.{key}", expected_para[pt_key],
                        resolved_value.get("points"))
    if lines_key in expected_para:
        cmp.check(f"{path}.{key}.tag", "lines",
                  next(iter(resolved_value)))
        cmp.check_close(f"{path}.{key}", expected_para[lines_key],
                        resolved_value.get("lines"))


def _compare_line_spacing(cmp, path, expected_para, resolved_value):
    rule = expected_para["space_within_rule"]
    tag = "multiple" if rule == "multiple" else "points"
    cmp.check(f"{path}.line_spacing.tag", tag, next(iter(resolved_value)))
    cmp.check_close(f"{path}.line_spacing", expected_para["space_within"],
                    resolved_value.get(tag))


def _bullet_relative_size(bullet: dict, run_size_pt) -> float:
    """Resolver bullet size -> COM RelativeSize (default 1.0)."""
    if bullet.get("size_pct") is not None:
        return bullet["size_pct"]
    if bullet.get("size_pt") is not None and run_size_pt:
        return bullet["size_pt"] / run_size_pt
    return 1.0


def _compare_bullet(cmp, path, expected_bullet, bullet, first_run_font):
    cmp.check(f"{path}.type", _COM_BULLET_TYPE[expected_bullet["type"]],
              bullet["type"])
    cmp.check(f"{path}.visible", expected_bullet["visible"],
              bullet["type"] != "none")
    if bullet["type"] == "none" or expected_bullet["type"] == 0:
        return
    if "char_code" in expected_bullet:
        cmp.check(f"{path}.char_code", expected_bullet["char_code"],
                  bullet.get("char_code"))
        cmp.check(f"{path}.char", expected_bullet["char"],
                  bullet.get("char"))
    if "font_name" in expected_bullet:
        actual_font = (first_run_font["name"]
                       if bullet["font_follows_text"] else bullet["font"])
        cmp.check(f"{path}.font_name", expected_bullet["font_name"],
                  actual_font)
    if "relative_size" in expected_bullet:
        cmp.check_close(
            f"{path}.relative_size", expected_bullet["relative_size"],
            _bullet_relative_size(bullet, first_run_font["size_pt"]))
    if "color_rgb" in expected_bullet:
        source = expected_bullet.get("color_source")
        cmp.check(f"{path}.color_source", source == "follow_text",
                  bullet["color_follows_text"])
        actual_color = (first_run_font["color_hex"]
                        if bullet["color_follows_text"]
                        else bullet["color_hex"])
        cmp.check(f"{path}.color", expected_bullet["color_rgb"],
                  actual_color)


def _compare_paragraph(cmp, path, expected_para, slide, shape, paragraph,
                       theme):
    resolved = resolve_paragraph(
        TextContext(slide, shape, paragraph, None, theme))
    cmp.check(f"{path}.level", expected_para["level"],
              resolved["indent_level"])
    cmp.check(f"{path}.alignment",
              _COM_ALIGNMENT[expected_para["alignment"]],
              resolved["alignment"])
    _compare_spacing(cmp, path, expected_para, "space_before",
                     resolved["space_before"])
    _compare_spacing(cmp, path, expected_para, "space_after",
                     resolved["space_after"])
    _compare_line_spacing(cmp, path, expected_para,
                          resolved["line_spacing"])

    runs = list(paragraph.runs)
    expected_runs = expected_para["runs"]
    cmp.check(f"{path}.run_count", len(expected_runs), len(runs))
    resolved_fonts = []
    for j, (expected_run, run) in enumerate(zip(expected_runs, runs)):
        font = resolve_run_font(
            TextContext(slide, shape, paragraph, run, theme))
        resolved_fonts.append(font)
        _compare_run(cmp, f"{path}.runs[{j}]", expected_run, font, run)
        _TOTALS["runs"] += 1

    if resolved_fonts:
        _compare_bullet(cmp, f"{path}.bullet", expected_para["bullet"],
                        resolved["bullet"], resolved_fonts[0])
    _TOTALS["paragraphs"] += 1


def _compare_geometry(cmp, path, expected_geo, resolved):
    geometry = resolved["geometry"]
    cmp.check_close(f"{path}.left_pt", expected_geo["left_pt"],
                    geometry["left_pt"])
    cmp.check_close(f"{path}.top_pt", expected_geo["top_pt"],
                    geometry["top_pt"])
    cmp.check_close(f"{path}.width_pt", expected_geo["width_pt"],
                    geometry["width_pt"])
    cmp.check_close(f"{path}.height_pt", expected_geo["height_pt"],
                    geometry["height_pt"])
    cmp.check_close(f"{path}.rotation", expected_geo["rotation"],
                    resolved["rotation_deg"])
    cmp.check(f"{path}.preset",
              _COM_AUTO_SHAPE[expected_geo["auto_shape_type"]],
              geometry["preset"])
    expected_adj = expected_geo["adjustments"]
    actual_adj = resolved["adjustments"]
    cmp.check(f"{path}.adjustment_count", len(expected_adj),
              len(actual_adj))
    for k, (want, got) in enumerate(zip(expected_adj, actual_adj)):
        cmp.check_close(f"{path}.adjustments[{k}]", want, got)


def _compare_line(cmp, path, expected_line, resolved_line):
    cmp.check(f"{path}.visible", expected_line["visible"],
              resolved_line["visible"])
    if not expected_line["visible"]:
        return
    cmp.check_close(f"{path}.weight_pt", expected_line["weight_pt"],
                    resolved_line["weight_pt"])
    cmp.check(f"{path}.dash", _COM_DASH[expected_line["dash_style"]],
              resolved_line["dash"])
    cmp.check(f"{path}.color", expected_line["color_rgb"],
              resolved_line["color_hex"])


def _compare_fill(cmp, path, expected_fill, resolved_fill):
    if not expected_fill["visible"]:
        cmp.check(f"{path}.type", "none", resolved_fill["type"])
        return
    cmp.check(f"{path}.type", _COM_FILL[expected_fill["type"]],
              resolved_fill["type"])
    if "color_rgb" in expected_fill:
        cmp.check(f"{path}.color", expected_fill["color_rgb"],
                  resolved_fill["color_hex"])


def _compare_shape(cmp, path, expected_shape, slide, shape, theme):
    cmp.check(f"{path}.name", expected_shape["name"], shape.name)
    cmp.check(f"{path}.is_placeholder", expected_shape["is_placeholder"],
              shape.is_placeholder)
    resolved = resolve_shape(shape, slide, theme)
    _compare_geometry(cmp, f"{path}.geometry", expected_shape["geometry"],
                      resolved)
    _compare_line(cmp, f"{path}.line", expected_shape["line"],
                  resolved["line"])
    _compare_fill(cmp, f"{path}.fill", expected_shape["fill"],
                  resolved["fill"])

    paragraphs = (list(shape.text_frame.paragraphs)
                  if shape.has_text_frame else [])
    expected_paras = expected_shape["paragraphs"]
    cmp.check(f"{path}.paragraph_count", len(expected_paras),
              len(paragraphs))
    for i, (expected_para, paragraph) in enumerate(
            zip(expected_paras, paragraphs)):
        _compare_paragraph(cmp, f"{path}.paragraphs[{i}]", expected_para,
                           slide, shape, paragraph, theme)
    _TOTALS["shapes"] += 1


# -- the exhaustive walk ------------------------------------------------------

def _cases():
    cases = []
    for fixture in FIXTURES:
        expected = load_expected(fixture)
        for slide_entry in expected["slides"]:
            cases.append((fixture, slide_entry["index"]))
    return cases


@pytest.fixture(scope="module")
def decks():
    """One parsed Presentation per fixture (module-scoped for speed)."""
    opened = {}
    for fixture in FIXTURES:
        if not fixture_missing(f"{fixture}.pptx"):
            opened[fixture] = Presentation(
                str(fixture_path(f"{fixture}.pptx")))
    return opened


@pytest.mark.parametrize("fixture,slide_number", _cases())
def test_slide_reproduces_every_recorded_value(decks, fixture,
                                               slide_number):
    if fixture not in decks:
        pytest.skip(f"fixture {fixture}.pptx not present")
    prs = decks[fixture]
    expected = load_expected(fixture)
    slide_entry = next(entry for entry in expected["slides"]
                       if entry["index"] == slide_number)
    slide = prs.slides[slide_number - 1]  # COM indices are 1-based
    theme = ThemeContext.for_slide(slide)

    cmp = Comparison()
    cmp.check(f"{fixture}[{slide_number}].layout_name",
              slide_entry["layout_name"], slide.slide_layout.name)
    master_name = slide.slide_layout.slide_master.element.find(
        qn("p:cSld")).get("name")
    cmp.check(f"{fixture}[{slide_number}].master_name",
              slide_entry["master_name"], master_name)

    shapes = list(slide.shapes)
    cmp.check(f"{fixture}[{slide_number}].shape_count",
              len(slide_entry["shapes"]), len(shapes))
    for i, (expected_shape, shape) in enumerate(
            zip(slide_entry["shapes"], shapes)):
        _compare_shape(cmp, f"{fixture}[{slide_number}].shapes[{i}]",
                       expected_shape, slide, shape, theme)

    _TOTALS["checks"] += cmp.checks
    assert not cmp.mismatches, (
        f"{len(cmp.mismatches)} mismatch(es) of {cmp.checks} checks:\n"
        + "\n".join(cmp.mismatches)
    )


@pytest.mark.parametrize("fixture", FIXTURES)
def test_slide_dimensions(decks, fixture):
    if fixture not in decks:
        pytest.skip(f"fixture {fixture}.pptx not present")
    expected = load_expected(fixture)["slide_size_pt"]
    prs = decks[fixture]
    assert abs(prs.slide_width / 12700 - expected["width"]) <= TOLERANCE
    assert abs(prs.slide_height / 12700 - expected["height"]) <= TOLERANCE


@pytest.mark.parametrize("fixture", FIXTURES)
def test_no_extraction_anomalies_recorded(fixture):
    """The expected-values files must be anomaly-free ground truth."""
    assert load_expected(fixture)["anomalies"] == []


def test_zz_assertion_volume(decks):
    """Guard against silent under-walking: the exhaustive walk must have
    compared a large, known-order-of-magnitude number of values."""
    if len(decks) < len(FIXTURES):
        pytest.skip("not all fixture decks present")
    assert _TOTALS["checks"] > 1500, _TOTALS
    assert _TOTALS["runs"] >= 130, _TOTALS
    assert _TOTALS["paragraphs"] >= 122, _TOTALS
    assert _TOTALS["shapes"] >= 55, _TOTALS
