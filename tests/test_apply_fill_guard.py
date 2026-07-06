"""Regression tests: shape-fill rule guards (Step 3 repair pass).

Two related defects in ``utils.style_apply._apply_shape_defaults``:

* MATERIAL -- the off-palette fill rule tested the RESOLVED effective
  hex for membership in the 12 base scheme hexes, so theme-linked
  TINTED fills (``a:schemeClr`` + lumMod/lumOff children) resolved
  off-scheme and were destructively repainted as hardcoded
  ``a:srgbClr`` on CONFORMANT decks: applying the learned profile to
  its own corpus rewrote house_01 (3 fills) and house_02 (1 fill),
  destroying their theme links.
* CONFIRMED -- the rule fired whenever the effective fill was
  off-scheme instead of ``!= the house fill``, so a profile whose
  house fill is itself off-scheme (realistic: panel fills are often
  theme tints) rewrote the already-house-true fill on EVERY apply,
  violating the unconditional second-apply no-op contract.

The fix restricts the snap to fills hardcoded as a literal,
untransformed ``a:srgbClr`` whose hex is off-palette AND differs from
the house fill; theme-linked / transformed / inherited fills are never
candidates. These tests pin both guards. The still-caught case (a
literal off-palette srgb, seeded violation v8) stays covered by
tests/test_apply_house_profile.py::test_v8_off_palette_fill_snapped_to_house_fill.
"""

import copy
import shutil
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from tests.conftest import fixture_path, house_corpus_paths
from tests.test_apply_house_profile import (
    DEVIANT,
    HOUSE_PROFILE,
    _shape_named,
)
from utils.style_apply import apply_house_profile

#: A fill hex that is NOT one of the 12 Meridian scheme hexes.
OFF_SCHEME_HOUSE_FILL = "EEF2F5"

#: Seeded violation v8: ColumnPanelRight carries a literal srgb fill.
V8_LITERAL_HEX = "8E44AD"


def _off_scheme_fill_profile() -> dict:
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_off_scheme_fill"
    profile["shape_defaults"]["fill"] = {
        "value": f"#{OFF_SCHEME_HOUSE_FILL}"}
    return profile


def _deviant_copy(tmp_path):
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    copy_path = tmp_path / "deviant.pptx"
    shutil.copyfile(source, copy_path)
    return copy_path


def _fill_changes(summary):
    return [change for change in summary["changes"]
            if change["property"] == "fill.color"]


# ---------------------------------------------------------------------------
# MATERIAL: the learned profile applied to its own corpus is a no-op
# ---------------------------------------------------------------------------

def test_apply_learned_profile_to_own_corpus_is_zero_writes(
        house_profile, tmp_path):
    """Conformant house decks must come out of apply untouched."""
    for source in house_corpus_paths():
        copy_path = tmp_path / f"copy_{Path(source).name}"
        shutil.copyfile(source, copy_path)
        summary = apply_house_profile(Presentation(str(copy_path)),
                                      house_profile)
        assert summary["writes"] == 0, (
            f"apply rewrote conformant corpus deck {source}: "
            f"{summary['changes']}")
        assert summary["slides_touched"] == []


def test_scheme_tint_fill_xml_untouched_by_apply(house_profile):
    """A ``schemeClr accent1`` + lumMod/lumOff fill resolves off the 12
    base hexes but is theme-true; apply must not repaint it."""
    corpus = house_corpus_paths()
    prs = Presentation(corpus[0])  # house_01: AccentRule on slide 3
    shape = _shape_named(prs.slides[2], "AccentRule")
    before = etree.tostring(shape._element)

    apply_house_profile(prs, house_profile)

    assert etree.tostring(shape._element) == before
    fill = shape._element.find(qn("p:spPr")).find(qn("a:solidFill"))
    scheme = fill.find(qn("a:schemeClr"))
    assert scheme is not None and scheme.get("val") == "accent1"
    assert scheme.find(qn("a:lumMod")) is not None


# ---------------------------------------------------------------------------
# CONFIRMED: off-scheme house fill must stay idempotent
# ---------------------------------------------------------------------------

def test_off_scheme_house_fill_snaps_v8_once_then_second_apply_noop(
        tmp_path):
    profile = _off_scheme_fill_profile()
    copy_path = _deviant_copy(tmp_path)

    prs = Presentation(str(copy_path))
    first = apply_house_profile(prs, profile)
    assert _fill_changes(first) == [{
        "slide": 4, "shape": "ColumnPanelRight", "property": "fill.color",
        "from": V8_LITERAL_HEX, "to": OFF_SCHEME_HOUSE_FILL,
    }]

    saved = tmp_path / "applied.pptx"
    prs.save(str(saved))
    second = apply_house_profile(Presentation(str(saved)), profile)
    assert second["writes"] == 0
    assert second["changes"] == []
    assert second["slides_touched"] == []


def test_literal_fill_already_equal_to_off_scheme_house_fill_untouched(
        tmp_path):
    """A hardcoded srgb fill that IS the house fill is conformant even
    when the house fill sits outside the 12 scheme hexes."""
    profile = _off_scheme_fill_profile()
    copy_path = _deviant_copy(tmp_path)
    prs = Presentation(str(copy_path))

    panel = _shape_named(prs.slides[3], "ColumnPanelRight")
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor.from_string(OFF_SCHEME_HOUSE_FILL)

    summary = apply_house_profile(prs, profile)
    assert _fill_changes(summary) == []
    srgb = (panel._element.find(qn("p:spPr"))
            .find(qn("a:solidFill")).find(qn("a:srgbClr")))
    assert srgb.get("val").upper() == OFF_SCHEME_HOUSE_FILL
