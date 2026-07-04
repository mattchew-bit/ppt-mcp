"""Shared helpers for the Step 1 slide-clone tests.

Not collected by pytest (no ``test_`` prefix). Provides:
    * inline deck builders (image / chart / hyperlink decks) used only for
      mechanics tests -- fixture decks stay PowerPoint-authored,
    * package-level assertions (rel validity, [Content_Types].xml coverage),
    * slide comparison helpers (signature + canonical XML equality).
"""

import io
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path, PurePosixPath

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A16_CREATION_ID = (
    "{http://schemas.microsoft.com/office/drawing/2014/main}creationId"
)
P14_CREATION_ID = (
    "{http://schemas.microsoft.com/office/powerpoint/2010/main}creationId"
)
CT_NS = "{http://schemas.openxmlformats.org/package/2006/content-types}"
REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"
P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


# ---------------------------------------------------------------- deck builders


def make_image_deck():
    """In-memory deck: one blank slide with a picture and a text box.

    Inline python-pptx construction is acceptable here: these tests cover
    rel-rewriting mechanics, not style inheritance (which uses the
    PowerPoint-authored fixtures).
    """
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.util import Inches

    buf = io.BytesIO()
    PILImage.new("RGB", (8, 8), (200, 30, 30)).save(buf, "PNG")
    buf.seek(0)

    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(buf, Inches(1), Inches(1))
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
    box.text_frame.text = "image deck"
    return prs


def make_two_image_deck():
    """In-memory deck: one slide with two *different* pictures.

    Regression deck: both images must land in the destination as distinct
    parts with distinct partnames (partname allocation must see parts
    created earlier in the same copy operation).
    """
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide = prs.slides.add_slide(blank)
    for offset, color in enumerate([(200, 30, 30), (30, 30, 200)]):
        buf = io.BytesIO()
        PILImage.new("RGB", (8, 8), color).save(buf, "PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(1 + offset * 2), Inches(1))
    return prs


def make_chart_deck():
    """In-memory deck with a single chart slide (negative-test only)."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide = prs.slides.add_slide(blank)
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("series 1", (1.0, 2.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(4), Inches(3),
        data,
    )
    return prs


def make_hyperlink_deck(url: str):
    """In-memory deck: one slide whose text run carries an external hyperlink."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "click me"
    run.hyperlink.address = url
    return prs


def make_movie_deck():
    """In-memory deck: one blank slide bearing a movie shape.

    python-pptx movie shapes always emit an action-only
    ``<a:hlinkClick r:id="" action="ppaction://media"/>`` -- the empty-rId
    convention that must never be treated as a dangling relationship.
    The video blob content is irrelevant to rel mechanics, so a stub
    passes; python-pptx stores it opaquely.
    """
    from pptx import Presentation
    from pptx.util import Inches

    fake_mp4 = b"\x00\x00\x00\x18ftypmp42" + b"\x00" * 64
    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_movie(
        io.BytesIO(fake_mp4),
        Inches(1), Inches(1), Inches(3), Inches(2),
        mime_type="video/mp4",
    )
    return prs


def make_slide_jump_deck():
    """In-memory deck: slide 0 has a text run jumping to slide 1 (RT.SLIDE).

    Slide 1 bears a picture so a blob-only sweep of the jump target would
    leave a dangling r:embed -- the regression this deck exists to catch.
    """
    from lxml import etree
    from PIL import Image as PILImage
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn
    from pptx.util import Inches

    prs = Presentation()
    blank = _layout_by_name(prs, "Blank")
    slide0 = prs.slides.add_slide(blank)
    slide1 = prs.slides.add_slide(blank)

    buf = io.BytesIO()
    PILImage.new("RGB", (8, 8), (30, 200, 30)).save(buf, "PNG")
    buf.seek(0)
    slide1.shapes.add_picture(buf, Inches(1), Inches(1))

    box = slide0.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "jump to slide 2"
    rid = slide0.part.relate_to(slide1.part, RT.SLIDE)
    rPr = run._r.get_or_add_rPr()
    hlink = etree.SubElement(rPr, qn("a:hlinkClick"))
    hlink.set(qn("r:id"), rid)
    hlink.set("action", "ppaction://hlinksldjump")
    return prs


def _layout_by_name(prs, name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise LookupError(f"no layout named {name!r} in deck")


def rename_all_layouts(prs, prefix: str = "zzz_nonmatching_") -> None:
    """Give every layout in ``prs`` a name no other deck will match."""
    counter = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            layout.name = f"{prefix}{counter}"
            counter += 1


# ------------------------------------------------------------ slide comparison


def slide_signature(slide) -> list:
    """Order-sensitive shape signature: type, name, geometry, text."""
    signature = []
    for shape in slide.shapes:
        signature.append({
            "shape_type": str(shape.shape_type),
            "name": shape.name,
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "text": shape.text_frame.text if shape.has_text_frame else None,
        })
    return signature


def canonical_cSld(slide_xml: bytes) -> str:
    """Canonical XML of a slide's ``p:cSld`` with creationIds removed.

    Covers shape tree, text, geometry, and any locally-defined background.
    """
    root = ET.fromstring(slide_xml)
    _strip_creation_ids_inplace(root)
    cSld = root.find(f"{P_NS}cSld")
    assert cSld is not None, "slide XML has no p:cSld element"
    return ET.canonicalize(ET.tostring(cSld, encoding="unicode"))


def _strip_creation_ids_inplace(root) -> None:
    parents = {child: parent for parent in root.iter() for child in parent}
    doomed = [
        el for el in root.iter()
        if el.tag in (A16_CREATION_ID, P14_CREATION_ID)
    ]
    for el in doomed:
        parents[el].remove(el)
    # prune ext wrappers emptied by the removal above
    for el in list(root.iter()):
        if el.tag.endswith("}ext") and len(el) == 0 and el in parents:
            parents[el].remove(el)


# ------------------------------------------------------- package-level checks


def slide_part_names(pptx_path) -> list:
    """Sorted slide part names (``ppt/slides/slideN.xml``) in a saved deck."""
    with zipfile.ZipFile(pptx_path) as archive:
        return sorted(
            name for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )


def read_zip(pptx_path, member: str) -> bytes:
    with zipfile.ZipFile(pptx_path) as archive:
        return archive.read(member)


def assert_slide_rels_valid(pptx_path) -> None:
    """Every r:id/r:embed in every slide resolves; every internal rel target exists."""
    with zipfile.ZipFile(pptx_path) as archive:
        names = set(archive.namelist())
        for name in sorted(names):
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            rels_name = f"ppt/slides/_rels/{PurePosixPath(name).name}.rels"
            rels = {}
            if rels_name in names:
                rels_root = ET.fromstring(archive.read(rels_name))
                for rel in rels_root.iter(f"{REL_NS}Relationship"):
                    rels[rel.get("Id")] = (rel.get("Target"), rel.get("TargetMode"))

            root = ET.fromstring(archive.read(name))
            # Empty r:id values are the OOXML convention for action-only
            # hyperlinks (e.g. python-pptx movie shapes emit
            # <a:hlinkClick r:id="" action="ppaction://media"/>) and do not
            # reference a relationship.
            referenced = {
                value
                for el in root.iter()
                for key, value in el.attrib.items()
                if key.startswith("{" + R_NS + "}") and value
            }
            dangling = referenced - set(rels)
            assert not dangling, f"{name}: r:id refs missing from rels: {dangling}"

            for rid, (target, mode) in rels.items():
                if mode == "External":
                    continue
                resolved = _resolve_target("ppt/slides", target)
                assert resolved in names, (
                    f"{name}: rel {rid} target {target!r} -> {resolved!r} "
                    "missing from package"
                )


def _resolve_target(base_dir: str, target: str) -> str:
    parts = list(PurePosixPath(base_dir).parts)
    for segment in PurePosixPath(target).parts:
        if segment == "..":
            parts.pop()
        elif segment != ".":
            parts.append(segment)
    return "/".join(parts)


def assert_content_types_cover_all_parts(pptx_path) -> None:
    """[Content_Types].xml must cover every part in the package."""
    with zipfile.ZipFile(pptx_path) as archive:
        names = archive.namelist()
        root = ET.fromstring(archive.read("[Content_Types].xml"))
    defaults = {
        el.get("Extension").lower()
        for el in root.iter(f"{CT_NS}Default")
    }
    overrides = {el.get("PartName") for el in root.iter(f"{CT_NS}Override")}
    for member in names:
        if member == "[Content_Types].xml" or member.endswith("/"):
            continue
        extension = member.rsplit(".", 1)[-1].lower() if "." in member else ""
        covered = f"/{member}" in overrides or extension in defaults
        assert covered, f"part /{member} has no [Content_Types].xml entry"


def creation_id_values(pptx_path) -> list:
    """All a16:creationId id values across every slide of a saved deck."""
    values = []
    for name in slide_part_names(pptx_path):
        root = ET.fromstring(read_zip(pptx_path, name))
        for el in root.iter(A16_CREATION_ID):
            values.append(el.get("id"))
    return values


def save_to(prs, directory, filename: str) -> Path:
    path = Path(directory) / filename
    prs.save(str(path))
    return path
