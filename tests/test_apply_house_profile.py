"""End-to-end tests: house-profile/1 apply on the deviant corpus deck.

A hand-written house-profile/1 dict (coded strictly to the pinned
schema contract, values from ``corpus_truth.json``) is applied to a
copy of ``deviant_01.pptx``. Assertions cover, per
``deviations.json``:

* every DETERMINISTIC violation is corrected (v1 size, v3 bullet char,
  v4 hardcoded srgb -> schemeClr, v6 space_after, v7 border weight,
  v8 off-palette fill, v9 font name) -- verified by re-running the
  RESOLVED analysis on the saved deck;
* GEOMETRY violations (v2 off-grid panel, v5 footer straggler) are
  untouched -- lint territory, never apply territory;
* slideMaster / slideLayout / theme parts are byte-identical;
* already-conformant content carries no new explicit overrides
  (element-level XML snapshots + a second-apply no-op check);
* the legacy flat-profile path of ``apply_style_profile`` still works.
"""

import copy
import json
import shutil
from types import SimpleNamespace

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import fixture_path
from utils.resolve_analysis import build_resolved_analysis
from utils.style_apply import (
    apply_house_profile,
    is_house_profile,
    leaf_points,
    leaf_value,
)

DEVIANT = "house_corpus/deviant_01.pptx"

EM_DASH = "—"
MIDDLE_DOT = "·"

#: Hand-written house-profile/1 dict -- pinned schema contract, values
#: from tests/fixtures/house_corpus/corpus_truth.json. Deliberately
#: independent of the profile-builder implementation.
HOUSE_PROFILE = {
    "schema_version": "house-profile/1",
    "name": "meridian_house_test",
    "source_decks": [f"house_{i:02d}.pptx" for i in range(1, 6)],
    "slide_size": {
        "width": {"value": 13.3333, "unit": "in"},
        "height": {"value": 7.5, "unit": "in"},
    },
    "typography": {
        "title": {
            "font": {"value": "Georgia"},
            "size": {"value": 30.0, "unit": "pt"},
            "bold": {"value": False},
            "color": {"value": "#14324F"},
        },
        "body": {
            "font": {"value": "Calibri"},
            "size": {"value": 14.0, "unit": "pt"},
            "color": {"value": "#20262B"},
        },
        "footer": {
            "font": {"value": "Calibri"},
            "size": {"value": 11.0, "unit": "pt"},
            "color": {"value": "#3E5C76"},
        },
    },
    "paragraph": {
        "space_before": {"value": 2.0, "unit": "pt"},
        "space_after": {"value": 8.0, "unit": "pt"},
        "line_spacing": {"value": 1.2, "unit": None},
        "bullets": {
            "l1": {"char": {"value": EM_DASH},
                   "size_pct": {"value": 95, "unit": None}},
            "l2": {"char": {"value": MIDDLE_DOT},
                   "size_pct": {"value": 90, "unit": None}},
            "l3": {"char": {"value": ">"},
                   "size_pct": {"value": 90, "unit": None}},
        },
    },
    "palette": {
        "scheme": {
            "dk1": {"value": "#20262B"}, "lt1": {"value": "#FAF9F6"},
            "dk2": {"value": "#14324F"}, "lt2": {"value": "#DCE3E8"},
            "accent1": {"value": "#1B7F79"},
            "accent2": {"value": "#D97C2B"},
            "accent3": {"value": "#3E5C76"},
            "accent4": {"value": "#8A5A83"},
            "accent5": {"value": "#5E8C61"},
            "accent6": {"value": "#B23A48"},
            "hlink": {"value": "#176B87"},
            "folHlink": {"value": "#6D5D8F"},
        },
        "usage": [
            {"color": {"value": "#20262B"}, "share": 0.42, "role": "text"},
            {"color": {"value": "#DCE3E8"}, "share": 0.21, "role": "fill"},
        ],
    },
    "shape_defaults": {
        "border": {
            "weight": {"value": 1.25, "unit": "pt"},
            "color": {"value": "#14324F"},
            "dash": {"value": "dash"},
        },
        "corner_radius": {"value": 0.12, "unit": None},
        "fill": {"value": "#DCE3E8"},
    },
    # Analysis-side sections: present to prove apply IGNORES them.
    "grid": {
        "edges": {
            "left": [{"value": 0.8333, "unit": "in"},
                     {"value": 5.0, "unit": "in"},
                     {"value": 9.1667, "unit": "in"}],
            "right": [{"value": 4.1667, "unit": "in"},
                      {"value": 8.3333, "unit": "in"},
                      {"value": 12.5, "unit": "in"}],
            "center": [{"value": 2.5, "unit": "in"},
                       {"value": 6.6667, "unit": "in"},
                       {"value": 10.8333, "unit": "in"}],
        },
        "tolerance": {"value": 0.0556, "unit": "in"},
    },
    "distributions": {
        "font_sizes_pt": {"values": [11.0, 14.0, 20.0, 30.0],
                          "shares": [0.25, 0.45, 0.15, 0.15]},
    },
}

#: Part-name prefixes apply must NEVER touch.
PROTECTED_PREFIXES = ("/ppt/slideMasters/", "/ppt/slideLayouts/",
                      "/ppt/theme/")

#: (slide_index, shape_name) elements that are fully conformant and must
#: come out of apply XML-identical (incl. the geometry-violation shapes,
#: which are lint territory).
CONFORMANT_SNAPSHOTS = (
    (1, "Title 1"),
    (1, "TakeawayPanel"),
    (1, "FooterNote"),
    (2, "BodyContent"),
    (2, "OffGridPanel"),
    (4, "StragglerNote"),
    (4, "ContactPanel"),
)


def _shape_named(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise AssertionError(f"no shape named {name!r}")


def _protected_blobs(prs):
    return {
        str(part.partname): part.blob
        for part in prs.part.package.iter_parts()
        if str(part.partname).startswith(PROTECTED_PREFIXES)
    }


def _snapshot_elements(prs):
    return {
        (index, name): etree.tostring(_shape_named(prs.slides[index],
                                                   name)._element)
        for index, name in CONFORMANT_SNAPSHOTS
    }


def _slide_record(report, slide_number):
    for record in report["slides"]:
        if record["slide_number"] == slide_number:
            return record
    raise AssertionError(f"slide {slide_number} missing from report")


def _shape_record(report, slide_number, name):
    for shape in _slide_record(report, slide_number)["shapes"]:
        if shape["name"] == name:
            return shape
    raise AssertionError(f"shape {name!r} missing from slide "
                         f"{slide_number} report")


@pytest.fixture(scope="module")
def applied(tmp_path_factory):
    """Apply the profile once; share the outcome across this module."""
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    workdir = tmp_path_factory.mktemp("house_apply")
    copy_path = workdir / "deviant_copy.pptx"
    shutil.copyfile(source, copy_path)

    prs = Presentation(str(copy_path))
    protected_before = _protected_blobs(prs)
    snapshots_before = _snapshot_elements(prs)
    summary = apply_house_profile(prs, HOUSE_PROFILE)
    protected_after = _protected_blobs(prs)
    snapshots_after = _snapshot_elements(prs)

    saved_path = workdir / "applied.pptx"
    prs.save(str(saved_path))
    reopened = Presentation(str(saved_path))
    report = build_resolved_analysis(reopened, detail="full",
                                     max_bytes=5_000_000)
    assert report["truncated"] is False
    return SimpleNamespace(
        summary=summary,
        protected_before=protected_before,
        protected_after=protected_after,
        snapshots_before=snapshots_before,
        snapshots_after=snapshots_after,
        saved_path=saved_path,
        reopened=reopened,
        report=report,
    )


# ---------------------------------------------------------------------------
# Deterministic violations corrected (re-run resolved analysis)
# ---------------------------------------------------------------------------

def test_v1_offscale_font_size_snapped_to_scale(applied):
    body = _shape_record(applied.report, 2, "BodyContent")
    run = body["paragraphs"][1]["runs"][0]
    assert run["font"]["size_pt"] == 14.0


def test_v9_wrong_font_reset_to_house_body_font(applied):
    body = _shape_record(applied.report, 2, "BodyContent")
    run = body["paragraphs"][2]["runs"][0]
    assert run["font"]["name"] == "Calibri"


def test_v3_wrong_bullet_char_reset_to_em_dash(applied):
    body = _shape_record(applied.report, 2, "BodyContent")
    bullet = body["paragraphs"][3]["bullet"]
    assert bullet["type"] == "char"
    assert bullet["char"] == EM_DASH
    # Resolver convention: buSzPct as a fraction (0.95 == 95%). The
    # size was already conformant (inherited) and must stay 95%.
    assert bullet["size_pct"] == 0.95


def test_v6_wrong_space_after_reset(applied):
    body = _shape_record(applied.report, 2, "BodyContent")
    paragraph = body["paragraphs"][0]
    assert paragraph["space_after"] == {"points": 8.0}
    # Untouched siblings of the fix:
    assert paragraph["space_before"] == {"points": 2.0}
    assert paragraph["line_spacing"] == {"multiple": 1.2}


def test_v7_border_weight_reset(applied):
    panel = _shape_record(applied.report, 3, "TakeawayPanel")
    assert panel["line"]["weight_pt"] == 1.25
    assert panel["line"]["dash"] == "dash"
    assert panel["line"]["color_hex"] == "14324F"


def test_v8_off_palette_fill_snapped_to_house_fill(applied):
    panel = _shape_record(applied.report, 4, "ColumnPanelRight")
    assert panel["fill"] == {"type": "solid", "color_hex": "DCE3E8"}


def test_v4_hardcoded_srgb_relinked_to_scheme(applied):
    shape = _shape_named(applied.reopened.slides[3], "ColumnPanelLeft")
    run = shape.text_frame.paragraphs[0].runs[0]
    fill = run._r.find(qn("a:rPr")).find(qn("a:solidFill"))
    assert fill.find(qn("a:srgbClr")) is None
    assert fill.find(qn("a:schemeClr")).get("val") == "dk1"
    # Visual no-op: the resolved color is unchanged.
    record = _shape_record(applied.report, 4, "ColumnPanelLeft")
    assert record["paragraphs"][0]["runs"][0]["font"]["color_hex"] == "20262B"


# ---------------------------------------------------------------------------
# Geometry violations are lint territory -- never applied
# ---------------------------------------------------------------------------

def test_v2_off_grid_panel_geometry_untouched(applied):
    panel = _shape_record(applied.report, 3, "OffGridPanel")
    assert panel["geometry"]["left_pt"] == 682.0


def test_v5_footer_straggler_untouched(applied):
    straggler = _shape_record(applied.report, 5, "StragglerNote")
    assert straggler["geometry"]["left_pt"] == 400.0
    assert straggler["geometry"]["top_pt"] == 505.0


# ---------------------------------------------------------------------------
# Hard rule: masters / layouts / theme byte-identical
# ---------------------------------------------------------------------------

def test_masters_layouts_theme_byte_identical(applied):
    assert applied.protected_before, "no protected parts found -- bad test"
    assert set(applied.protected_before) == set(applied.protected_after)
    for partname, blob in applied.protected_before.items():
        assert applied.protected_after[partname] == blob, (
            f"apply mutated protected part {partname}")


# ---------------------------------------------------------------------------
# Minimal-diff: conformant content gets no new explicit overrides
# ---------------------------------------------------------------------------

def test_conformant_elements_xml_identical(applied):
    for key, before in applied.snapshots_before.items():
        assert applied.snapshots_after[key] == before, (
            f"apply modified conformant shape {key}")


def test_no_override_added_to_inherited_title_run(applied):
    title = _shape_named(applied.reopened.slides[1], "Title 1")
    r_pr = title.text_frame.paragraphs[0].runs[0]._r.find(qn("a:rPr"))
    assert r_pr.get("sz") is None
    assert r_pr.get("b") is None
    assert r_pr.find(qn("a:latin")) is None
    assert r_pr.find(qn("a:solidFill")) is None


def test_second_apply_is_a_noop(applied):
    prs = Presentation(str(applied.saved_path))
    summary = apply_house_profile(prs, HOUSE_PROFILE)
    assert summary["writes"] == 0
    assert summary["slides_touched"] == []
    assert summary["changes"] == []


def test_summary_reports_exactly_the_seven_deterministic_fixes(applied):
    summary = applied.summary
    assert summary["writes"] == 7
    assert summary["changes_truncated"] is False
    fixed = {(change["slide"], change["shape"], change["property"])
             for change in summary["changes"]}
    assert fixed == {
        (2, "BodyContent", "font.size_pt"),
        (2, "BodyContent", "font.name"),
        (2, "BodyContent", "bullet.char"),
        (2, "BodyContent", "paragraph.space_after"),
        (3, "TakeawayPanel", "line.weight_pt"),
        (4, "ColumnPanelLeft", "font.color_source"),
        (4, "ColumnPanelRight", "fill.color"),
    }
    assert summary["slides_touched"] == [2, 3, 4]


# ---------------------------------------------------------------------------
# Schema helpers
# ---------------------------------------------------------------------------

def test_is_house_profile():
    assert is_house_profile(HOUSE_PROFILE)
    assert not is_house_profile({"schema_version": "house-profile/2"})
    assert not is_house_profile("meridian")
    assert not is_house_profile(None)


def test_leaf_readers():
    assert leaf_value({"value": 14.0, "unit": "pt"}) == 14.0
    assert leaf_value("Calibri") == "Calibri"
    assert leaf_points({"value": 1.0, "unit": "in"}) == 72.0
    assert leaf_points({"value": 8.0, "unit": "pt"}) == 8.0
    assert leaf_points(8.0) == 8.0
    assert leaf_points(None) is None
    with pytest.raises(ValueError):
        leaf_value({"unit": "pt"})
    with pytest.raises(ValueError):
        leaf_points({"value": 8.0, "unit": "em"})
    with pytest.raises(ValueError):
        leaf_points({"value": "eight", "unit": "pt"})


def test_apply_rejects_non_house_profiles():
    with pytest.raises(ValueError):
        apply_house_profile(object(), {"name": "legacy-ish"})
    with pytest.raises(ValueError):
        apply_house_profile(None, HOUSE_PROFILE)
    bad = copy.deepcopy(HOUSE_PROFILE)
    bad["typography"] = ["not", "a", "dict"]
    with pytest.raises(ValueError):
        apply_house_profile(object(), bad)


# ---------------------------------------------------------------------------
# MCP tool layer: house path + legacy regression
# ---------------------------------------------------------------------------

class _RecorderApp:
    """Minimal FastMCP stand-in capturing registered tool functions."""

    def __init__(self):
        self.tools = {}

    def tool(self, *args, **kwargs):
        def decorator(fn):
            self.tools[fn.__name__] = fn
            return fn

        return decorator


def _style_tools(presentations, current_id):
    from tools.style_tools import register_style_tools

    app = _RecorderApp()
    register_style_tools(app, presentations, lambda: current_id)
    return app.tools


@pytest.fixture
def deviant_copy(tmp_path):
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    copy_path = tmp_path / "deviant.pptx"
    shutil.copyfile(source, copy_path)
    return copy_path


def test_tool_load_and_apply_house_profile(tmp_path, deviant_copy):
    profile_path = tmp_path / "meridian.json"
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_tool_test"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")

    prs = Presentation(str(deviant_copy))
    tools = _style_tools({"deck": prs}, "deck")

    loaded = tools["load_style_profile"](file_path=str(profile_path))
    assert "error" not in loaded
    assert loaded["profile_name"] == "meridian_tool_test"
    assert loaded["schema_version"] == "house-profile/1"

    result = tools["apply_style_profile"](profile_name="meridian_tool_test")
    assert "error" not in result
    assert result["writes"] == 7
    assert result["slides_touched"] == [2, 3, 4]
    assert "message" in result


def test_tool_apply_house_profile_requires_presentation(tmp_path):
    profile_path = tmp_path / "meridian.json"
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_no_deck"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    tools = _style_tools({}, None)
    tools["load_style_profile"](file_path=str(profile_path))
    result = tools["apply_style_profile"](profile_name="meridian_no_deck")
    assert "error" in result


def test_tool_load_house_profile_requires_name(tmp_path):
    profile = copy.deepcopy(HOUSE_PROFILE)
    del profile["name"]
    profile_path = tmp_path / "anon.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")
    tools = _style_tools({}, None)
    assert "error" in tools["load_style_profile"](
        file_path=str(profile_path))


def test_tool_apply_house_profile_distinguishes_invalid_input(
        tmp_path, deviant_copy):
    """Regression: the ValueError arm must return a DISTINCT payload
    (it used to duplicate the generic except-Exception message)."""
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_invalid_section"
    profile["typography"] = ["not", "a", "dict"]  # _validate_profile trips
    profile_path = tmp_path / "invalid.json"
    profile_path.write_text(json.dumps(profile), encoding="utf-8")

    prs = Presentation(str(deviant_copy))
    tools = _style_tools({"deck": prs}, "deck")
    loaded = tools["load_style_profile"](file_path=str(profile_path))
    assert "error" not in loaded
    result = tools["apply_style_profile"](
        profile_name="meridian_invalid_section")
    assert "error" in result
    assert result["error"].startswith(
        "House-profile apply rejected invalid input")


def test_apply_style_profile_body_within_house_line_limit():
    """Regression: the tool body grew past the <50-line house rule when
    the house branch landed; the legacy loop now lives in a helper."""
    import ast
    import inspect
    import textwrap

    tools = _style_tools({}, None)
    source = textwrap.dedent(inspect.getsource(tools["apply_style_profile"]))
    function_def = ast.parse(source).body[0]
    body = function_def.body
    if (isinstance(body[0], ast.Expr)
            and isinstance(body[0].value, ast.Constant)
            and isinstance(body[0].value.value, str)):
        body = body[1:]  # the docstring does not count against the rule
    statement_lines = body[-1].end_lineno - body[0].lineno + 1
    assert statement_lines < 50, (
        f"apply_style_profile body is {statement_lines} lines; house "
        "rule caps functions at <50 -- extract a helper")


def test_tool_legacy_profile_path_unchanged(deviant_copy):
    """The old flat-profile flow must behave exactly as before."""
    prs = Presentation(str(deviant_copy))
    tools = _style_tools({"deck": prs}, "deck")

    created = tools["create_style_profile"](
        file_path=str(deviant_copy), profile_name="legacy_regression")
    assert "error" not in created
    assert created["profile_name"] == "legacy_regression"

    result = tools["apply_style_profile"](profile_name="legacy_regression")
    assert "error" not in result
    assert result["font_applied"] == created["primary_font"]
    assert result["runs_modified"] > 0
    assert "title_size" in result and "body_size" in result
