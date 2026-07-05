"""Tests for ``utils/clone_utils.py`` -- the slide clone / rel-rewriting engine.

Covers the python-pptx issue-#132 pitfall matrix:
    * duplicate_slide on every slide of every Step 0 fixture -- shape counts,
      text, geometry, and cSld XML match the source exactly; all rels valid,
    * copy_slide between two copies of the same fixture -- XML equivalence,
      notesSlide exclusion, no duplicate a16:creationId, ph idx matches layout,
    * copy_slide with no matching layout name -> ValueError,
    * layout name match with mismatched p:ph idx set -> ValueError
      (issue-#132 false-positive pitfall),
    * chart-bearing slide -> NotImplementedError naming the shape,
    * unhandled-rel final sweep -> leaf target part copied and
      [Content_Types].xml covers every part in the saved package; non-leaf
      sweep targets (own relationships) -> NotImplementedError,
    * movie slides (empty-rId action hyperlink + RT.MEDIA/RT.VIDEO rels)
      survive both duplicate_slide and copy_slide,
    * internal slide-jump hyperlinks: kept by duplicate_slide, dropped
      (text kept, no phantom slide part) by copy_slide.
"""

import shutil
import zipfile

import pytest
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from tests.conftest import FIXTURE_DECKS, fixture_path, skip_if_fixture_missing
from tests.support_clone import (
    assert_content_types_cover_all_parts,
    assert_slide_rels_valid,
    canonical_cSld,
    creation_id_values,
    make_chart_deck,
    make_hyperlink_deck,
    make_image_deck,
    make_movie_deck,
    make_slide_jump_deck,
    make_two_image_deck,
    read_zip,
    rename_all_layouts,
    save_to,
    slide_part_names,
    slide_signature,
)
from utils.clone_utils import copy_slide, duplicate_slide


def _deck_params():
    return [
        pytest.param(deck, marks=skip_if_fixture_missing(deck), id=deck)
        for deck in FIXTURE_DECKS
    ]


def _slide_xml_by_position(pptx_path, position: int) -> bytes:
    """Slide XML of the slide at 0-based ``position`` in presentation order."""
    prs = Presentation(str(pptx_path))
    return prs.slides[position].part.blob


# ------------------------------------------------------------- duplicate_slide


@pytest.mark.parametrize("deck_name", _deck_params())
def test_duplicate_every_slide_roundtrip(deck_name, open_deck, tmp_path):
    prs = open_deck(deck_name)
    source_count = len(prs.slides)
    source_signatures = [slide_signature(s) for s in prs.slides]

    for index in range(source_count):
        new_slide = duplicate_slide(prs, index)
        assert new_slide is not None

    path = save_to(prs, tmp_path, f"dup_{deck_name}")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == source_count * 2

    for index in range(source_count):
        duplicate = reopened.slides[source_count + index]
        assert slide_signature(duplicate) == source_signatures[index], (
            f"{deck_name} slide {index}: duplicate signature differs from source"
        )
        source_xml = reopened.slides[index].part.blob
        dup_xml = duplicate.part.blob
        assert canonical_cSld(dup_xml) == canonical_cSld(source_xml), (
            f"{deck_name} slide {index}: duplicate cSld XML differs from source"
        )

    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


@pytest.mark.parametrize("deck_name", _deck_params())
def test_duplicate_slide_keeps_layout(deck_name, open_deck, tmp_path):
    prs = open_deck(deck_name)
    source_layouts = [s.slide_layout.name for s in prs.slides]
    count = len(prs.slides)
    for index in range(count):
        duplicate_slide(prs, index)

    path = save_to(prs, tmp_path, f"dup_layout_{deck_name}")
    reopened = Presentation(str(path))
    duplicate_layouts = [
        s.slide_layout.name for s in list(reopened.slides)[count:]
    ]
    assert duplicate_layouts == source_layouts


@pytest.mark.parametrize("deck_name", _deck_params())
def test_duplicate_slide_no_duplicate_creation_ids(deck_name, open_deck, tmp_path):
    prs = open_deck(deck_name)
    duplicate_slide(prs, 0)
    path = save_to(prs, tmp_path, f"dup_ids_{deck_name}")
    ids = creation_id_values(path)
    assert len(ids) == len(set(ids)), f"duplicate a16:creationId values: {ids}"


def test_duplicate_slide_invalid_index_raises():
    prs = make_image_deck()
    with pytest.raises(ValueError, match="[Ss]lide index"):
        duplicate_slide(prs, 5)
    with pytest.raises(ValueError, match="[Ss]lide index"):
        duplicate_slide(prs, -1)


def test_duplicate_slide_shares_image_part(tmp_path):
    prs = make_image_deck()
    duplicate_slide(prs, 0)
    path = save_to(prs, tmp_path, "dup_image.pptx")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == 2
    media = [
        n for n in __import__("zipfile").ZipFile(path).namelist()
        if n.startswith("ppt/media/")
    ]
    assert len(media) == 1, "same-deck duplicate must share, not re-add, images"
    assert_slide_rels_valid(path)


def test_duplicate_chart_slide_rejected():
    prs = make_chart_deck()
    chart_shape_names = [s.name for s in prs.slides[0].shapes if s.has_chart]
    with pytest.raises(NotImplementedError) as excinfo:
        duplicate_slide(prs, 0)
    assert chart_shape_names[0] in str(excinfo.value)


# ------------------------------------------------------------------ copy_slide


@pytest.mark.parametrize("deck_name", _deck_params())
def test_copy_slide_between_copies_of_same_fixture(deck_name, tmp_path):
    """Copy each slide from a fixture into a second copy of the same fixture."""
    src_prs = Presentation(str(fixture_path(deck_name)))
    dst_file = tmp_path / f"dst_{deck_name}"
    shutil.copyfile(fixture_path(deck_name), dst_file)
    dst_prs = Presentation(str(dst_file))
    original_count = len(dst_prs.slides)

    for index in range(len(src_prs.slides)):
        copy_slide(src_prs, index, dst_prs)

    path = save_to(dst_prs, tmp_path, f"copied_{deck_name}")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == original_count + len(src_prs.slides)

    for index in range(len(src_prs.slides)):
        copied = reopened.slides[original_count + index]
        source_xml = src_prs.slides[index].part.blob
        assert canonical_cSld(copied.part.blob) == canonical_cSld(source_xml), (
            f"{deck_name} slide {index}: copied cSld XML differs from source"
        )
        assert copied.slide_layout.name == src_prs.slides[index].slide_layout.name

    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


@pytest.mark.parametrize("deck_name", _deck_params())
def test_copy_slide_excludes_notes_and_creation_ids(deck_name, tmp_path):
    src_prs = Presentation(str(fixture_path(deck_name)))
    dst_file = tmp_path / f"dst_{deck_name}"
    shutil.copyfile(fixture_path(deck_name), dst_file)
    dst_prs = Presentation(str(dst_file))

    new_slide = copy_slide(src_prs, 0, dst_prs)
    notes_rels = [
        rel for rel in new_slide.part.rels.values()
        if rel.reltype.endswith("/notesSlide")
    ]
    assert not notes_rels, "copied slide must not carry a notesSlide rel"

    path = save_to(dst_prs, tmp_path, f"copied_ids_{deck_name}")
    ids = creation_id_values(path)
    assert len(ids) == len(set(ids)), f"duplicate a16:creationId values: {ids}"


@pytest.mark.parametrize("deck_name", _deck_params())
def test_copy_slide_ph_idx_all_match_layout(deck_name, tmp_path):
    src_prs = Presentation(str(fixture_path(deck_name)))
    dst_file = tmp_path / f"dst_{deck_name}"
    shutil.copyfile(fixture_path(deck_name), dst_file)
    dst_prs = Presentation(str(dst_file))

    for index in range(len(src_prs.slides)):
        copy_slide(src_prs, index, dst_prs)

    path = save_to(dst_prs, tmp_path, f"copied_ph_{deck_name}")
    reopened = Presentation(str(path))
    for slide in reopened.slides:
        layout_idxs = {
            ph.placeholder_format.idx for ph in slide.slide_layout.placeholders
        }
        for ph in slide.placeholders:
            assert ph.placeholder_format.idx in layout_idxs, (
                f"ph idx {ph.placeholder_format.idx} on slide has no matching "
                f"placeholder in layout {slide.slide_layout.name!r}"
            )


def test_copy_slide_no_matching_layout_raises_value_error():
    src_prs = Presentation(str(fixture_path(FIXTURE_DECKS[0])))
    dst_prs = Presentation()  # default template
    rename_all_layouts(dst_prs)
    with pytest.raises(ValueError, match="layout"):
        copy_slide(src_prs, 0, dst_prs)


def test_copy_chart_slide_rejected():
    src_prs = make_chart_deck()
    dst_prs = Presentation()
    chart_shape_names = [s.name for s in src_prs.slides[0].shapes if s.has_chart]
    with pytest.raises(NotImplementedError) as excinfo:
        copy_slide(src_prs, 0, dst_prs)
    assert chart_shape_names[0] in str(excinfo.value)


def test_copy_slide_invalid_index_raises():
    src_prs = make_image_deck()
    dst_prs = Presentation()
    with pytest.raises(ValueError, match="[Ss]lide index"):
        copy_slide(src_prs, 99, dst_prs)


def test_copy_slide_remaps_image_and_dedupes(tmp_path):
    src_prs = make_image_deck()
    dst_prs = Presentation()
    copy_slide(src_prs, 0, dst_prs)
    copy_slide(src_prs, 0, dst_prs)  # second copy must reuse the image part

    path = save_to(dst_prs, tmp_path, "copied_image.pptx")
    import zipfile

    media = [
        n for n in zipfile.ZipFile(path).namelist() if n.startswith("ppt/media/")
    ]
    assert len(media) == 1, "identical images must be deduped by hash"
    reopened = Presentation(str(path))
    pictures = [s for s in reopened.slides[0].shapes if s.shape_type == 13]
    assert pictures, "copied slide lost its picture"
    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


def test_copy_slide_two_distinct_images_get_distinct_parts(tmp_path):
    """Two different images on one slide must not collide on a partname."""
    import zipfile

    src_prs = make_two_image_deck()
    dst_prs = Presentation()
    copy_slide(src_prs, 0, dst_prs)

    path = save_to(dst_prs, tmp_path, "copied_two_images.pptx")
    with zipfile.ZipFile(path) as archive:
        media = [n for n in archive.namelist() if n.startswith("ppt/media/")]
    assert len(media) == 2, f"expected 2 distinct media parts, got {media}"
    assert len(set(media)) == 2, f"duplicate partnames in package: {media}"

    reopened = Presentation(str(path))
    blobs = {
        shape.image.blob
        for shape in reopened.slides[0].shapes
        if shape.shape_type == 13
    }
    assert len(blobs) == 2, "the two copied pictures lost their distinct images"
    assert_slide_rels_valid(path)


def test_copy_slide_preserves_external_hyperlink(tmp_path):
    url = "https://example.com/step1"
    src_prs = make_hyperlink_deck(url)
    dst_prs = Presentation()
    copy_slide(src_prs, 0, dst_prs)

    path = save_to(dst_prs, tmp_path, "copied_hyperlink.pptx")
    reopened = Presentation(str(path))
    addresses = [
        run.hyperlink.address
        for shape in reopened.slides[0].shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    ]
    assert url in addresses
    assert_slide_rels_valid(path)


def test_copy_slide_unhandled_rel_sweep(tmp_path):
    """An unhandled rel type must be swept: target part copied, types patched."""
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    deck_name = FIXTURE_DECKS[0]
    src_prs = Presentation(str(fixture_path(deck_name)))
    dst_file = tmp_path / f"dst_{deck_name}"
    shutil.copyfile(fixture_path(deck_name), dst_file)
    dst_prs = Presentation(str(dst_file))

    tag_blob = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<p:tagLst xmlns:p='
        b'"http://schemas.openxmlformats.org/presentationml/2006/main"/>'
    )
    tag_part = Part(
        PackURI("/ppt/tags/tag1.xml"),
        "application/vnd.openxmlformats-officedocument.presentationml.tags+xml",
        src_prs.part.package,
        tag_blob,
    )
    tags_reltype = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags"
    )
    src_prs.slides[0].part.relate_to(tag_part, tags_reltype)

    new_slide = copy_slide(src_prs, 0, dst_prs)
    swept = [
        rel for rel in new_slide.part.rels.values() if rel.reltype == tags_reltype
    ]
    assert swept, "unhandled rel was not recreated on the copied slide"
    assert swept[0].target_part.blob == tag_blob

    path = save_to(dst_prs, tmp_path, "copied_sweep.pptx")
    tag_members = [
        n for n in __import__("zipfile").ZipFile(path).namelist()
        if n.startswith("ppt/tags/")
    ]
    assert tag_members, "swept part missing from saved package"
    assert_content_types_cover_all_parts(path)
    assert_slide_rels_valid(path)


def test_copy_slide_input_validation():
    src_prs = make_image_deck()
    with pytest.raises(ValueError, match="[Pp]resentation"):
        copy_slide(src_prs, 0, "not a presentation")
    with pytest.raises(ValueError, match="[Pp]resentation"):
        copy_slide("not a presentation", 0, src_prs)


# ------------------------------------------------------------- media (movies)


def test_duplicate_slide_with_movie(tmp_path):
    """Regression: movie shapes carry <a:hlinkClick r:id=""> (action-only,
    empty rId) which must not be treated as a dangling relationship."""
    prs = make_movie_deck()
    baseline = save_to(prs, tmp_path, "movie_baseline.pptx")
    with zipfile.ZipFile(baseline) as archive:
        baseline_media = sorted(
            n for n in archive.namelist() if n.startswith("ppt/media/")
        )

    duplicate_slide(prs, 0)

    path = save_to(prs, tmp_path, "dup_movie.pptx")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == 2
    with zipfile.ZipFile(path) as archive:
        media = sorted(n for n in archive.namelist() if n.startswith("ppt/media/"))
    assert media == baseline_media, (
        "same-deck duplicate must share, not re-add, media parts"
    )
    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


def test_copy_slide_with_movie_recreates_media(tmp_path):
    """Regression: empty-rId action hyperlink must not abort the copy, and
    the _copy_media_rel handler must re-create the video part."""
    src_prs = make_movie_deck()
    dst_prs = Presentation()
    new_slide = copy_slide(src_prs, 0, dst_prs)

    media_rels = [
        rel for rel in new_slide.part.rels.values()
        if rel.reltype in (RT.MEDIA, RT.VIDEO)
    ]
    assert media_rels, "copied movie slide lost its media rels"

    path = save_to(dst_prs, tmp_path, "copied_movie.pptx")
    with zipfile.ZipFile(path) as archive:
        media = [n for n in archive.namelist() if n.startswith("ppt/media/")]
    videos = [n for n in media if n.endswith(".mp4")]
    assert videos, f"video part missing from destination package: {media}"
    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


# ------------------------------------------------- internal slide-jump links


def test_duplicate_slide_keeps_internal_slide_jump_link(tmp_path):
    """Same-deck duplicate shares the jump-target slide, so the link stays."""
    prs = make_slide_jump_deck()
    duplicate_slide(prs, 0)

    path = save_to(prs, tmp_path, "dup_jump.pptx")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == 3
    duplicate = reopened.slides[2]
    jump_rels = [
        rel for rel in duplicate.part.rels.values() if rel.reltype == RT.SLIDE
    ]
    assert jump_rels, "duplicate lost its internal slide-jump rel"
    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


def test_copy_slide_drops_internal_slide_jump_link(tmp_path):
    """Regression: cross-deck copy of a slide-jump link must drop the link
    (keeping the text) instead of blob-copying the target slide as a phantom
    part with dangling refs (the PowerPoint repair-prompt failure)."""
    src_prs = make_slide_jump_deck()
    dst_prs = Presentation()
    new_slide = copy_slide(src_prs, 0, dst_prs)

    jump_rels = [
        rel for rel in new_slide.part.rels.values() if rel.reltype == RT.SLIDE
    ]
    assert not jump_rels, "internal slide-jump rel must be dropped on copy"

    path = save_to(dst_prs, tmp_path, "copied_jump.pptx")
    reopened = Presentation(str(path))
    assert len(reopened.slides) == 1, "destination must gain exactly one slide"
    assert len(slide_part_names(path)) == 1, (
        "phantom slide part in package (slide part outside sldIdLst)"
    )
    texts = [
        shape.text_frame.text
        for shape in reopened.slides[0].shapes
        if shape.has_text_frame
    ]
    assert "jump to slide 2" in texts, "link text must be kept"
    assert_slide_rels_valid(path)
    assert_content_types_cover_all_parts(path)


# --------------------------------------------- layout-match ph idx validation


def test_copy_slide_layout_name_false_positive_rejected():
    """Issue-#132 pitfall: a destination layout matching by *name* but
    lacking the slide's p:ph idx values must be rejected, not silently
    bound (placeholder inheritance would fall back to master defaults)."""
    src_prs = Presentation()
    title_layout = src_prs.slide_masters[0].slide_layouts[0]  # "Title Slide"
    src_prs.slides.add_slide(title_layout)  # placeholders idx {0, 1}

    dst_prs = Presentation()
    dst_layouts = [
        layout
        for master in dst_prs.slide_masters
        for layout in master.slide_layouts
    ]
    # Rename the genuine match away, then masquerade the placeholder-less
    # Blank layout under the source layout's name.
    for layout in dst_layouts:
        if layout.name == title_layout.name:
            layout.name = "renamed away"
    blank = next(l for l in dst_layouts if l.name == "Blank")
    blank.name = title_layout.name

    with pytest.raises(ValueError, match="template lineage"):
        copy_slide(src_prs, 0, dst_prs)
    assert len(dst_prs.slides) == 0, "failed copy must leave destination untouched"


# ------------------------------------------------------- non-leaf sweep guard


def test_copy_slide_sweep_refuses_target_with_child_rels():
    """The final sweep blob-copies leaf parts only; a swept target carrying
    its own relationships must be refused (a shallow copy would strand its
    internal r:id references)."""
    from pptx.opc.package import Part
    from pptx.opc.packuri import PackURI

    src_prs = make_image_deck()
    dst_prs = Presentation()

    tags_ct = (
        "application/vnd.openxmlformats-officedocument.presentationml.tags+xml"
    )
    tags_reltype = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags"
    )
    child = Part(
        PackURI("/ppt/tags/tag8.xml"), tags_ct, src_prs.part.package, b"<child/>"
    )
    parent = Part(
        PackURI("/ppt/tags/tag9.xml"), tags_ct, src_prs.part.package, b"<parent/>"
    )
    parent.relate_to(child, tags_reltype)
    src_prs.slides[0].part.relate_to(parent, tags_reltype)

    with pytest.raises(NotImplementedError, match="relationships"):
        copy_slide(src_prs, 0, dst_prs)
    assert len(dst_prs.slides) == 0, "failed copy must leave destination untouched"
