"""Unit tests for the Stage 2 property layer (``utils.resolve_utils`` +
``resolve_text_props`` + ``resolve_shape_props``).

The exhaustive fixture walk lives in ``test_resolver_fixtures.py``; these
tests cover behavior the COM-authored fixtures cannot exercise --
normAutofit fontScale scaling, bullet masking, style-reference ordering
and indexing, line attribute merging -- by editing an in-memory copy of a
fixture deck (never saved back) before resolving, plus input validation.
"""

import pytest
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import fixture_path, skip_if_fixture_missing
from utils.resolve_shape_props import (
    resolve_shape_fill,
    resolve_shape_geometry,
    resolve_shape_line,
)
from utils.resolve_text_props import (
    autofit_font_scale,
    scaled_size_pt,
    split_deck_defaults,
    xml_bool,
)
from utils.resolve_theme import ThemeContext
from utils.resolve_utils import (
    TextContext,
    resolve_paragraph,
    resolve_run_font,
    resolve_shape,
)

pytestmark = skip_if_fixture_missing("theme_only.pptx")

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _deck():
    return Presentation(str(fixture_path("theme_only.pptx")))


def _float_shape(prs, name="FloatNoteHubs"):
    slide = prs.slides[3]
    return slide, next(s for s in slide.shapes if s.name == name)


def _add_run_size(run, hundredths):
    rpr = run._r.find(qn("a:rPr"))
    assert rpr is not None
    rpr.set("sz", str(hundredths))


def _set_autofit(shape, font_scale):
    body_pr = shape._element.find(qn("p:txBody") + "/" + qn("a:bodyPr"))
    autofit = body_pr.makeelement(qn("a:normAutofit"), {})
    if font_scale is not None:
        autofit.set("fontScale", font_scale)
    body_pr.append(autofit)


# ---------------------------------------------------------------------------
# normAutofit fontScale (ShapeCrawler PortionFontSize rule)
# ---------------------------------------------------------------------------

class TestNormAutofitScaling:
    def test_scales_run_level_size(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        _add_run_size(run, 2400)
        _set_autofit(shape, "62500")
        font = resolve_run_font(TextContext(slide, shape, paragraph, run))
        assert font["size_pt"] == pytest.approx(24.0 * 0.625)

    def test_percent_string_font_scale(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        _set_autofit(shape, "62.5%")
        assert autofit_font_scale(shape) == pytest.approx(0.625)

    def test_inherited_size_is_not_scaled(self):
        """Sizes from deeper cascade levels are template values --
        ShapeCrawler applies fontScale only to run-rPr sizes."""
        prs = _deck()
        slide, shape = _float_shape(prs)
        _set_autofit(shape, "50000")
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]  # run rPr carries no sz in this fixture
        font = resolve_run_font(TextContext(slide, shape, paragraph, run))
        assert font["size_pt"] == 18.0  # defaultTextStyle size, unscaled

    def test_autofit_without_font_scale_is_ignored(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        _set_autofit(shape, None)
        assert autofit_font_scale(shape) is None

    @pytest.mark.parametrize("bad", [0, -100, "x"])
    def test_scaled_size_pt_validates(self, bad):
        with pytest.raises((ValueError, TypeError)):
            scaled_size_pt(bad, None)

    def test_scaled_size_pt_rejects_bad_scale(self):
        with pytest.raises(ValueError):
            scaled_size_pt(1800, 1.5)


# ---------------------------------------------------------------------------
# Bullet cascade semantics
# ---------------------------------------------------------------------------

class TestBulletCascade:
    def _content_paragraph(self, prs):
        slide = prs.slides[1]
        shape = next(s for s in slide.shapes
                     if s.name == "Content Placeholder 2")
        return slide, shape, shape.text_frame.paragraphs[0]

    def test_bunone_at_paragraph_masks_master_buchar(self):
        prs = _deck()
        slide, shape, paragraph = self._content_paragraph(prs)
        ppr = paragraph._p.get_or_add_pPr()
        ppr.insert(0, ppr.makeelement(qn("a:buNone"), {}))
        resolved = resolve_paragraph(TextContext(slide, shape, paragraph))
        assert resolved["bullet"] == {"type": "none"}

    def test_bullet_subproperties_cascade_independently(self):
        """A paragraph-level buChar override keeps the master's buSzPct
        and buFont (they inherit separately, like PowerPoint)."""
        prs = _deck()
        slide, shape, paragraph = self._content_paragraph(prs)
        ppr = paragraph._p.get_or_add_pPr()
        bu_char = ppr.makeelement(qn("a:buChar"), {"char": "*"})
        ppr.insert(0, bu_char)
        bullet = resolve_paragraph(
            TextContext(slide, shape, paragraph))["bullet"]
        assert bullet["char"] == "*"
        assert bullet["size_pct"] == pytest.approx(0.9)  # master lvl1
        assert bullet["font"] == "Arial"                 # master buFont

    def test_explicit_bullet_color(self):
        prs = _deck()
        slide, shape, paragraph = self._content_paragraph(prs)
        ppr = paragraph._p.get_or_add_pPr()
        bu_clr = ppr.makeelement(qn("a:buClr"), {})
        srgb = bu_clr.makeelement(qn("a:srgbClr"), {"val": "FF0000"})
        bu_clr.append(srgb)
        ppr.insert(0, bu_clr)
        bullet = resolve_paragraph(
            TextContext(slide, shape, paragraph))["bullet"]
        assert bullet["color_follows_text"] is False
        assert bullet["color_hex"] == "FF0000"

    def test_autonum_bullet(self):
        prs = _deck()
        slide, shape, paragraph = self._content_paragraph(prs)
        ppr = paragraph._p.get_or_add_pPr()
        autonum = ppr.makeelement(
            qn("a:buAutoNum"), {"type": "arabicPeriod", "startAt": "3"})
        ppr.insert(0, autonum)
        bullet = resolve_paragraph(
            TextContext(slide, shape, paragraph))["bullet"]
        assert bullet["type"] == "autonum"
        assert bullet["scheme"] == "arabicPeriod"
        assert bullet["start_at"] == 3

    def test_empty_buchar_reports_no_char_code(self):
        """Regression: a malformed ``a:buChar char=""`` used to raise
        IndexError from ``ord(char[0])``; never trust file content."""
        prs = _deck()
        slide, shape, paragraph = self._content_paragraph(prs)
        ppr = paragraph._p.get_or_add_pPr()
        ppr.insert(0, ppr.makeelement(qn("a:buChar"), {"char": ""}))
        bullet = resolve_paragraph(
            TextContext(slide, shape, paragraph))["bullet"]
        assert bullet["type"] == "char"
        assert bullet["char"] is None
        assert bullet["char_code"] is None


# ---------------------------------------------------------------------------
# Shape style references
# ---------------------------------------------------------------------------

class TestStyleReferences:
    def test_own_lst_style_color_beats_font_ref(self):
        """POI order: the shape chain wins over p:style/fontRef."""
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        tx_body = shape._element.find(qn("p:txBody"))
        lst_style = tx_body.find(qn("a:lstStyle"))
        lvl1 = etree.SubElement(lst_style, qn("a:lvl1pPr"))
        def_rpr = etree.SubElement(lvl1, qn("a:defRPr"))
        fill = etree.SubElement(def_rpr, qn("a:solidFill"))
        etree.SubElement(fill, qn("a:srgbClr")).set("val", "112233")
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        font = resolve_run_font(TextContext(slide, shape, paragraph, run))
        assert font["color_hex"] == "112233"

    def test_font_ref_color_beats_deck_defaults(self):
        """Fixture-arbitrated: FloatPanel text renders lt1 (fontRef),
        not the defaultTextStyle tx1 color."""
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        font = resolve_run_font(TextContext(slide, shape, paragraph, run))
        assert font["color_hex"] == theme.scheme_color_hex("lt1")
        assert font["color_hex"] != theme.scheme_color_hex("tx1")

    def test_fill_ref_indexes_theme_fill_styles(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        sp_pr = shape._element.find(qn("p:spPr"))
        sp_pr.remove(sp_pr.find(qn("a:solidFill")))
        fill = resolve_shape_fill(shape, theme)
        # fillRef idx=1 -> fillStyleLst[0] = solidFill phClr; phClr is
        # the fillRef's own accent1 color.
        assert fill["type"] == "solid"
        assert fill["color_hex"] == theme.scheme_color_hex("accent1")

    def test_fill_ref_background_offset(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        sp_pr = shape._element.find(qn("p:spPr"))
        sp_pr.remove(sp_pr.find(qn("a:solidFill")))
        ref = shape._element.find(qn("p:style") + "/" + qn("a:fillRef"))
        ref.set("idx", "1001")
        fill = resolve_shape_fill(shape, theme)
        assert fill["type"] in {"solid", "gradient"}  # bgFillStyleLst[0]

    def test_fill_ref_idx_out_of_range_raises(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        sp_pr = shape._element.find(qn("p:spPr"))
        sp_pr.remove(sp_pr.find(qn("a:solidFill")))
        ref = shape._element.find(qn("p:style") + "/" + qn("a:fillRef"))
        ref.set("idx", "99")
        with pytest.raises(ValueError):
            resolve_shape_fill(shape, theme)

    def test_line_width_merges_from_theme_when_explicit_has_none(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        explicit_ln = shape._element.find(
            qn("p:spPr") + "/" + qn("a:ln"))
        del explicit_ln.attrib["w"]
        theme_ln = theme.line_styles()[1]  # lnRef idx=2
        expected_w = int(theme_ln.get("w")) / 12700.0
        line = resolve_shape_line(shape, theme)
        assert line["visible"] is True
        assert line["weight_pt"] == pytest.approx(expected_w)

    def test_line_from_ln_ref_alone(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        theme = ThemeContext.for_slide(slide)
        sp_pr = shape._element.find(qn("p:spPr"))
        sp_pr.remove(sp_pr.find(qn("a:ln")))
        line = resolve_shape_line(shape, theme)
        assert line["visible"] is True
        assert line["color_hex"] is not None  # phClr = shaded accent1


# ---------------------------------------------------------------------------
# Geometry
# ---------------------------------------------------------------------------

class TestGeometry:
    def test_round_rect_default_adjustment_when_avlst_empty(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelPlan")
        geom = shape._element.find(
            qn("p:spPr") + "/" + qn("a:prstGeom"))
        av_lst = geom.find(qn("a:avLst"))
        for gd in list(av_lst):
            av_lst.remove(gd)
        resolved = resolve_shape_geometry(shape, slide)
        assert resolved["adjustments"] == [pytest.approx(0.16667)]

    def test_unknown_preset_without_avlst_reports_no_adjustments(self):
        prs = _deck()
        slide, shape = _float_shape(prs, "FloatPanelRun")
        resolved = resolve_shape_geometry(shape, slide)
        assert resolved["preset"] == "ellipse"
        assert resolved["adjustments"] == []

    def test_rotation(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        xfrm = shape._element.find(qn("p:spPr") + "/" + qn("a:xfrm"))
        xfrm.set("rot", "2700000")
        resolved = resolve_shape_geometry(shape, slide)
        assert resolved["rotation_deg"] == pytest.approx(45.0)

    def test_resolve_shape_derives_slide_from_part(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        resolved = resolve_shape(shape)
        assert resolved["geometry"]["left_pt"] == pytest.approx(60.0)


# ---------------------------------------------------------------------------
# Malformed-deck attribute handling
# ---------------------------------------------------------------------------

class TestMalformedAttributes:
    """Regression: required integer attributes missing from a malformed
    deck must raise an explicit ValueError, not a bare TypeError from
    ``int(None)`` (house error-handling pattern, cf. resolve_colors)."""

    @staticmethod
    def _ppr(inner_xml: str):
        return etree.fromstring(f'<a:pPr xmlns:a="{_A}">{inner_xml}</a:pPr>')

    def test_spcpts_without_val_raises_value_error(self):
        from utils.resolve_text_props import extract_space_before

        ppr = self._ppr("<a:spcBef><a:spcPts/></a:spcBef>")
        with pytest.raises(ValueError, match="val"):
            extract_space_before(ppr)

    def test_spcpct_without_val_raises_value_error(self):
        from utils.resolve_text_props import extract_line_spacing

        ppr = self._ppr("<a:lnSpc><a:spcPct/></a:lnSpc>")
        with pytest.raises(ValueError, match="val"):
            extract_line_spacing(ppr)

    @pytest.mark.parametrize("inner", ["<a:buSzPct/>", "<a:buSzPts/>"])
    def test_bullet_size_without_val_raises_value_error(self, inner):
        from utils.resolve_text_props import extract_bullet_size

        with pytest.raises(ValueError, match="val"):
            extract_bullet_size(self._ppr(inner))

    def test_offset_without_coordinate_raises_value_error(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        xfrm = shape._element.find(qn("p:spPr") + "/" + qn("a:xfrm"))
        del xfrm.find(qn("a:off")).attrib["x"]
        with pytest.raises(ValueError, match="'x'"):
            resolve_shape_geometry(shape, slide)

    def test_extent_without_size_raises_value_error(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        xfrm = shape._element.find(qn("p:spPr") + "/" + qn("a:xfrm"))
        del xfrm.find(qn("a:ext")).attrib["cy"]
        with pytest.raises(ValueError, match="'cy'"):
            resolve_shape_geometry(shape, slide)

    def test_childless_solidfill_falls_through_cascade(self):
        """Regression: merely READING ``run.font.color`` makes
        python-pptx inject an empty ``a:rPr/a:solidFill``; the resolver
        must treat it as "no color specified" and keep descending, not
        raise "has no color child element"."""
        prs = _deck()
        slide, shape = _float_shape(prs)
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        baseline = resolve_run_font(
            TextContext(slide, shape, paragraph, run))["color_hex"]
        _ = run.font.color.type  # python-pptx read-side mutation
        injected = run._r.find(qn("a:rPr") + "/" + qn("a:solidFill"))
        assert injected is not None and len(injected) == 0
        font = resolve_run_font(TextContext(slide, shape, paragraph, run))
        assert font["color_hex"] == baseline


# ---------------------------------------------------------------------------
# Validation + small helpers
# ---------------------------------------------------------------------------

class TestValidation:
    def test_text_context_requires_slide_and_shape(self):
        with pytest.raises(ValueError):
            TextContext(None, object())
        with pytest.raises(ValueError):
            TextContext(object(), None)

    def test_resolve_run_font_requires_run(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        paragraph = shape.text_frame.paragraphs[0]
        with pytest.raises(ValueError):
            resolve_run_font(TextContext(slide, shape, paragraph, None))

    def test_resolve_paragraph_requires_paragraph(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        with pytest.raises(ValueError):
            resolve_paragraph(TextContext(slide, shape))

    def test_resolve_shape_requires_shape(self):
        with pytest.raises(ValueError):
            resolve_shape(None)

    def test_xml_bool(self):
        assert xml_bool("1") is True
        assert xml_bool("false") is False
        assert xml_bool(None) is None
        with pytest.raises(ValueError):
            xml_bool("maybe")

    def test_split_deck_defaults_partitions_by_name(self):
        prs = _deck()
        slide, shape = _float_shape(prs)
        paragraph = shape.text_frame.paragraphs[0]
        from utils.resolve_core import build_text_cascade_sources

        sources = build_text_cascade_sources(slide, shape, paragraph,
                                             paragraph.runs[0])
        chain, tail = split_deck_defaults(sources)
        assert all(s.name not in {"presentation-defaultTextStyle",
                                  "theme-txDef"} for s in chain)
        assert all(s.name in {"presentation-defaultTextStyle",
                              "theme-txDef"} for s in tail)
        assert len(chain) + len(tail) == len(sources)
