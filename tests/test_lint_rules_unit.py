"""Per-rule unit tests on minimal synthetic decks (Step 5, gate c).

Each test builds the smallest python-pptx deck that trips exactly one
rule and lints it against a hand-written house-profile/1 dict (coded
straight to the pinned schema contract -- deliberately independent of
the profile builder). Synthetic decks are fine HERE: these tests
exercise rule logic on resolved values, not the inheritance resolver
(which has its own PowerPoint-authored fixtures).
"""

import copy

import pytest
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from utils.lint_engine import lint_against_profile

#: Minimal pinned-contract profile the synthetic decks lint against.
BASE_PROFILE = {
    "schema_version": "house-profile/1",
    "name": "unit_test_house",
    "source_decks": ["synthetic.pptx"],
    "slide_size": {"width": {"value": 10.0, "unit": "in"},
                   "height": {"value": 7.5, "unit": "in"}},
    "typography": {
        "body": {"font": {"value": "Calibri"},
                 "size": {"value": 14.0, "unit": "pt"},
                 "bold": {"value": False},
                 "color": {"value": "#20262B"}},
    },
    "paragraph": {
        "space_before": {"value": 2.0, "unit": "pt"},
        "space_after": {"value": 8.0, "unit": "pt"},
        "line_spacing": {"value": 1.0, "unit": None},
        "bullets": {"l1": {"char": {"value": "—"}}},
    },
    "palette": {
        "scheme": {"dk1": {"value": "#20262B"},
                   "lt1": {"value": "#FFFFFF"}},
        "usage": [{"color": "#DDDDDD", "share": 0.5, "role": "fill"}],
    },
    "shape_defaults": {
        "border": {"weight": {"value": 1.0, "unit": "pt"},
                   "color": {"value": "#20262B"},
                   "dash": {"value": "solid"}},
    },
    "grid": {"edges": {"left": [], "right": [], "center": []},
             "unit": "in", "tolerance": {"value": 0.06, "unit": "in"}},
    "archetypes": {},
    "images": {"count_per_slide": {"mean": 0, "max": 0},
               "size_distribution": [], "zones": []},
    "distributions": {
        "font_sizes_pt": {"values": [11.0, 14.0], "shares": [0.5, 0.5]},
        "space_after_pt": {"values": [8.0], "shares": [1.0]},
        "palette_shares": {"values": ["#20262B"], "shares": [1.0]},
    },
}


def _profile(**overrides):
    profile = copy.deepcopy(BASE_PROFILE)
    profile.update(overrides)
    return profile


def _deck(tmp_path, build):
    """Build a one-slide (or more) deck; returns its path as str."""
    prs = Presentation()
    build(prs)
    path = tmp_path / "synthetic.pptx"
    prs.save(str(path))
    return str(path)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, text, left=1, top=1, width=4, height=1,
             size_pt=14.0, name=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = "Calibri"
    if name:
        box.name = name
    return box


def _findings(deck_path, profile=None, rule_id=None,
              severities=("error", "warn")):
    findings = lint_against_profile(deck_path, profile or _profile())
    return [
        f for f in findings
        if (rule_id is None or f["rule_id"] == rule_id)
        and f["severity"] in severities
    ]


# ------------------------------------------------------------ registry

def test_registry_ids_unique_and_documented():
    from utils.lint_rules import RULES, RULES_BY_ID

    ids = [rule.rule_id for rule in RULES]
    assert len(ids) == len(set(ids))
    assert len(RULES_BY_ID) == len(RULES)
    for rule in RULES:
        assert rule.severity in ("error", "warn", "info")
        assert rule.doc.strip()
        assert callable(rule.fn)


def test_registry_covers_the_plan_catalog():
    from utils.lint_rules import RULES_BY_ID

    catalog = {
        "font-scale", "font-family", "bullet-style", "spacing",
        "color-palette", "hardcoded-color", "border-style", "off-grid",
        "straggler-textbox", "off-slide", "archetype-geometry",
        "autofit-shrink", "image-distortion", "image-dpi",
        "footer-presence", "proofing-language", "overlap",
        "empty-slide", "tiny-font", "text-overflow-predicted",
    }
    assert catalog <= set(RULES_BY_ID)


# ------------------------------------------------------------ text rules

def test_font_scale_flags_offscale_run(tmp_path):
    deck = _deck(tmp_path, lambda prs: _textbox(
        _blank_slide(prs), "off scale", size_pt=13.0))
    hits = _findings(deck, rule_id="font-scale")
    assert len(hits) == 1
    assert hits[0]["actual"] == pytest.approx(13.0)
    assert "{11, 14}" in hits[0]["message"]


def test_font_scale_passes_onscale_run(tmp_path):
    deck = _deck(tmp_path, lambda prs: _textbox(
        _blank_slide(prs), "on scale", size_pt=14.0))
    assert _findings(deck, rule_id="font-scale") == []


def test_font_family_flags_foreign_font(tmp_path):
    def build(prs):
        box = _textbox(_blank_slide(prs), "wrong font")
        box.text_frame.paragraphs[0].runs[0].font.name = "Comic Sans MS"

    hits = _findings(_deck(tmp_path, build), rule_id="font-family")
    assert [h["actual"] for h in hits] == ["Comic Sans MS"]


def test_font_family_flags_explicit_ea_typeface(tmp_path):
    def build(prs):
        box = _textbox(_blank_slide(prs), "latin ok, ea foreign")
        rpr = box.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr()
        ea = etree.SubElement(rpr, qn("a:ea"))
        ea.set("typeface", "MS Mincho")

    hits = _findings(_deck(tmp_path, build), rule_id="font-family")
    assert len(hits) == 1
    assert hits[0]["property"] == "font.ea"
    assert hits[0]["actual"] == "MS Mincho"
    assert hits[0]["severity"] == "warn"


def test_font_family_flags_inconsistent_bullet_font(tmp_path):
    def build(prs):
        slide = _blank_slide(prs)
        box = _textbox(slide, "one", height=2)
        frame = box.text_frame
        for text in ("two", "three"):
            para = frame.add_paragraph()
            run = para.add_run()
            run.text = text
            run.font.size = Pt(14)
            run.font.name = "Calibri"
        for index, para in enumerate(frame.paragraphs):
            ppr = para._p.get_or_add_pPr()
            bufont = etree.SubElement(ppr, qn("a:buFont"))
            bufont.set("typeface",
                       "Wingdings" if index == 2 else "Arial")
            buchar = etree.SubElement(ppr, qn("a:buChar"))
            buchar.set("char", "—")

    hits = _findings(_deck(tmp_path, build), rule_id="font-family")
    assert len(hits) == 1
    assert hits[0]["property"] == "bullet.font"
    assert hits[0]["actual"] == "Wingdings"
    assert hits[0]["expected"] == "Arial"


def test_hardcoded_color_flags_theme_equal_srgb_text(tmp_path):
    def build(prs):
        box = _textbox(_blank_slide(prs), "hardcoded")
        run = box.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(0x20, 0x26, 0x2B)

    hits = _findings(_deck(tmp_path, build), rule_id="hardcoded-color")
    assert len(hits) == 1
    assert hits[0]["expected"] == "schemeClr dk1"
    assert hits[0]["actual"] == "srgbClr 20262B"


def test_color_palette_flags_off_palette_fill_and_text(tmp_path):
    def build(prs):
        slide = _blank_slide(prs)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2),
            Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0xFF)
        shape.line.fill.background()
        box = _textbox(slide, "loud", top=3)
        box.text_frame.paragraphs[0].runs[0].font.color.rgb = (
            RGBColor(0x12, 0x34, 0x56))

    hits = _findings(_deck(tmp_path, build), rule_id="color-palette")
    assert {h["property"] for h in hits} == {"fill.color_hex",
                                             "font.color_hex"}
    assert {h["actual"] for h in hits} == {"#FF00FF", "#123456"}


def test_color_palette_allows_usage_colors(tmp_path):
    def build(prs):
        shape = _blank_slide(prs).shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2),
            Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.fill.background()

    assert _findings(_deck(tmp_path, build),
                     rule_id="color-palette") == []


def test_proofing_language_flags_minority_lang(tmp_path):
    def build(prs):
        slide = _blank_slide(prs)
        for index, text in enumerate(("one", "two", "three", "quatre")):
            box = _textbox(slide, text, top=1 + index)
            rpr = box.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr()
            rpr.set("lang", "fr-FR" if text == "quatre" else "en-US")

    hits = _findings(_deck(tmp_path, build), rule_id="proofing-language")
    assert len(hits) == 1
    assert hits[0]["actual"] == "fr-FR"
    assert hits[0]["expected"] == "en-US"


def test_tiny_font_flags_below_floor(tmp_path):
    deck = _deck(tmp_path, lambda prs: _textbox(
        _blank_slide(prs), "tiny", size_pt=6.0))
    hits = _findings(deck, rule_id="tiny-font")
    assert len(hits) == 1
    assert hits[0]["actual"] == pytest.approx(6.0)


def test_autofit_shrink_flags_stale_font_scale(tmp_path):
    def build(prs):
        box = _textbox(_blank_slide(prs), "shrunk")
        bodypr = box.text_frame._txBody.find(qn("a:bodyPr"))
        norm = etree.SubElement(bodypr, qn("a:normAutofit"))
        norm.set("fontScale", "62500")

    hits = _findings(_deck(tmp_path, build), rule_id="autofit-shrink")
    assert len(hits) == 1
    assert hits[0]["actual"] == pytest.approx(62.5)


def test_empty_slide_flagged(tmp_path):
    deck = _deck(tmp_path, lambda prs: _blank_slide(prs))
    hits = _findings(deck, rule_id="empty-slide")
    assert len(hits) == 1
    assert hits[0]["slide"] == 1


# ------------------------------------------------------------ geometry

def test_border_style_flags_weight_color_dash(tmp_path):
    def build(prs):
        shape = _blank_slide(prs).shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2),
            Inches(1))
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shape.line.width = Pt(3.0)

    hits = _findings(_deck(tmp_path, build), rule_id="border-style")
    properties = {h["property"] for h in hits}
    assert "line.weight_pt" in properties
    assert "line.color_hex" in properties
    weight = next(h for h in hits if h["property"] == "line.weight_pt")
    assert weight["actual"] == pytest.approx(3.0)
    assert weight["expected"] == pytest.approx(1.0)


def test_off_grid_flags_unanchored_shape(tmp_path):
    grid = {"edges": {"left": [1.0], "right": [5.0], "center": [3.0]},
            "unit": "in", "tolerance": {"value": 0.06, "unit": "in"}}

    def build(prs):
        slide = _blank_slide(prs)
        _textbox(slide, "anchored", left=1.0, top=1, width=4)
        _textbox(slide, "adrift", left=2.2, top=3, width=1.0,
                 name="Adrift")

    hits = _findings(_deck(tmp_path, build), profile=_profile(grid=grid),
                     rule_id="off-grid")
    assert [h["shape"] for h in hits] == ["Adrift"]
    assert "in" in hits[0]["message"]  # distance-to-gridline message


def test_straggler_textbox_flags_unanchored_footer_text(tmp_path):
    grid = {"edges": {"left": [1.0], "right": [9.0], "center": []},
            "unit": "in", "tolerance": {"value": 0.06, "unit": "in"}}

    def build(prs):
        slide = _blank_slide(prs)
        # sanctioned furniture: left edge on the grid, inside footer zone
        _textbox(slide, "page 1", left=1.0, top=7.0, width=1, height=0.3,
                 name="FooterPage")
        _textbox(slide, "draft", left=4.3, top=7.0, width=2, height=0.3,
                 name="Straggler")

    hits = _findings(_deck(tmp_path, build), profile=_profile(grid=grid),
                     rule_id="straggler-textbox")
    assert [h["shape"] for h in hits] == ["Straggler"]


def test_off_slide_flags_bleeding_shape(tmp_path):
    def build(prs):
        _textbox(_blank_slide(prs), "gone", left=9.5, top=1, width=3,
                 name="Bleeder")

    hits = _findings(_deck(tmp_path, build), rule_id="off-slide")
    assert [h["shape"] for h in hits] == ["Bleeder"]


def test_overlap_flags_floating_pair_not_placeholders(tmp_path):
    def build(prs):
        slide = _blank_slide(prs)
        for name, left in (("A", 1.0), ("B", 1.5)):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(left), Inches(1),
                Inches(2), Inches(1))
            shape.name = name

    hits = _findings(_deck(tmp_path, build), rule_id="overlap")
    assert len(hits) == 1
    assert "'A'" in hits[0]["message"] and "'B'" in hits[0]["message"]


def test_image_dpi_and_distortion(tmp_path):
    from PIL import Image as PILImage

    asset = tmp_path / "tiny.png"
    PILImage.new("RGB", (50, 40), (10, 20, 30)).save(str(asset))

    def build(prs):
        _blank_slide(prs).shapes.add_picture(
            str(asset), Inches(1), Inches(1), Inches(5), Inches(1))

    deck = _deck(tmp_path, build)
    dpi_hits = _findings(deck, rule_id="image-dpi")
    assert len(dpi_hits) == 1
    assert dpi_hits[0]["actual"] < 96
    distortion_hits = _findings(deck, rule_id="image-distortion")
    assert len(distortion_hits) == 1


def test_footer_presence_and_archetype_geometry(tmp_path):
    archetypes = {
        "content": {
            "title_band": {"x": {"value": 5.0, "unit": "in"},
                           "y": {"value": 5.0, "unit": "in"},
                           "w": {"value": 3.0, "unit": "in"},
                           "h": {"value": 0.5, "unit": "in"}},
            "count": 3,
        },
    }

    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text_frame.text = "A working slide"
        body = slide.placeholders[1]
        body.text_frame.text = "Body content line"
        run = body.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)

    deck = _deck(tmp_path, build)
    profile = _profile(archetypes=archetypes)
    footer_hits = _findings(deck, profile=profile,
                            rule_id="footer-presence")
    assert len(footer_hits) == 1
    geometry_hits = _findings(deck, profile=profile,
                              rule_id="archetype-geometry")
    assert len(geometry_hits) == 1
    assert geometry_hits[0]["property"] == "archetype.title_band"


def test_unknown_archetype_reports_info_not_warn(tmp_path):
    archetypes = {
        "two_column": {
            "title_band": {"x": {"value": 0.5, "unit": "in"},
                           "y": {"value": 0.3, "unit": "in"},
                           "w": {"value": 9.0, "unit": "in"},
                           "h": {"value": 0.8, "unit": "in"}},
            "count": 2,
        },
    }
    deck = _deck(tmp_path, lambda prs: _textbox(
        _blank_slide(prs), "just a floating note", size_pt=14.0))
    hits = _findings(deck, profile=_profile(archetypes=archetypes),
                     rule_id="archetype-geometry",
                     severities=("error", "warn", "info"))
    assert len(hits) == 1
    assert hits[0]["severity"] == "info"


# ------------------------------------------------------------ ordering

def test_findings_sorted_severity_then_slide(tmp_path):
    def build(prs):
        _blank_slide(prs)                       # slide 1: empty (warn)
        slide = _blank_slide(prs)               # slide 2: error + warn
        _textbox(slide, "off scale", size_pt=13.0)
        _textbox(slide, "tiny", top=3, size_pt=6.0)

    findings = lint_against_profile(_deck(tmp_path, build), _profile())
    rank = {"error": 0, "warn": 1, "info": 2}
    keys = [(rank[f["severity"]], f["slide"]) for f in findings]
    assert keys == sorted(keys)
    assert findings[0]["severity"] == "error"
