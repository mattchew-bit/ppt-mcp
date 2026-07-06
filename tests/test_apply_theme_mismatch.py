"""Regression tests: house-profile apply on a THEME-MISMATCHED deck.

Repair-pass defect (Step 3): rule/observed hex -> ``a:schemeClr`` token
was mapped through the PROFILE's ``palette.scheme`` without checking
that the TARGET deck's theme resolves the token to the same hex. On a
deck whose theme differs from the profile scheme (the deviant-deck use
case) apply wrote schemeClr tokens that RENDER the deck theme's color
instead of the rule color -- and, because the effective value therefore
never converged on the rule, every write repeated on the second apply,
violating the unconditional no-op-on-second-apply contract.

These tests force the mismatch by pointing profile scheme slots at
hexes the deviant deck's Meridian theme does NOT carry, then assert:

* rule colors land as explicit ``a:srgbClr`` (never a mismatched
  token) so the RENDERED color equals the rule color;
* ``_relink_hardcoded_srgb`` is gated on deck-scheme hex equality (a
  hardcoded srgb is left alone when relinking would change the render);
* the second apply is an unconditional no-op.

The matched-theme behavior (token written, relink performed) stays
covered by tests/test_apply_house_profile.py.
"""

import copy
import shutil
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import fixture_path
from tests.test_apply_house_profile import (
    DEVIANT,
    HOUSE_PROFILE,
    _shape_named,
    _shape_record,
)
from utils.resolve_analysis import build_resolved_analysis
from utils.style_apply import apply_house_profile

#: A hex NO slot of the deviant deck's Meridian theme resolves to.
FOREIGN_HEX = "ABCDEF"

#: Meridian theme truth (corpus_truth.json): accent5 really is 5E8C61,
#: so a profile that claims accent5 == #20262B mismatches the deck.
DECK_ACCENT5 = "5E8C61"


def _mismatch_profile() -> dict:
    """HOUSE_PROFILE variant whose scheme claims accent2 == FOREIGN_HEX
    and whose body-text / l1-bullet / border color rules demand it."""
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_theme_mismatch"
    profile["palette"]["scheme"]["accent2"] = {"value": f"#{FOREIGN_HEX}"}
    profile["typography"]["body"]["color"] = {"value": f"#{FOREIGN_HEX}"}
    profile["paragraph"]["bullets"]["l1"]["color"] = {
        "value": f"#{FOREIGN_HEX}"}
    profile["shape_defaults"]["border"]["color"] = {
        "value": f"#{FOREIGN_HEX}"}
    return profile


def _relink_bait_profile() -> dict:
    """HOUSE_PROFILE variant where the deck's hardcoded ``srgbClr
    20262B`` (violation v4) maps to a scheme token -- accent5 -- that
    the DECK theme resolves to a different hex, so relinking it would
    change the rendered color."""
    profile = copy.deepcopy(HOUSE_PROFILE)
    profile["name"] = "meridian_relink_bait"
    profile["palette"]["scheme"]["dk1"] = {"value": "#101314"}
    profile["palette"]["scheme"]["accent5"] = {"value": "#20262B"}
    return profile


def _apply_and_reopen(tmp_path_factory, profile, label):
    source = fixture_path(DEVIANT)
    if not source.is_file():
        pytest.skip(f"fixture {DEVIANT} not present")
    workdir = tmp_path_factory.mktemp(label)
    copy_path = workdir / "deviant_copy.pptx"
    shutil.copyfile(source, copy_path)
    prs = Presentation(str(copy_path))
    summary = apply_house_profile(prs, profile)
    saved_path = workdir / "applied.pptx"
    prs.save(str(saved_path))
    reopened = Presentation(str(saved_path))
    report = build_resolved_analysis(reopened, detail="full",
                                     max_bytes=5_000_000)
    assert report["truncated"] is False
    return SimpleNamespace(summary=summary, saved_path=saved_path,
                           reopened=reopened, report=report)


@pytest.fixture(scope="module")
def mismatched(tmp_path_factory):
    return _apply_and_reopen(tmp_path_factory, _mismatch_profile(),
                             "theme_mismatch")


@pytest.fixture(scope="module")
def relink_baited(tmp_path_factory):
    return _apply_and_reopen(tmp_path_factory, _relink_bait_profile(),
                             "relink_bait")


# ---------------------------------------------------------------------------
# Rule colors render rule-true on the mismatched deck
# ---------------------------------------------------------------------------

def test_body_color_rule_renders_rule_hex_not_deck_token_color(mismatched):
    body = _shape_record(mismatched.report, 2, "BodyContent")
    for paragraph in body["paragraphs"]:
        for run in paragraph["runs"]:
            assert run["font"]["color_hex"] == FOREIGN_HEX


def test_body_color_written_as_explicit_srgb_not_scheme_token(mismatched):
    shape = _shape_named(mismatched.reopened.slides[1], "BodyContent")
    run = shape.text_frame.paragraphs[0].runs[0]
    fill = run._r.find(qn("a:rPr")).find(qn("a:solidFill"))
    assert fill is not None
    assert fill.find(qn("a:schemeClr")) is None
    assert fill.find(qn("a:srgbClr")).get("val").upper() == FOREIGN_HEX


def test_bullet_color_rule_renders_rule_hex(mismatched):
    body = _shape_record(mismatched.report, 2, "BodyContent")
    level_one_bullets = [
        paragraph["bullet"] for paragraph in body["paragraphs"]
        if paragraph["indent_level"] == 1
        and paragraph["bullet"].get("type") == "char"
    ]
    assert level_one_bullets, "expected level-1 char bullets in BodyContent"
    for bullet in level_one_bullets:
        assert bullet["color_hex"] == FOREIGN_HEX


def test_border_color_rule_renders_rule_hex(mismatched):
    panel = _shape_record(mismatched.report, 3, "TakeawayPanel")
    assert panel["line"]["color_hex"] == FOREIGN_HEX


def test_second_apply_is_a_noop_on_mismatched_theme(mismatched):
    prs = Presentation(str(mismatched.saved_path))
    summary = apply_house_profile(prs, _mismatch_profile())
    assert summary["writes"] == 0
    assert summary["slides_touched"] == []
    assert summary["changes"] == []


# ---------------------------------------------------------------------------
# srgb->schemeClr relink is gated on deck-scheme hex equality
# ---------------------------------------------------------------------------

def test_relink_skipped_when_deck_token_resolves_differently(relink_baited):
    assert all(change["property"] != "font.color_source"
               for change in relink_baited.summary["changes"])
    shape = _shape_named(relink_baited.reopened.slides[3], "ColumnPanelLeft")
    run = shape.text_frame.paragraphs[0].runs[0]
    fill = run._r.find(qn("a:rPr")).find(qn("a:solidFill"))
    assert fill.find(qn("a:schemeClr")) is None
    assert fill.find(qn("a:srgbClr")).get("val").upper() == "20262B"
    # Render check: the deck theme's accent5 must NOT have leaked in.
    record = _shape_record(relink_baited.report, 4, "ColumnPanelLeft")
    run_color = record["paragraphs"][0]["runs"][0]["font"]["color_hex"]
    assert run_color == "20262B"
    assert run_color != DECK_ACCENT5


def test_relink_bait_apply_fixes_other_deviations_then_noops(relink_baited):
    # The 6 non-relink deterministic fixes still land...
    assert relink_baited.summary["writes"] == 6
    # ...and the second apply is an unconditional no-op.
    prs = Presentation(str(relink_baited.saved_path))
    summary = apply_house_profile(prs, _relink_bait_profile())
    assert summary["writes"] == 0
    assert summary["changes"] == []
