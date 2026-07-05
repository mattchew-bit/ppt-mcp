"""Read-only guarantee regression tests for the effective-style resolver.

The resolver's contract (resolve_core / resolve_utils / resolve_analysis
docstrings) is STRICTLY read-only: analyzing a deck must never mutate any
slide, layout, master, theme or presentation part. Two shipped defects
violated it while the value-checking suite stayed green, because the
injected elements (`<a:pPr/>`, empty ``<p:txBody>``) do not change any
resolved value:

    1. ``indent_level_of`` called ``paragraph.level``, whose python-pptx
       implementation goes through ``get_or_add_pPr()`` and INSERTS an
       empty ``<a:pPr/>`` into every paragraph that has none.
    2. ``_shape_record`` guarded on ``shape.has_text_frame`` (True for
       every ``p:sp``) then touched ``shape.text_frame``, whose
       ``get_or_add_txBody()`` CREATES a txBody with a fabricated empty
       paragraph on shapes whose XML legally omits ``p:txBody``.

This module pins the guarantee directly: per-fixture round-trip byte
snapshots of every package part across a full ``build_resolved_analysis``
run, plus one targeted regression test per fixed defect.
"""

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tests.conftest import FIXTURE_DECKS, fixture_path
from utils.resolve_analysis import _shape_record, build_resolved_analysis
from utils.resolve_core import indent_level_of
from utils.resolve_theme import ThemeContext
from utils.resolve_utils import resolve_shape


def _open_fixture(deck: str) -> Presentation:
    path = fixture_path(deck)
    if not path.is_file():
        pytest.skip(f"fixture {deck} not present in tests/fixtures/")
    return Presentation(str(path))


def _part_blobs(prs) -> dict:
    """{partname: serialized bytes} for every part in the package.

    ``blob`` serializes an XML part's CURRENT lxml tree, so in-memory
    mutations show up even though nothing was saved to disk.
    """
    return {
        str(part.partname): part.blob
        for part in prs.part.package.iter_parts()
    }


# ---------------------------------------------------------------------------
# Round-trip immutability: full analysis run leaves every part byte-equal
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("deck", FIXTURE_DECKS)
def test_full_analysis_is_byte_level_read_only(deck):
    prs = _open_fixture(deck)
    before = _part_blobs(prs)

    build_resolved_analysis(prs, detail="full")

    after = _part_blobs(prs)
    assert sorted(after) == sorted(before), (
        f"{deck}: analysis added/removed package parts"
    )
    dirty = [name for name, blob in after.items() if blob != before[name]]
    assert dirty == [], (
        f"{deck}: build_resolved_analysis mutated part(s): {dirty}"
    )


# ---------------------------------------------------------------------------
# Defect 1 regression: indent_level_of must not insert <a:pPr/>
# ---------------------------------------------------------------------------

def test_indent_level_of_does_not_insert_ppr():
    prs = _open_fixture("theme_only.pptx")
    paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]
    p_element = paragraph._p

    # Force the defect precondition: a paragraph with no explicit a:pPr.
    existing = p_element.find(qn("a:pPr"))
    if existing is not None:
        p_element.remove(existing)

    assert indent_level_of(paragraph) == 1
    assert p_element.find(qn("a:pPr")) is None, (
        "indent_level_of inserted an <a:pPr/> into the paragraph XML"
    )


def test_indent_level_of_reads_explicit_lvl():
    prs = _open_fixture("theme_only.pptx")
    paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]
    p_element = paragraph._p

    existing = p_element.find(qn("a:pPr"))
    if existing is not None:
        p_element.remove(existing)
    p_pr = p_element.makeelement(qn("a:pPr"), {"lvl": "2"}, None)
    p_element.insert(0, p_pr)

    assert indent_level_of(paragraph) == 3  # lvl attr is 0-based

    p_pr.attrib.pop("lvl")  # pPr present but no lvl attr -> level 1
    assert indent_level_of(paragraph) == 1


# ---------------------------------------------------------------------------
# Defect 2 regression: _shape_record must not fabricate <p:txBody>
# ---------------------------------------------------------------------------

def test_shape_record_does_not_fabricate_txbody():
    prs = _open_fixture("theme_only.pptx")
    slide = prs.slides[0]
    shape = slide.shapes[0]
    sp_element = shape._element

    # Strip the txBody in-memory: a p:sp without p:txBody is legal OOXML.
    tx_body = sp_element.find(qn("p:txBody"))
    assert tx_body is not None, "fixture shape unexpectedly has no txBody"
    sp_element.remove(tx_body)

    theme = ThemeContext.for_master(slide.slide_layout.slide_master)
    record = _shape_record(slide, shape, theme, {"fonts"})

    assert sp_element.find(qn("p:txBody")) is None, (
        "_shape_record created a <p:txBody> on a txBody-less shape"
    )
    assert "paragraphs" not in record, (
        "_shape_record reported paragraphs fabricated from a created txBody"
    )


def test_shape_record_still_reports_real_paragraphs():
    prs = _open_fixture("theme_only.pptx")
    slide = prs.slides[0]
    shape = slide.shapes[0]
    theme = ThemeContext.for_master(slide.slide_layout.slide_master)

    record = _shape_record(slide, shape, theme, {"fonts"})

    assert "paragraphs" in record
    assert len(record["paragraphs"]) >= 1


# ---------------------------------------------------------------------------
# Defect 4 regression: resolve_shape boundary error for non-slide parts
# ---------------------------------------------------------------------------

def test_resolve_shape_rejects_layout_shape_without_slide():
    prs = _open_fixture("theme_only.pptx")
    layout_shape = prs.slides[0].slide_layout.shapes[0]

    with pytest.raises(ValueError, match="slide"):
        resolve_shape(layout_shape)


def test_resolve_shape_defaults_slide_from_slide_part():
    prs = _open_fixture("theme_only.pptx")
    shape = prs.slides[0].shapes[0]

    result = resolve_shape(shape)  # slide derived from the shape's part

    assert "geometry" in result and "line" in result and "fill" in result
