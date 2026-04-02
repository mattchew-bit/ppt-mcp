"""
Style analysis utilities for PowerPoint presentations.

Extracts font, color, layout, and text hierarchy patterns from existing
presentations to create reusable style profiles.

Ported from charleslukowski/ppt_mcp (MIT) with sklearn/numpy dependency
removed — uses pure-Python frequency counting instead of KMeans clustering.
"""

import json
import logging
from collections import Counter, defaultdict
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.text.text import TextFrame

logger = logging.getLogger(__name__)


@dataclass
class FontInfo:
    """Font styling extracted from a presentation."""
    family: str
    size_pt: float
    bold: bool = False
    italic: bool = False
    color_rgb: Optional[Tuple[int, int, int]] = None
    frequency: int = 1


@dataclass
class ColorInfo:
    """Color usage extracted from a presentation."""
    rgb: Tuple[int, int, int]
    hex_code: str
    context: str  # 'text', 'fill', 'unknown'
    frequency: int = 1


@dataclass
class StyleProfile:
    """Complete style profile for a presentation."""
    name: str
    description: str
    source_file: str
    primary_font: str
    title_font_size: float
    body_font_size: float
    color_palette: List[Dict[str, Any]]
    font_usage: Dict[str, int]
    size_patterns: Dict[str, int]  # str keys for JSON compat
    common_positions: List[Tuple[float, float]]
    slide_dimensions: Tuple[float, float]
    consistency_score: float
    text_hierarchy: Dict[str, Any] = field(default_factory=dict)


def analyze_presentation(file_path: str) -> Dict[str, Any]:
    """Analyze a PowerPoint file and return comprehensive style data."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Presentation not found: {file_path}")

    prs = Presentation(str(path))

    fonts = _analyze_fonts(prs)
    colors = _analyze_colors(prs)
    layouts = _analyze_layouts(prs)
    hierarchy = _analyze_text_hierarchy(prs)
    shapes = _analyze_shapes(prs)

    result = {
        "file_path": str(path),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width.inches,
        "slide_height": prs.slide_height.inches,
        "fonts": fonts,
        "colors": colors,
        "layouts": layouts,
        "text_hierarchy": hierarchy,
        "shapes": shapes,
        "consistency_score": 0.0,
    }
    result["consistency_score"] = _consistency_score(result)
    return result


def create_profile(analysis: Dict[str, Any], name: str) -> StyleProfile:
    """Create a StyleProfile from analysis results."""
    fonts = analysis["fonts"]
    hierarchy = analysis.get("text_hierarchy", {})

    primary_font = fonts.get("primary_font", "Calibri")
    common_sizes = fonts.get("common_sizes", [])
    default_size = common_sizes[0][0] if common_sizes else 18.0

    title_size = default_size + 6
    if hierarchy.get("title", {}).get("common_sizes"):
        title_size = hierarchy["title"]["common_sizes"][0][0]

    body_size = default_size
    if hierarchy.get("body", {}).get("common_sizes"):
        body_size = hierarchy["body"]["common_sizes"][0][0]

    palette = []
    for hex_code, freq in analysis["colors"].get("primary_palette", [])[:10]:
        hex_clean = hex_code.lstrip("#")
        rgb = tuple(int(hex_clean[i : i + 2], 16) for i in (0, 2, 4))
        contexts = analysis["colors"].get("color_contexts", {}).get(hex_code, ["unknown"])
        primary_ctx = max(set(contexts), key=contexts.count)
        palette.append({
            "rgb": rgb,
            "hex": hex_code,
            "context": primary_ctx,
            "frequency": freq,
        })

    return StyleProfile(
        name=name,
        description=f"Style profile from {analysis['file_path']}",
        source_file=analysis["file_path"],
        primary_font=primary_font,
        title_font_size=title_size,
        body_font_size=body_size,
        color_palette=palette,
        font_usage=fonts.get("font_usage", {}),
        size_patterns={str(k): v for k, v in fonts.get("size_patterns", {}).items()},
        common_positions=analysis["layouts"].get("common_positions", []),
        slide_dimensions=(analysis["slide_width"], analysis["slide_height"]),
        consistency_score=analysis["consistency_score"],
        text_hierarchy=analysis.get("text_hierarchy", {}),
    )


def save_profile(profile: StyleProfile, file_path: str) -> None:
    """Save a StyleProfile to a JSON file."""
    data = asdict(profile)
    Path(file_path).write_text(json.dumps(data, indent=2, default=str))
    logger.info("Saved style profile '%s' to %s", profile.name, file_path)


def load_profile(file_path: str) -> StyleProfile:
    """Load a StyleProfile from a JSON file."""
    data = json.loads(Path(file_path).read_text())
    # Convert tuples back from lists
    if "slide_dimensions" in data:
        data["slide_dimensions"] = tuple(data["slide_dimensions"])
    data["common_positions"] = [tuple(p) for p in data.get("common_positions", [])]
    for c in data.get("color_palette", []):
        if "rgb" in c:
            c["rgb"] = tuple(c["rgb"])
    return StyleProfile(**data)


# -- Internal analysis functions -----------------------------------------------

def _analyze_fonts(prs: Presentation) -> Dict[str, Any]:
    font_usage: Dict[str, int] = defaultdict(int)
    size_patterns: Dict[float, int] = defaultdict(int)
    font_details: Dict[str, Dict] = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                _extract_font_info(shape.text_frame, font_usage, font_details, size_patterns)

    primary_font = max(font_usage, key=font_usage.get) if font_usage else "Calibri"
    common_sizes = sorted(size_patterns.items(), key=lambda x: x[1], reverse=True)[:5]

    return {
        "primary_font": primary_font,
        "font_usage": dict(font_usage),
        "common_sizes": common_sizes,
        "size_patterns": dict(size_patterns),
    }


def _extract_font_info(
    tf: TextFrame,
    font_usage: Dict,
    font_details: Dict,
    size_patterns: Dict,
) -> None:
    for para in tf.paragraphs:
        for run in para.runs:
            name = run.font.name
            if not name:
                continue
            font_usage[name] += 1

            if name not in font_details:
                font_details[name] = {"sizes": set(), "bold": 0, "italic": 0, "colors": set()}

            if run.font.size:
                size_pt = run.font.size.pt
                font_details[name]["sizes"].add(size_pt)
                size_patterns[size_pt] += 1

            if run.font.bold:
                font_details[name]["bold"] += 1
            if run.font.italic:
                font_details[name]["italic"] += 1

            try:
                if hasattr(run.font.color, "rgb") and run.font.color.rgb:
                    rgb = run.font.color.rgb
                    font_details[name]["colors"].add((rgb[0], rgb[1], rgb[2]))
            except (AttributeError, TypeError):
                pass


def _analyze_colors(prs: Presentation) -> Dict[str, Any]:
    color_usage: Dict[str, int] = defaultdict(int)
    color_contexts: Dict[str, List[str]] = defaultdict(list)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if hasattr(run.font.color, "rgb") and run.font.color.rgb:
                                rgb = run.font.color.rgb
                                key = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                color_usage[key] += 1
                                color_contexts[key].append("text")
                        except (AttributeError, TypeError):
                            pass

            if hasattr(shape, "fill"):
                try:
                    if shape.fill.type == 1:  # Solid fill
                        if hasattr(shape.fill, "fore_color") and hasattr(shape.fill.fore_color, "rgb"):
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                key = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                                color_usage[key] += 1
                                color_contexts[key].append("fill")
                except (AttributeError, TypeError):
                    pass

    primary_palette = sorted(color_usage.items(), key=lambda x: x[1], reverse=True)[:10]

    return {
        "primary_palette": primary_palette,
        "color_usage": dict(color_usage),
        "color_contexts": {k: v for k, v in color_contexts.items()},
        "total_unique_colors": len(color_usage),
    }


def _analyze_layouts(prs: Presentation) -> Dict[str, Any]:
    positions: List[Tuple[float, float]] = []
    sizes: List[Tuple[float, float]] = []
    lefts: List[float] = []
    tops: List[float] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            w = shape.width.inches if shape.width else 0
            h = shape.height.inches if shape.height else 0
            positions.append((round(left, 2), round(top, 2)))
            sizes.append((round(w, 2), round(h, 2)))
            lefts.append(left)
            tops.append(top)

    # Most common positions by frequency (replaces KMeans clustering)
    common_positions = [pos for pos, _ in Counter(positions).most_common(5)]
    common_sizes = [s for s, _ in Counter(sizes).most_common(5)]

    avg_left = sum(lefts) / len(lefts) if lefts else 0
    avg_top = sum(tops) / len(tops) if tops else 0

    return {
        "average_margins": {"left": round(avg_left, 2), "top": round(avg_top, 2)},
        "common_positions": common_positions,
        "common_sizes": common_sizes,
        "total_shapes": len(positions),
    }


def _analyze_text_hierarchy(prs: Presentation) -> Dict[str, Any]:
    patterns: Dict[str, Dict] = {
        "title": {"fonts": defaultdict(int), "sizes": []},
        "subtitle": {"fonts": defaultdict(int), "sizes": []},
        "body": {"fonts": defaultdict(int), "sizes": []},
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not (hasattr(shape, "text_frame") and shape.text_frame):
                continue

            text_type = _classify_text_type(shape)
            if text_type not in patterns:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        patterns[text_type]["fonts"][run.font.name] += 1
                    if run.font.size:
                        patterns[text_type]["sizes"].append(run.font.size.pt)

    result = {}
    for text_type, data in patterns.items():
        if data["fonts"]:
            result[text_type] = {
                "primary_font": max(data["fonts"], key=data["fonts"].get),
                "common_sizes": Counter(data["sizes"]).most_common(3),
            }
    return result


def _classify_text_type(shape) -> str:
    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return "body"
    top = shape.top.inches if shape.top else 0
    text_len = len(shape.text_frame.text)
    if top < 2 and text_len < 100:
        return "title"
    if top < 3 and text_len < 200:
        return "subtitle"
    return "body"


def _analyze_shapes(prs: Presentation) -> Dict[str, Any]:
    shape_types: Dict[str, int] = defaultdict(int)
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_types[type(shape).__name__] += 1
    return {
        "shape_distribution": dict(shape_types),
        "total_shapes": sum(shape_types.values()),
    }


def _consistency_score(analysis: Dict[str, Any]) -> float:
    factors = []

    font_count = len(analysis["fonts"]["font_usage"])
    factors.append(max(0.0, 1 - (font_count - 1) * 0.1))

    color_count = analysis["colors"]["total_unique_colors"]
    factors.append(max(0.0, 1 - (color_count - 5) * 0.05))

    size_count = len(analysis["fonts"]["size_patterns"])
    factors.append(max(0.0, 1 - (size_count - 3) * 0.1))

    return round(sum(factors) / len(factors), 2) if factors else 0.0
