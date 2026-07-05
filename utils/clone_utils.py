"""Slide cloning with relationship rewriting for PowerPoint MCP Server.

Public API:
    * ``duplicate_slide(prs, slide_index)`` -- same-deck duplicate, appended.
    * ``copy_slide(src_prs, slide_index, dst_prs)`` -- cross-deck copy,
      appended to the destination. v1 requires the two decks to share
      template lineage: the copied slide is re-bound to the destination
      layout **matched by layout name**.

Ported semantics (MIT sources, re-implemented -- not copied):
    * pptx-automizer (github.com/singerla/pptx-automizer): per-rel-type
      handlers (image / media / hyperlink) plus its ``assertRelatedContent``
      pattern -- after the typed handlers run, a final sweep copies any
      remaining unhandled rel target so no broken rel survives (broken rels
      are the #1 cause of PowerPoint repair prompts). ``[Content_Types].xml``
      coverage comes for free: python-pptx regenerates it from part content
      types at save time, and the sweep guarantees every target is a real
      part in the destination package.
    * python-pptx PR #1141 (scanny/python-pptx): ``SlidePart`` clone with
      relationships re-created in ascending rId order. This module goes one
      step further and rewrites the r:id / r:embed attributes in the slide
      XML from an old->new rId mapping, so decks with non-contiguous rIds
      survive cloning too.

v1 scope guards:
    * charts, SmartArt diagrams, OLE objects, and ActiveX controls are
      rejected with ``NotImplementedError`` naming the offending shape,
    * notesSlide rels are never copied (issue-#132 pitfall) -- speaker
      notes are dropped by both ``duplicate_slide`` and ``copy_slide``,
    * internal slide-jump hyperlinks (``ppaction://hlinksldjump``, rel type
      RT.SLIDE) are dropped on cross-deck copy (text kept); same-deck
      duplicates keep them since the target slide is shared,
    * the final sweep only blob-copies *leaf* parts; unhandled rel targets
      that carry their own relationships raise ``NotImplementedError``,
    * ``a16:creationId`` / ``p14:creationId`` values are stripped from the
      clone so PowerPoint regenerates them (duplicate-ID corruption pitfall),
    * a locally-defined slide background travels with the deep-copied slide
      XML (``p:bg`` lives inside ``p:cSld``) -- no ``Slide.background``
      element surgery, which is broken upstream (python-pptx issue #1126).
"""

import copy
import io
import re
from typing import Dict, List, Optional, Tuple

from pptx.media import Video
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.parts.slide import SlidePart
from pptx.presentation import Presentation

# ---------------------------------------------------------------------------
# Namespace / rel-type constants
# ---------------------------------------------------------------------------

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A16_CREATION_ID = (
    "{http://schemas.microsoft.com/office/drawing/2014/main}creationId"
)
_P14_CREATION_ID = (
    "{http://schemas.microsoft.com/office/powerpoint/2010/main}creationId"
)

_MEDIA_RELTYPES = frozenset({
    RT.MEDIA,
    RT.VIDEO,
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio",
})

# Explicitly rejected in v1: rel graphs too gnarly to rewrite safely.
_UNSUPPORTED_RELTYPES: Dict[str, str] = {
    RT.CHART: "chart",
    "http://schemas.microsoft.com/office/2014/relationships/chartEx": "chart",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "diagramData": "SmartArt diagram",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "diagramLayout": "SmartArt diagram",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "diagramColors": "SmartArt diagram",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "diagramQuickStyle": "SmartArt diagram",
    "http://schemas.microsoft.com/office/2007/relationships/diagramDrawing":
        "SmartArt diagram",
    RT.OLE_OBJECT: "OLE object",
    RT.PACKAGE: "OLE object",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "control": "ActiveX control",
    "http://schemas.microsoft.com/office/2006/relationships/"
    "activeXControlBinary": "ActiveX control",
}

_SHAPE_TAGS = frozenset({
    f"{{{_P_NS}}}sp",
    f"{{{_P_NS}}}pic",
    f"{{{_P_NS}}}graphicFrame",
    f"{{{_P_NS}}}grpSp",
    f"{{{_P_NS}}}cxnSp",
})

_HYPERLINK_TAGS = frozenset({
    f"{{{_A_NS}}}hlinkClick",
    f"{{{_A_NS}}}hlinkHover",
})


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def duplicate_slide(prs, slide_index: int):
    """Append a duplicate of ``prs.slides[slide_index]`` and return it.

    Same-deck clone following python-pptx PR #1141 semantics: the slide XML
    is deep-copied while related parts (images, media, layout) are *shared*
    with the source slide. Speaker notes are not copied (the duplicate
    carries no notesSlide rel). Slides bearing charts, SmartArt, OLE
    objects, or ActiveX controls raise ``NotImplementedError`` (v1 scope
    guard).
    """
    _require_presentation(prs, "prs")
    slide = _slide_at(prs, slide_index)
    _reject_unsupported_rels(slide.part, slide_index)

    new_element = copy.deepcopy(slide.part._element)
    _strip_creation_ids(new_element)

    package = prs.part.package
    new_part = SlidePart(
        package.next_partname("/ppt/slides/slide%d.xml"),
        CT.PML_SLIDE,
        package,
        new_element,
    )
    return _attach_and_populate(prs, slide.part, new_part, share_parts=True)


def copy_slide(src_prs, slide_index: int, dst_prs):
    """Copy ``src_prs.slides[slide_index]`` into ``dst_prs`` and return it.

    v1 constraint: source and destination must share template lineage. The
    copied slide is re-bound to the destination layout matched by layout
    name; when no layout of that name exists a ``ValueError`` explains that
    the destination deck should be started from the same template.

    Image and media parts are re-created in the destination package and
    deduplicated by content hash; external hyperlinks are preserved;
    internal slide-jump hyperlinks are dropped (text kept); notesSlide
    rels are excluded; any other leaf rel target is copied by the final
    safety sweep so no broken rel survives.
    """
    _require_presentation(src_prs, "src_prs")
    _require_presentation(dst_prs, "dst_prs")
    slide = _slide_at(src_prs, slide_index)
    layout_part = _match_layout_by_name(src_prs, slide, dst_prs)
    _reject_unsupported_rels(slide.part, slide_index)

    new_element = copy.deepcopy(slide.part._element)
    _strip_creation_ids(new_element)

    package = dst_prs.part.package
    new_part = SlidePart(
        package.next_partname("/ppt/slides/slide%d.xml"),
        CT.PML_SLIDE,
        package,
        new_element,
    )
    return _attach_and_populate(
        dst_prs, slide.part, new_part, share_parts=False, layout_part=layout_part
    )


# ---------------------------------------------------------------------------
# Input validation
# ---------------------------------------------------------------------------


def _require_presentation(obj, param_name: str) -> None:
    if not isinstance(obj, Presentation):
        raise ValueError(
            f"{param_name} must be a python-pptx Presentation, "
            f"got {type(obj).__name__}"
        )


def _slide_at(prs, slide_index: int):
    slide_count = len(prs.slides)
    if not isinstance(slide_index, int) or isinstance(slide_index, bool):
        raise ValueError(f"slide index must be an int, got {slide_index!r}")
    if not 0 <= slide_index < slide_count:
        raise ValueError(
            f"Invalid slide index: {slide_index}. "
            f"Available slides: 0-{slide_count - 1}"
        )
    return prs.slides[slide_index]


# ---------------------------------------------------------------------------
# v1 scope guard: reject charts / SmartArt / OLE / ActiveX
# ---------------------------------------------------------------------------


def _reject_unsupported_rels(slide_part, slide_index: int) -> None:
    for rid, rel in _rels_ascending(slide_part):
        label = _UNSUPPORTED_RELTYPES.get(rel.reltype)
        if label is None:
            continue
        shape_name = _shape_name_for_rid(slide_part._element, rid)
        raise NotImplementedError(
            f"Slide {slide_index} contains shape '{shape_name}' with a "
            f"{label}; copying charts, SmartArt, OLE objects, and ActiveX "
            "controls is not supported in v1. Remove or replace that shape "
            "and retry."
        )


def _shape_name_for_rid(slide_element, rid: str) -> str:
    """Name (p:cNvPr @name) of the shape whose XML references ``rid``."""
    holder = None
    for element in slide_element.iter():
        for key, value in element.attrib.items():
            if key.startswith(f"{{{_R_NS}}}") and value == rid:
                holder = element
                break
        if holder is not None:
            break
    while holder is not None and holder.tag not in _SHAPE_TAGS:
        holder = holder.getparent()
    if holder is None:
        return "(unknown shape)"
    cNvPr = holder.find(f".//{{{_P_NS}}}cNvPr")
    if cNvPr is None or not cNvPr.get("name"):
        return "(unnamed shape)"
    return cNvPr.get("name")


# ---------------------------------------------------------------------------
# creationId stripping (duplicate-ID corruption pitfall)
# ---------------------------------------------------------------------------


def _strip_creation_ids(slide_element) -> None:
    """Remove a16/p14 creationId elements so PowerPoint regenerates them."""
    doomed = [
        element
        for element in slide_element.iter()
        if element.tag in (_A16_CREATION_ID, _P14_CREATION_ID)
    ]
    for element in doomed:
        parent = element.getparent()
        parent.remove(element)
        # Prune the wrapping <ext uri="..."> if the removal emptied it.
        if parent.tag.endswith("}ext") and len(parent) == 0:
            grandparent = parent.getparent()
            if grandparent is not None:
                grandparent.remove(parent)


# ---------------------------------------------------------------------------
# Layout re-binding (v1: matched by layout name)
# ---------------------------------------------------------------------------


def _match_layout_by_name(src_prs, slide, dst_prs):
    """Destination layout part whose name matches the source slide's layout.

    When several destination masters carry the same layout name (common in
    multi-master decks), the layout under the master at the same position as
    the source master is preferred. A name match alone is not enough
    (issue-#132 pitfall): the matched layout must also carry every ``p:ph``
    idx the slide inherits from its source layout, otherwise placeholder
    inheritance silently falls back to master defaults.
    """
    src_layout = slide.slide_layout
    layout_name = src_layout.name
    src_master_index = _master_index(src_prs, src_layout)

    candidates: List[Tuple[int, object]] = [
        (master_index, layout)
        for master_index, master in enumerate(dst_prs.slide_masters)
        for layout in master.slide_layouts
        if layout.name == layout_name
    ]
    if not candidates:
        available = sorted({
            layout.name
            for master in dst_prs.slide_masters
            for layout in master.slide_layouts
        })
        raise ValueError(
            f"Destination has no slide layout named {layout_name!r} "
            f"(available layouts: {', '.join(available)}). copy_slide v1 "
            "requires both decks to share template lineage -- start the "
            "destination presentation from the same template as the source."
        )
    # Same-master candidate first, then document order (stable sort).
    ordered = sorted(candidates, key=lambda item: item[0] != src_master_index)
    inherited_idxs = _inherited_placeholder_idxs(slide)
    for _, layout in ordered:
        layout_idxs = {ph.placeholder_format.idx for ph in layout.placeholders}
        if inherited_idxs <= layout_idxs:
            return layout.part
    raise ValueError(
        f"Destination layout {layout_name!r} matches by name but lacks "
        f"placeholder idx values {sorted(inherited_idxs)} used by the source "
        "slide, so placeholder inheritance would silently break. copy_slide "
        "v1 requires both decks to share template lineage -- start the "
        "destination presentation from the same template as the source."
    )


def _inherited_placeholder_idxs(slide) -> set:
    """``p:ph`` idx values the slide actually inherits from its layout.

    Placeholder idxs present on the slide but absent from its *source*
    layout are pre-existing orphans (they never inherited layout formatting)
    and are excluded, so they do not block an otherwise-valid re-binding.
    """
    slide_idxs = {ph.placeholder_format.idx for ph in slide.placeholders}
    src_layout_idxs = {
        ph.placeholder_format.idx for ph in slide.slide_layout.placeholders
    }
    return slide_idxs & src_layout_idxs


def _master_index(prs, layout) -> int:
    layout_master_element = layout.slide_master.part._element
    for index, master in enumerate(prs.slide_masters):
        if master.part._element is layout_master_element:
            return index
    return 0


# ---------------------------------------------------------------------------
# Relationship re-creation (the rel-rewriting engine)
# ---------------------------------------------------------------------------


def _rels_ascending(part) -> List[Tuple[str, object]]:
    """Rels sorted by ascending numeric rId (PR #1141 ordering)."""
    def _numeric(rid: str) -> int:
        digits = re.sub(r"\D", "", rid)
        return int(digits) if digits else 0

    return sorted(part.rels.items(), key=lambda item: _numeric(item[0]))


def _recreate_rels(
    src_part,
    new_part,
    share_parts: bool,
    layout_part=None,
) -> Dict[str, str]:
    """Re-create every rel of ``src_part`` on ``new_part``.

    Returns an old-rId -> new-rId mapping to rewrite the slide XML with.
    ``share_parts=True`` (same-deck duplicate) relates to the existing
    target parts; ``share_parts=False`` (cross-deck copy) re-creates each
    target in the destination package via the typed handlers, with the
    final generic sweep covering every rel type not explicitly handled.
    """
    rid_map: Dict[str, str] = {}
    for old_rid, rel in _rels_ascending(src_part):
        if rel.reltype == RT.NOTES_SLIDE:
            continue  # never copied (issue-#132 pitfall)
        new_rid = _recreate_one_rel(rel, new_part, share_parts, layout_part)
        if new_rid is not None:
            rid_map[old_rid] = new_rid
        else:
            _strip_hyperlink_refs(new_part._element, old_rid)
    return rid_map


def _recreate_one_rel(
    rel, new_part, share_parts: bool, layout_part
) -> Optional[str]:
    """Handler dispatch for one source rel; returns the new rId or None."""
    if rel.is_external:
        return new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
    if share_parts:
        return new_part.relate_to(rel.target_part, rel.reltype)
    if rel.reltype == RT.SLIDE_LAYOUT:
        if layout_part is None:
            raise RuntimeError("cross-deck copy requires a destination layout")
        return new_part.relate_to(layout_part, RT.SLIDE_LAYOUT)
    if rel.reltype == RT.IMAGE:
        return _copy_image_rel(rel, new_part)
    if rel.reltype in _MEDIA_RELTYPES:
        return _copy_media_rel(rel, new_part)
    if rel.reltype == RT.SLIDE:
        # Internal slide-jump hyperlink (a:hlinkClick
        # action="ppaction://hlinksldjump"): its rel targets another slide
        # part, and that slide is not part of the copy. Drop the link
        # action, keep the text (v1 behavior). Note: ECMA-376 hyperlink
        # rels (RT.HYPERLINK) are always external, so slide jumps are the
        # only internal-hyperlink case; external links are handled by the
        # is_external branch above.
        return None
    return _sweep_copy_rel(rel, new_part)


def _copy_image_rel(rel, new_part) -> str:
    """Image handler: re-create in destination, deduplicated by SHA1 hash."""
    package = new_part.package
    image_part = package.get_or_add_image_part(io.BytesIO(rel.target_part.blob))
    return new_part.relate_to(image_part, RT.IMAGE)


def _copy_media_rel(rel, new_part) -> str:
    """Media handler: re-create audio/video part, deduplicated by SHA1 hash."""
    source_part = rel.target_part
    media = Video.from_blob(source_part.blob, source_part.content_type)
    media_part = new_part.package.get_or_add_media_part(media)
    return new_part.relate_to(media_part, rel.reltype)


def _sweep_copy_rel(rel, new_part) -> str:
    """Final safety sweep (pptx-automizer ``assertRelatedContent`` pattern).

    Any *leaf* rel type without a dedicated handler still gets its target
    part copied into the destination package, so the saved deck never
    contains a rel pointing at a missing part. Content-type entries
    regenerate from the part's content type when python-pptx serializes
    the package.

    Targets that carry their own relationships are refused: a blob-only
    copy would strand the target's internal r:id references (dangling
    in-package refs -- the exact repair-prompt failure the sweep exists to
    prevent).
    """
    target = rel.target_part
    if len(target.rels) > 0:
        raise NotImplementedError(
            f"Slide has a relationship of type {rel.reltype!r} to part "
            f"'{target.partname}', which carries its own relationships; "
            "copying such non-leaf parts is not supported in v1. Remove "
            "the referencing element and retry."
        )
    target_part = _get_or_copy_part(new_part.package, target)
    return new_part.relate_to(target_part, rel.reltype)


def _get_or_copy_part(dst_package, src_part):
    """Existing identical part in ``dst_package``, or a fresh blob copy."""
    for part in dst_package.iter_parts():
        if (
            part.content_type == src_part.content_type
            and part.blob == src_part.blob
        ):
            return part
    partname = dst_package.next_partname(_partname_template(str(src_part.partname)))
    return Part(partname, src_part.content_type, dst_package, src_part.blob)


def _partname_template(partname: str) -> str:
    """``/ppt/tags/tag1.xml`` -> ``/ppt/tags/tag%d.xml`` (for next_partname)."""
    match = re.match(r"^(.*?)(\d+)?(\.[A-Za-z0-9]+)$", partname)
    if match is None:
        return f"{partname}%d"
    return f"{match.group(1)}%d{match.group(3)}"


# ---------------------------------------------------------------------------
# Slide-XML rId rewriting + integrity check
# ---------------------------------------------------------------------------


def _apply_rid_mapping(slide_element, rid_map: Dict[str, str]) -> None:
    """Rewrite every r-namespace attribute per the old->new rId mapping.

    A single pass is swap-safe: each attribute's *original* value is looked
    up in the mapping exactly once.
    """
    prefix = f"{{{_R_NS}}}"
    for element in slide_element.iter():
        for key, value in list(element.attrib.items()):
            if key.startswith(prefix) and value in rid_map:
                element.set(key, rid_map[value])


def _strip_hyperlink_refs(slide_element, rid: str) -> None:
    """Remove hlinkClick/hlinkHover elements referencing a dropped rel."""
    id_attr = f"{{{_R_NS}}}id"
    doomed = [
        element
        for element in slide_element.iter()
        if element.tag in _HYPERLINK_TAGS and element.get(id_attr) == rid
    ]
    for element in doomed:
        element.getparent().remove(element)


def _assert_no_dangling_rids(slide_element, new_part) -> None:
    """Every r:id referenced by the new slide XML must resolve in its rels.

    Empty r:id values are ignored: an empty rId is the OOXML convention for
    action-only hyperlinks -- python-pptx movie shapes always emit
    ``<a:hlinkClick r:id="" action="ppaction://media"/>``, and treating ""
    as a reference would falsely reject every media slide.
    """
    prefix = f"{{{_R_NS}}}"
    referenced = {
        value
        for element in slide_element.iter()
        for key, value in element.attrib.items()
        if key.startswith(prefix) and value
    }
    dangling = referenced - set(new_part.rels.keys())
    if dangling:
        raise RuntimeError(
            "slide clone produced dangling relationship references "
            f"({', '.join(sorted(dangling))}); aborting to avoid writing a "
            "deck PowerPoint would ask to repair"
        )


# ---------------------------------------------------------------------------
# Appending the new slide to a presentation
# ---------------------------------------------------------------------------


def _attach_and_populate(prs, src_slide_part, new_part, share_parts, layout_part=None):
    """Relate the new slide part first, then re-create its rels; then append.

    The presentation->slide rel must exist *before* image/media/sweep parts
    are created: python-pptx allocates partnames (``next_partname``,
    ``next_image_partname``) and deduplicates by hash by walking the rel
    graph from the package root, so parts created for a not-yet-reachable
    slide would be invisible -- two images copied in one call would collide
    on the same partname and corrupt the saved deck.

    On failure the presentation->slide rel is rolled back so the destination
    deck is left untouched (any parts created meanwhile become unreachable
    and are never serialized).
    """
    rid = prs.part.relate_to(new_part, RT.SLIDE)
    try:
        rid_map = _recreate_rels(src_slide_part, new_part, share_parts, layout_part)
        _apply_rid_mapping(new_part._element, rid_map)
        _assert_no_dangling_rids(new_part._element, new_part)
    except Exception:
        prs.part.rels.pop(rid)
        raise
    prs.slides._sldIdLst.add_sldId(rid)
    return new_part.slide
