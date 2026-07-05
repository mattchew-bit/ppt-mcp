"""Deterministic house-profile application engine (Step 3 write side).

Applies the DETERMINISTIC sections of a ``house-profile/1`` profile --
``typography`` / ``paragraph`` / ``palette`` / ``shape_defaults`` -- to
an open presentation. The analysis-side sections (``grid``,
``archetypes``, ``distributions``, ``images``) are consumed by the model
at generation time and enforced by the Step 5 lint; they are NEVER
applied here, and shape geometry is never touched.

Minimal-diff contract: every rule is checked against the EFFECTIVE
(inheritance-resolved, Step 2 resolver) value first and written only
where the effective value deviates. Content that is already conformant
through inheritance receives no new explicit overrides, so a second
apply of the same profile is a no-op.

Rule -> target mapping (v1, documented):
    * ``typography.title``  -- placeholders of type title/ctrTitle
    * ``typography.body``   -- placeholders of type body/obj; the size
      rule applies to indent level 1 only (deeper levels follow the
      template's own level scale); font/bold/color apply to all levels
    * ``typography.footer`` -- placeholders of type ftr/sldNum/dt
    * ``paragraph``         -- body-role placeholders; spacing rules at
      indent level 1, ``bullets.l<N>`` per matching indent level
    * ``palette.scheme``    -- run colors hardcoded as ``a:srgbClr``
      whose hex equals a scheme color are re-linked to the
      ``a:schemeClr`` token (visual no-op, theme-true XML). A token is
      written ONLY when the TARGET deck's theme resolves it to the same
      hex; on a deck whose theme deviates from the profile scheme the
      rule color is written as explicit ``a:srgbClr`` instead, so the
      rendered color is always the rule color and a second apply stays
      a no-op
    * ``shape_defaults``    -- non-placeholder ``p:sp`` shapes: border
      weight/color/dash where a border is visibly drawn, solid fills
      HARDCODED as an untransformed ``a:srgbClr`` whose hex is
      off-palette AND differs from the house fill snap to the house
      fill (theme-linked ``a:schemeClr`` fills -- including tints via
      lumMod/lumOff etc. -- transformed srgb fills and inherited fills
      are never touched: they are theme-true by construction, so
      repainting them would destroy the theme link on conformant
      content), roundRect corner radius

Leaf values are DTCG-shaped (``{"value": X, "unit": "pt"|"in"|null}`` or
``{"value": "#RRGGBB"}``); bare scalars are tolerated for robustness.
"""

from typing import Any, Dict, FrozenSet, List, NamedTuple, Optional

from pptx.oxml.ns import qn

from .resolve_theme import ThemeContext
from .style_roles import placeholder_role
from .resolve_utils import (
    TextContext,
    resolve_paragraph,
    resolve_run_font,
    resolve_shape,
)
from .style_write import (
    SCHEME_TOKENS,
    set_paragraph_bullet,
    set_paragraph_spacing,
    set_run_color,
    set_run_font,
    set_shape_border,
    set_shape_corner_radius,
    set_shape_fill,
)

HOUSE_PROFILE_SCHEMA = "house-profile/1"

#: Comparison tolerances (floats out of EMU / percent arithmetic).
_PT_TOLERANCE = 0.05
_MULTIPLE_TOLERANCE = 0.01
_PCT_TOLERANCE = 0.5
_ADJ_TOLERANCE = 0.005

_INCHES_TO_POINTS = 72.0

#: Cap on the per-change detail list in the apply summary.
_MAX_CHANGE_RECORDS = 200


def is_house_profile(profile: Any) -> bool:
    """True when ``profile`` is a house-profile/1 dict (vs legacy)."""
    return (isinstance(profile, dict)
            and profile.get("schema_version") == HOUSE_PROFILE_SCHEMA)


# ---------------------------------------------------------------------------
# DTCG-shaped leaf readers
# ---------------------------------------------------------------------------

def leaf_value(node: Any) -> Any:
    """DTCG leaf ``{"value": X, ...}`` -> ``X``; bare scalars pass through."""
    if isinstance(node, dict):
        if "value" not in node:
            raise ValueError(f"malformed profile leaf (no 'value'): {node!r}")
        return node["value"]
    return node


def leaf_points(node: Any) -> Optional[float]:
    """Dimension leaf -> points (``in`` converted, ``pt``/bare as-is)."""
    if node is None:
        return None
    value = leaf_value(node)
    if not isinstance(value, (int, float)):
        raise ValueError(f"profile dimension is not numeric: {node!r}")
    unit = node.get("unit") if isinstance(node, dict) else None
    if unit == "in":
        return float(value) * _INCHES_TO_POINTS
    if unit in (None, "pt"):
        return float(value)
    raise ValueError(f"unsupported profile dimension unit {unit!r}")


def _leaf_hex(node: Any) -> Optional[str]:
    if node is None:
        return None
    value = leaf_value(node)
    if not isinstance(value, str):
        raise ValueError(f"profile color is not a hex string: {node!r}")
    return value.lstrip("#").upper()


def _hex_or_none(value: Optional[str]) -> Optional[str]:
    return None if value is None else value.lstrip("#").upper()


def _close(a: Optional[float], b: Optional[float], tolerance: float) -> bool:
    if a is None or b is None:
        return a is b
    return abs(a - b) <= tolerance


# ---------------------------------------------------------------------------
# Palette map
# ---------------------------------------------------------------------------

def _scheme_by_hex(profile: Dict) -> Dict[str, str]:
    """Reverse map hex -> scheme token, canonical token order wins ties."""
    scheme = (profile.get("palette") or {}).get("scheme") or {}
    by_hex: Dict[str, str] = {}
    for token in SCHEME_TOKENS:
        if token not in scheme:
            continue
        hex_value = _leaf_hex(scheme[token])
        if hex_value is not None and hex_value not in by_hex:
            by_hex[hex_value] = token
    return by_hex


class _PaletteView(NamedTuple):
    """Profile palette hexes + the deck-theme-verified token map.

    ``tokens_by_hex`` maps a hex to a scheme token ONLY where the
    TARGET deck's theme resolves that token to the same hex -- writing
    the token is then a render no-op that restores the theme link. On
    a deck whose theme deviates from the profile scheme the entry is
    filtered out and callers fall back to explicit ``a:srgbClr``,
    which keeps the rendered color rule-true and the second apply a
    no-op. ``hexes`` is the FULL profile palette (membership test for
    the off-palette fill rule, deliberately deck-independent -- the
    house palette defines what "on palette" means).
    """
    hexes: FrozenSet[str]
    tokens_by_hex: Dict[str, str]


def _palette_view(scheme_by_hex: Dict[str, str],
                  theme: ThemeContext) -> _PaletteView:
    """Deck-verified palette view for one master's theme."""
    deck_scheme = theme.color_scheme
    verified = {
        hex_value: token
        for hex_value, token in scheme_by_hex.items()
        if deck_scheme.get(token) == hex_value
    }
    return _PaletteView(hexes=frozenset(scheme_by_hex),
                        tokens_by_hex=verified)


class _ChangeLog:
    """Write counter + capped per-change detail list."""

    def __init__(self):
        self.writes = 0
        self.changes: List[Dict] = []

    def record(self, slide_number: int, shape_name: str, prop: str,
               old: Any, new: Any) -> None:
        self.writes += 1
        if len(self.changes) < _MAX_CHANGE_RECORDS:
            self.changes.append({
                "slide": slide_number, "shape": shape_name,
                "property": prop, "from": old, "to": new,
            })


# ---------------------------------------------------------------------------
# Role classification
# ---------------------------------------------------------------------------

def _role_of(shape) -> Optional[str]:
    """``"title"`` / ``"body"`` / ``"footer"`` for placeholders, else None.

    Delegates to the shared exact-match table in ``utils.style_roles``
    so learn and apply agree on which placeholder types own a role.
    """
    if not shape.is_placeholder:
        return None
    ph_type = shape.placeholder_format.type
    return placeholder_role(None if ph_type is None else ph_type.name)


def _has_text_body(shape) -> bool:
    """True when the shape's XML carries ``p:txBody`` (see
    resolve_analysis: ``shape.text_frame`` on a body-less sp would
    CREATE one -- a mutation we must not make as a side effect)."""
    return shape._element.find(qn("p:txBody")) is not None


# ---------------------------------------------------------------------------
# Run-level rules
# ---------------------------------------------------------------------------

def _write_color(run, rule_hex: str, palette: _PaletteView) -> None:
    """Rule hex -> schemeClr token when the deck theme agrees, else srgbClr."""
    token = palette.tokens_by_hex.get(rule_hex)
    if token is not None:
        set_run_color(run, scheme=token)
    else:
        set_run_color(run, hex_value=rule_hex)


def _apply_run_typography(run, effective: Dict, rule: Dict, role: str,
                          level: int, palette: _PaletteView, log: _ChangeLog,
                          slide_number: int, shape_name: str) -> None:
    """Write only the run properties whose effective value deviates.

    The level-1-only restriction on the size rule is scoped to the
    ``body`` role (deeper body levels follow the template's own level
    scale); title/footer runs get the size rule at ANY indent level.
    """
    font_rule = rule.get("font") and leaf_value(rule["font"])
    if font_rule and effective["name"] != font_rule:
        set_run_font(run, name=font_rule)
        log.record(slide_number, shape_name, "font.name",
                   effective["name"], font_rule)
    size_rule = leaf_points(rule.get("size"))
    size_in_scope = role != "body" or level == 1
    if (size_rule is not None and size_in_scope
            and not _close(effective["size_pt"], size_rule, _PT_TOLERANCE)):
        set_run_font(run, size_pt=size_rule)
        log.record(slide_number, shape_name, "font.size_pt",
                   effective["size_pt"], size_rule)
    if "bold" in rule:
        bold_rule = bool(leaf_value(rule["bold"]))
        if bool(effective["bold"]) != bold_rule:
            set_run_font(run, bold=bold_rule)
            log.record(slide_number, shape_name, "font.bold",
                       effective["bold"], bold_rule)
    color_rule = _leaf_hex(rule.get("color"))
    if (color_rule is not None
            and _hex_or_none(effective["color_hex"]) != color_rule):
        _write_color(run, color_rule, palette)
        log.record(slide_number, shape_name, "font.color",
                   effective["color_hex"], color_rule)


def _relink_hardcoded_srgb(run, palette: _PaletteView,
                           log: _ChangeLog, slide_number: int,
                           shape_name: str) -> None:
    """``a:srgbClr`` equal to a scheme color -> ``a:schemeClr`` token.

    Visual no-op that restores the theme link (deviation kind
    ``hardcoded_srgb``). Gated on the DECK theme resolving the token to
    the same hex (``_PaletteView.tokens_by_hex``): on a deck whose
    theme deviates from the profile scheme, relinking would CHANGE the
    rendered color, so the explicit srgb is left alone. Transformed
    colors (srgbClr with lumMod etc. children) are left alone too --
    their final value is not the scheme hex.
    """
    r_pr = run._r.find(qn("a:rPr"))
    if r_pr is None:
        return
    fill = r_pr.find(qn("a:solidFill"))
    if fill is None:
        return
    srgb = fill.find(qn("a:srgbClr"))
    if srgb is None or len(srgb):
        return
    token = palette.tokens_by_hex.get((srgb.get("val") or "").upper())
    if token is None:
        return
    set_run_color(run, scheme=token)
    log.record(slide_number, shape_name, "font.color_source",
               f"srgbClr {srgb.get('val')}", f"schemeClr {token}")


# ---------------------------------------------------------------------------
# Paragraph-level rules (body role)
# ---------------------------------------------------------------------------

def _spacing_deviation(resolved: Dict, rule_pt: Optional[float]) -> bool:
    if rule_pt is None:
        return False
    if "points" in resolved:
        return not _close(resolved["points"], rule_pt, _PT_TOLERANCE)
    return True  # spcPct-tagged spacing vs a points rule always deviates


def _apply_paragraph_spacing(paragraph, resolved: Dict, rules: Dict,
                             log: _ChangeLog, slide_number: int,
                             shape_name: str) -> None:
    before_rule = leaf_points(rules.get("space_before"))
    after_rule = leaf_points(rules.get("space_after"))
    line_rule = rules.get("line_spacing") and leaf_value(
        rules["line_spacing"])
    before = (before_rule if _spacing_deviation(resolved["space_before"],
                                                before_rule) else None)
    after = (after_rule if _spacing_deviation(resolved["space_after"],
                                              after_rule) else None)
    line = None
    if line_rule is not None:
        resolved_line = resolved["line_spacing"]
        if ("multiple" not in resolved_line
                or not _close(resolved_line["multiple"], float(line_rule),
                              _MULTIPLE_TOLERANCE)):
            line = float(line_rule)
    if before is None and after is None and line is None:
        return
    set_paragraph_spacing(paragraph, space_before_pt=before,
                          space_after_pt=after, line_spacing=line)
    for prop, old, new in (("space_before", resolved["space_before"], before),
                           ("space_after", resolved["space_after"], after),
                           ("line_spacing", resolved["line_spacing"], line)):
        if new is not None:
            log.record(slide_number, shape_name, f"paragraph.{prop}",
                       old, new)


def _apply_bullet_rule(paragraph, resolved_bullet: Dict, rule: Dict,
                       palette: _PaletteView, log: _ChangeLog,
                       slide_number: int, shape_name: str) -> None:
    char_rule = rule.get("char") and leaf_value(rule["char"])
    char = None
    if char_rule and (resolved_bullet.get("type") != "char"
                      or resolved_bullet.get("char") != char_rule):
        char = char_rule
    size_rule = rule.get("size_pct") and leaf_value(rule["size_pct"])
    size_pct = None
    if size_rule is not None:
        # The resolver reports buSzPct as a FRACTION (0.95 == 95%);
        # the profile schema carries percent (95).
        resolved_fraction = resolved_bullet.get("size_pct")
        resolved_pct = (None if resolved_fraction is None
                        else resolved_fraction * 100.0)
        if not _close(resolved_pct, float(size_rule), _PCT_TOLERANCE):
            size_pct = float(size_rule)
    color_rule = _leaf_hex(rule.get("color"))
    color_kwargs: Dict[str, str] = {}
    if (color_rule is not None
            and _hex_or_none(resolved_bullet.get("color_hex")) != color_rule):
        token = palette.tokens_by_hex.get(color_rule)
        color_kwargs = ({"color_scheme": token} if token is not None
                        else {"color_hex": color_rule})
    if char is None and size_pct is None and not color_kwargs:
        return
    set_paragraph_bullet(paragraph, char=char, size_pct=size_pct,
                         **color_kwargs)
    if char is not None:
        log.record(slide_number, shape_name, "bullet.char",
                   resolved_bullet.get("char"), char)
    if size_pct is not None:
        log.record(slide_number, shape_name, "bullet.size_pct",
                   resolved_bullet.get("size_pct"), size_pct)
    if color_kwargs:
        log.record(slide_number, shape_name, "bullet.color",
                   resolved_bullet.get("color_hex"), color_rule)


# ---------------------------------------------------------------------------
# Shape defaults (non-placeholder p:sp)
# ---------------------------------------------------------------------------

def _apply_border_rule(shape, line: Dict, border_rule: Dict,
                       palette: _PaletteView, log: _ChangeLog,
                       slide_number: int) -> None:
    """Fix weight/color/dash on shapes that visibly draw a border."""
    if not line["visible"]:
        return
    weight_rule = leaf_points(border_rule.get("weight"))
    weight = (weight_rule if weight_rule is not None and not _close(
        line["weight_pt"], weight_rule, _PT_TOLERANCE) else None)
    color_rule = _leaf_hex(border_rule.get("color"))
    color_kwargs: Dict[str, str] = {}
    if (color_rule is not None
            and _hex_or_none(line["color_hex"]) != color_rule):
        token = palette.tokens_by_hex.get(color_rule)
        color_kwargs = ({"color_scheme": token} if token is not None
                        else {"color_hex": color_rule})
    dash_rule = border_rule.get("dash") and leaf_value(border_rule["dash"])
    dash = dash_rule if dash_rule and line["dash"] != dash_rule else None
    if weight is None and not color_kwargs and dash is None:
        return
    set_shape_border(shape, weight_pt=weight, dash=dash, **color_kwargs)
    if weight is not None:
        log.record(slide_number, shape.name, "line.weight_pt",
                   line["weight_pt"], weight)
    if color_kwargs:
        log.record(slide_number, shape.name, "line.color",
                   line["color_hex"], color_rule)
    if dash is not None:
        log.record(slide_number, shape.name, "line.dash",
                   line["dash"], dash)


def _literal_srgb_fill_hex(shape) -> Optional[str]:
    """Hex of the shape's OWN untransformed ``a:srgbClr`` solid fill.

    Returns ``None`` for theme-linked (``a:schemeClr``, incl. tints via
    lumMod/lumOff children), transformed-srgb, inherited (style
    ``fillRef`` / no ``spPr`` fill) and non-solid fills. Only a literal
    plain srgb fill is ever a candidate for the off-palette fill snap:
    everything else is theme-derived, and repainting it as a hardcoded
    hex would destroy the theme link on conformant content (repair-pass
    defect: schemeClr-tint fills resolved off the 12 base hexes and were
    rewritten on every apply, breaking the second-apply no-op contract).
    """
    sp_pr = shape._element.find(qn("p:spPr"))
    solid = None if sp_pr is None else sp_pr.find(qn("a:solidFill"))
    srgb = None if solid is None else solid.find(qn("a:srgbClr"))
    if srgb is None or len(srgb):
        return None
    value = srgb.get("val")
    return None if value is None else value.upper()


def _apply_shape_defaults(shape, resolved: Dict, defaults: Dict,
                          palette: _PaletteView, log: _ChangeLog,
                          slide_number: int) -> None:
    border_rule = defaults.get("border") or {}
    if border_rule:
        _apply_border_rule(shape, resolved["line"], border_rule,
                           palette, log, slide_number)
    fill_rule = _leaf_hex(defaults.get("fill"))
    literal_fill = _literal_srgb_fill_hex(shape)
    if (fill_rule is not None and resolved["fill"]["type"] == "solid"
            and literal_fill is not None
            and literal_fill != fill_rule
            and literal_fill not in palette.hexes):
        set_shape_fill(shape, hex_value=fill_rule)
        log.record(slide_number, shape.name, "fill.color",
                   literal_fill, fill_rule)
    radius_rule = defaults.get("corner_radius")
    if radius_rule is not None and resolved["geometry"]["preset"] == "roundRect":
        rule_fraction = float(leaf_value(radius_rule))
        adjustments = resolved["adjustments"]
        current = adjustments[0] if adjustments else None
        if not _close(current, rule_fraction, _ADJ_TOLERANCE):
            set_shape_corner_radius(shape, rule_fraction)
            log.record(slide_number, shape.name, "corner_radius",
                       current, rule_fraction)


# ---------------------------------------------------------------------------
# Per-shape / per-slide walk
# ---------------------------------------------------------------------------

def _apply_text_rules(slide, shape, theme: ThemeContext, profile: Dict,
                      palette: _PaletteView, log: _ChangeLog,
                      slide_number: int) -> None:
    role = _role_of(shape)
    typography_rule = (profile.get("typography") or {}).get(role) or {}
    paragraph_rules = (profile.get("paragraph") or {}) if role == "body" else {}
    bullet_rules = paragraph_rules.get("bullets") or {}
    for paragraph in shape.text_frame.paragraphs:
        resolved_para = resolve_paragraph(
            TextContext(slide, shape, paragraph, None, theme))
        level = resolved_para["indent_level"]
        if paragraph_rules and level == 1:
            _apply_paragraph_spacing(paragraph, resolved_para,
                                     paragraph_rules, log, slide_number,
                                     shape.name)
        bullet_rule = bullet_rules.get(f"l{level}")
        if bullet_rule:
            _apply_bullet_rule(paragraph, resolved_para["bullet"],
                               bullet_rule, palette, log,
                               slide_number, shape.name)
        for run in paragraph.runs:
            effective = resolve_run_font(
                TextContext(slide, shape, paragraph, run, theme))
            if typography_rule:
                _apply_run_typography(run, effective, typography_rule,
                                      role, level, palette, log,
                                      slide_number, shape.name)
            _relink_hardcoded_srgb(run, palette, log, slide_number,
                                   shape.name)


def _apply_to_shape(slide, shape, theme: ThemeContext, profile: Dict,
                    palette: _PaletteView, log: _ChangeLog,
                    slide_number: int) -> None:
    if _has_text_body(shape):
        _apply_text_rules(slide, shape, theme, profile, palette,
                          log, slide_number)
    defaults = profile.get("shape_defaults") or {}
    if (defaults and not shape.is_placeholder
            and shape._element.tag == qn("p:sp")):
        resolved = resolve_shape(shape, slide, theme)
        _apply_shape_defaults(shape, resolved, defaults, palette,
                              log, slide_number)


def _validate_profile(profile: Dict) -> None:
    if not is_house_profile(profile):
        raise ValueError(
            "apply_house_profile requires a house-profile/1 dict "
            f"(schema_version={HOUSE_PROFILE_SCHEMA!r})")
    for section in ("typography", "paragraph", "palette", "shape_defaults"):
        value = profile.get(section)
        if value is not None and not isinstance(value, dict):
            raise ValueError(
                f"profile section {section!r} must be a dict, "
                f"got {type(value).__name__}")


def apply_house_profile(prs, profile: Dict) -> Dict:
    """Apply a house-profile/1's deterministic rules to a presentation.

    Returns a summary: total writes, slides touched and a capped
    per-change list. Geometry, masters, layouts and themes are never
    modified; already-conformant values receive no explicit overrides.
    """
    if prs is None:
        raise ValueError("prs must not be None")
    _validate_profile(profile)
    scheme_by_hex = _scheme_by_hex(profile)
    log = _ChangeLog()
    context_cache: Dict[str, tuple] = {}
    slides_touched = set()
    for index, slide in enumerate(prs.slides):
        master = slide.slide_layout.slide_master
        key = str(master.part.partname)
        context = context_cache.get(key)
        if context is None:
            theme = ThemeContext.for_master(master)
            context = (theme, _palette_view(scheme_by_hex, theme))
            context_cache[key] = context
        theme, palette = context
        writes_before = log.writes
        for shape in slide.shapes:
            _apply_to_shape(slide, shape, theme, profile, palette,
                            log, index + 1)
        if log.writes > writes_before:
            slides_touched.add(index + 1)
    return {
        "schema_version": HOUSE_PROFILE_SCHEMA,
        "profile_name": profile.get("name"),
        "writes": log.writes,
        "slides_touched": sorted(slides_touched),
        "changes": log.changes,
        "changes_truncated": log.writes > len(log.changes),
    }
