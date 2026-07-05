"""Generic OOXML style-cascade walker for the effective-style resolver.

The foundation of Step 2 of the style-fidelity plan: when python-pptx
returns ``None`` because a value is inherited, this module walks the
ordered list of places PowerPoint would look and returns the first
explicit value -- ONE walker (``resolve_first``), parameterized by a
per-property extractor callable, reused by every property. That is the
Apache POI XSLF architecture (``CharacterPropertyFetcher`` /
``ParagraphPropertyFetcher``, Apache-2.0 -- re-implemented, not copied);
the *contents* of each cascade step port ShapeCrawler's verified lookup
logic (MIT -- ``ReferencedFont``, ``SCPShapeTree``, ``IndentFonts``).

CHOSEN CASCADE ORDER (the ShapeCrawler-vs-POI precedence caveat,
documented per plan; the Step 0 fixture decks arbitrate -- whatever
order reproduces PowerPoint's COM-reported effective values is correct):

Placeholder shapes (``build_text_cascade_sources`` with a placeholder):
    1. run          ``a:rPr`` on the run itself
    2. paragraph    ``a:pPr`` of the run's paragraph (its ``a:defRPr``
                    for character properties)
    3. shape        the shape's own ``p:txBody/a:lstStyle``
    4. layout-placeholder   matched layout placeholder's ``a:lstStyle``
    5. master-placeholder   matched master placeholder's ``a:lstStyle``
       (matched via the *layout* placeholder when one exists, else via
       the slide placeholder -- ShapeCrawler ``ReferencedFont`` walk)
    6. master-txStyles      the master's ``p:txStyles`` style picked by
       placeholder type: title/ctrTitle -> ``p:titleStyle``;
       body/subTitle/obj/absent-type -> ``p:bodyStyle``;
       everything else (dt/ftr/sldNum/...) -> ``p:otherStyle``
    7. presentation-defaultTextStyle    ``p:defaultTextStyle`` on
       ``ppt/presentation.xml``
    8. theme-txDef  the governing master's theme
       ``a:objectDefaults/a:txDef/a:lstStyle``

Non-placeholder shapes (floating text boxes -- most shapes on real
consulting slides) get the shorter list: 1, 2, 3, 7, 8. The master's
``p:otherStyle`` is deliberately NOT consulted for slide shapes (it
styles shapes on the master itself); ECMA-376 intent routes floating
boxes through ``defaultTextStyle`` -> theme ``txDef``.

Indent levels are 1-based throughout (level = ``a:pPr/@lvl`` + 1),
matching ShapeCrawler. Within a list-style source, ``a:lvlNpPr`` is
consulted first, then ``a:defPPr`` as a documented fallback.

Placeholder matching ports ``SCPShapeTree`` semantics exactly: first
pass requires BOTH type and idx raw-attribute equality (absent matches
absent), second pass falls back to type-only and needs an explicit
target type. No title/ctrTitle aliasing is applied at the matching
layer -- ctrTitle placeholders that miss the master shape fall through
to ``p:titleStyle`` at step 6, which is how ShapeCrawler handles it.

Strictly read-only: no slide, layout, master or theme element is ever
mutated. All lookups run on live lxml elements via python-pptx's
``._element`` escape hatch.
"""

from dataclasses import dataclass
from typing import Callable, Iterable, List, Optional, Tuple, TypeVar

from lxml import etree
from pptx.oxml.ns import qn

from .resolve_theme import (
    ThemeContext,
    presentation_default_text_style,
)

T = TypeVar("T")

# ---------------------------------------------------------------------------
# Style sources
# ---------------------------------------------------------------------------

#: ``StyleSource.kind`` values.
KIND_RUN_PROPS = "runProps"      # element is an a:rPr / a:defRPr
KIND_PARA_PROPS = "paraProps"    # element is an a:pPr
KIND_LIST_STYLE = "listStyle"    # element holds a:lvlNpPr children
                                 # (a:lstStyle, p:titleStyle,
                                 #  p:defaultTextStyle, txDef lstStyle...)

_VALID_KINDS = frozenset({KIND_RUN_PROPS, KIND_PARA_PROPS, KIND_LIST_STYLE})

MIN_INDENT_LEVEL = 1
MAX_INDENT_LEVEL = 9


def _validate_level(level: int) -> None:
    if not isinstance(level, int) or not (
        MIN_INDENT_LEVEL <= level <= MAX_INDENT_LEVEL
    ):
        raise ValueError(
            f"indent level must be an int in "
            f"{MIN_INDENT_LEVEL}..{MAX_INDENT_LEVEL}, got {level!r}"
        )


@dataclass(frozen=True)
class StyleSource:
    """One stop on the inheritance walk (immutable).

    ``name`` identifies the cascade step for debugging/tracing (e.g.
    ``"layout-placeholder"``); ``kind`` says how to dig per-level
    properties out of ``element``.
    """

    name: str
    kind: str
    element: etree._Element

    def __post_init__(self):
        if self.kind not in _VALID_KINDS:
            raise ValueError(f"unknown StyleSource kind: {self.kind!r}")
        if self.element is None:
            raise ValueError(f"StyleSource {self.name!r} has no element")

    def paragraph_props(self, level: int) -> Optional[etree._Element]:
        """The pPr-like element for ``level``, or ``None``.

        * runProps sources have no paragraph properties.
        * paraProps sources return their element as-is (it IS the pPr).
        * listStyle sources look up ``a:lvl{level}pPr``, then fall back
          to ``a:defPPr`` (ShapeCrawler ``IndentFonts`` lvl1-9 walk,
          plus the documented defPPr fallback).
        """
        _validate_level(level)
        if self.kind == KIND_RUN_PROPS:
            return None
        if self.kind == KIND_PARA_PROPS:
            return self.element
        level_props = self.element.find(qn(f"a:lvl{level}pPr"))
        if level_props is not None:
            return level_props
        return self.element.find(qn("a:defPPr"))

    def run_props(self, level: int) -> Optional[etree._Element]:
        """The rPr/defRPr-like element for ``level``, or ``None``."""
        if self.kind == KIND_RUN_PROPS:
            _validate_level(level)
            return self.element
        paragraph_props = self.paragraph_props(level)
        if paragraph_props is None:
            return None
        return paragraph_props.find(qn("a:defRPr"))


# ---------------------------------------------------------------------------
# The one generic walker (POI CharacterPropertyFetcher pattern)
# ---------------------------------------------------------------------------

def resolve_first(
    sources: Iterable[StyleSource],
    extractor: Callable[[StyleSource], Optional[T]],
) -> Optional[T]:
    """Walk ``sources`` in order; return the first non-``None`` extract.

    ``extractor`` receives each ``StyleSource`` and returns the property
    value if that source defines it explicitly, else ``None``. This is
    the single cascade walker every resolved property goes through.
    """
    if extractor is None:
        raise ValueError("extractor must not be None")
    for source in sources:
        value = extractor(source)
        if value is not None:
            return value
    return None


def resolve_run_property(
    sources: Iterable[StyleSource],
    level: int,
    extract: Callable[[etree._Element], Optional[T]],
) -> Optional[T]:
    """Character-property convenience wrapper around ``resolve_first``.

    ``extract`` receives the rPr/defRPr-like element of each source that
    has one for ``level``.
    """
    _validate_level(level)

    def _extractor(source: StyleSource) -> Optional[T]:
        run_props = source.run_props(level)
        return extract(run_props) if run_props is not None else None

    return resolve_first(sources, _extractor)


def resolve_paragraph_property(
    sources: Iterable[StyleSource],
    level: int,
    extract: Callable[[etree._Element], Optional[T]],
) -> Optional[T]:
    """Paragraph-property convenience wrapper around ``resolve_first``."""
    _validate_level(level)

    def _extractor(source: StyleSource) -> Optional[T]:
        paragraph_props = source.paragraph_props(level)
        return extract(paragraph_props) if paragraph_props is not None else None

    return resolve_first(sources, _extractor)


# ---------------------------------------------------------------------------
# Placeholder matching (SCPShapeTree semantics)
# ---------------------------------------------------------------------------

def placeholder_of(sp_element: etree._Element) -> Optional[etree._Element]:
    """The ``p:ph`` element of a ``p:sp``, or ``None`` if not a placeholder."""
    if sp_element is None:
        raise ValueError("sp_element must not be None")
    return sp_element.find(
        qn("p:nvSpPr") + "/" + qn("p:nvPr") + "/" + qn("p:ph")
    )


def _placeholder_shapes(
    shape_tree: etree._Element,
) -> List[Tuple[etree._Element, etree._Element]]:
    """(sp, ph) pairs for every direct-child placeholder ``p:sp``."""
    pairs = []
    for sp_element in shape_tree.findall(qn("p:sp")):
        ph_element = placeholder_of(sp_element)
        if ph_element is not None:
            pairs.append((sp_element, ph_element))
    return pairs


def match_placeholder_in_tree(
    target_ph: etree._Element, shape_tree: etree._Element
) -> Optional[etree._Element]:
    """Find the layout/master ``p:sp`` a slide placeholder inherits from.

    Ported SCPShapeTree semantics:
      1. match on BOTH type and idx (raw attribute values; an absent
         attribute only matches an absent attribute),
      2. fall back to type-only, which requires the target to carry an
         explicit type.
    Returns the matched ``p:sp`` element or ``None``.
    """
    if target_ph is None:
        raise ValueError("target_ph must not be None")
    if shape_tree is None:
        raise ValueError("shape_tree must not be None")
    target_type = target_ph.get("type")
    target_idx = target_ph.get("idx")
    candidates = _placeholder_shapes(shape_tree)

    for sp_element, ph_element in candidates:
        if (ph_element.get("type") == target_type
                and ph_element.get("idx") == target_idx):
            return sp_element

    if target_type is None:
        return None
    for sp_element, ph_element in candidates:
        if ph_element.get("type") == target_type:
            return sp_element
    return None


# ---------------------------------------------------------------------------
# Master txStyles selection
# ---------------------------------------------------------------------------

#: Placeholder types routed to p:titleStyle / p:bodyStyle; everything
#: else (dt, ftr, sldNum, pic, tbl, chart, media, clipArt, ...) gets
#: p:otherStyle. An absent type attribute means "body" per ECMA-376.
_TITLE_PH_TYPES = frozenset({"title", "ctrTitle"})
_BODY_PH_TYPES = frozenset({"body", "subTitle", "obj"})


def master_text_style_for_ph_type(
    master_element: etree._Element, ph_type: Optional[str]
) -> Optional[etree._Element]:
    """The ``p:txStyles`` child of a master that styles a placeholder type.

    Returns the ``p:titleStyle`` / ``p:bodyStyle`` / ``p:otherStyle``
    element (each holds ``a:lvlNpPr`` children), or ``None`` when the
    master has no ``p:txStyles`` block.
    """
    if master_element is None:
        raise ValueError("master_element must not be None")
    tx_styles = master_element.find(qn("p:txStyles"))
    if tx_styles is None:
        return None
    if ph_type in _TITLE_PH_TYPES:
        tag = "p:titleStyle"
    elif ph_type is None or ph_type in _BODY_PH_TYPES:
        tag = "p:bodyStyle"
    else:
        tag = "p:otherStyle"
    return tx_styles.find(qn(tag))


# ---------------------------------------------------------------------------
# Cascade source builders
# ---------------------------------------------------------------------------

def indent_level_of(paragraph) -> int:
    """1-based indent level of a python-pptx paragraph (lvl attr + 1).

    Reads ``a:pPr/@lvl`` straight off the paragraph XML. Deliberately
    NOT ``paragraph.level``: python-pptx implements that property via
    ``get_or_add_pPr()``, which INSERTS an empty ``<a:pPr/>`` into the
    slide part when the paragraph has none -- a mutation this module's
    strict read-only guarantee forbids.
    """
    if paragraph is None:
        raise ValueError("paragraph must not be None")
    p_pr = paragraph._p.find(qn("a:pPr"))
    if p_pr is None:
        return MIN_INDENT_LEVEL
    lvl = p_pr.get("lvl")
    if lvl is None:
        return MIN_INDENT_LEVEL
    return int(lvl) + 1


def _list_style_source(name: str, container, tag_path: str) -> Optional[StyleSource]:
    """Build a listStyle StyleSource from ``container.find(tag_path)``."""
    if container is None:
        return None
    element = container.find(tag_path)
    if element is None:
        return None
    return StyleSource(name=name, kind=KIND_LIST_STYLE, element=element)


def _txbody_list_style_source(
    name: str, sp_element: Optional[etree._Element]
) -> Optional[StyleSource]:
    """listStyle source from a ``p:sp``'s ``p:txBody/a:lstStyle``."""
    if sp_element is None:
        return None
    return _list_style_source(
        name, sp_element, qn("p:txBody") + "/" + qn("a:lstStyle")
    )


def _placeholder_chain_sources(shape, slide) -> List[StyleSource]:
    """Cascade steps 4-6 for a placeholder shape (layout/master/txStyles)."""
    sources: List[StyleSource] = []
    sp_element = shape._element
    slide_ph = placeholder_of(sp_element)
    layout = slide.slide_layout
    master = layout.slide_master

    layout_tree = layout.shapes._spTree
    layout_sp = match_placeholder_in_tree(slide_ph, layout_tree)
    layout_source = _txbody_list_style_source("layout-placeholder", layout_sp)
    if layout_source is not None:
        sources.append(layout_source)

    # Master placeholder is matched via the layout shape's ph when one
    # exists (ShapeCrawler ReferencedFont), else via the slide's ph.
    master_lookup_ph = slide_ph
    if layout_sp is not None:
        layout_ph = placeholder_of(layout_sp)
        if layout_ph is not None:
            master_lookup_ph = layout_ph
    master_tree = master.shapes._spTree
    master_sp = match_placeholder_in_tree(master_lookup_ph, master_tree)
    master_source = _txbody_list_style_source("master-placeholder", master_sp)
    if master_source is not None:
        sources.append(master_source)

    tx_style = master_text_style_for_ph_type(
        master.element, slide_ph.get("type")
    )
    if tx_style is not None:
        sources.append(StyleSource(
            name="master-txStyles", kind=KIND_LIST_STYLE, element=tx_style,
        ))
    return sources


def _deck_default_sources(
    slide, theme: Optional[ThemeContext] = None
) -> List[StyleSource]:
    """Cascade steps 7-8: presentation defaults, then theme txDef.

    ``theme`` is the prebuilt ``ThemeContext`` of the slide's governing
    master; when ``None`` one is built here (a full theme-part parse),
    so callers resolving more than one property should thread one in.
    """
    sources: List[StyleSource] = []
    default_text_style = presentation_default_text_style(slide)
    if default_text_style is not None:
        sources.append(StyleSource(
            name="presentation-defaultTextStyle",
            kind=KIND_LIST_STYLE,
            element=default_text_style,
        ))
    if theme is None:
        theme = ThemeContext.for_slide(slide)
    tx_def_style = theme.text_default_list_style()
    if tx_def_style is not None:
        sources.append(StyleSource(
            name="theme-txDef", kind=KIND_LIST_STYLE, element=tx_def_style,
        ))
    return sources


def build_text_cascade_sources(
    slide, shape, paragraph=None, run=None,
    theme: Optional[ThemeContext] = None,
) -> Tuple[StyleSource, ...]:
    """Ordered cascade sources for text inside ``shape`` on ``slide``.

    ``shape`` / ``slide`` / ``paragraph`` / ``run`` are python-pptx
    objects; ``paragraph`` and ``run`` are optional -- when provided,
    their ``a:pPr`` / ``a:rPr`` lead the list. Placeholder shapes get
    the full 8-step chain, floating shapes the short one (see module
    docstring for the documented order). The returned tuple is what
    ``resolve_run_property`` / ``resolve_paragraph_property`` walk.

    ``theme`` is an optional prebuilt ``ThemeContext`` for the slide's
    governing master; passing it avoids re-parsing the theme part for
    the txDef tail on every cascade build (walk-the-deck callers must
    thread it through).
    """
    if slide is None or shape is None:
        raise ValueError("slide and shape must not be None")
    if theme is not None and not isinstance(theme, ThemeContext):
        raise ValueError(
            f"theme must be a ThemeContext or None, got {type(theme)!r}"
        )
    sources: List[StyleSource] = []

    if run is not None:
        rpr_element = run._r.find(qn("a:rPr"))
        if rpr_element is not None:
            sources.append(StyleSource(
                name="run", kind=KIND_RUN_PROPS, element=rpr_element,
            ))
    if paragraph is not None:
        ppr_element = paragraph._p.find(qn("a:pPr"))
        if ppr_element is not None:
            sources.append(StyleSource(
                name="paragraph", kind=KIND_PARA_PROPS, element=ppr_element,
            ))

    shape_source = _txbody_list_style_source("shape-lstStyle", shape._element)
    if shape_source is not None:
        sources.append(shape_source)

    is_placeholder = placeholder_of(shape._element) is not None
    if is_placeholder:
        sources.extend(_placeholder_chain_sources(shape, slide))
    sources.extend(_deck_default_sources(slide, theme))
    return tuple(sources)
