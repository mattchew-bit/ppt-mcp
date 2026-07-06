"""Style-lint engine (Step 5 of the style-fidelity plan).

``lint_against_profile(deck_path, profile)`` compares the
inheritance-RESOLVED analysis of a deck (``utils.resolve_analysis`` --
resolution is never reimplemented here) against a ``house-profile/1``
dict (``utils.profile_schema``) and returns an ordered list of finding
dicts::

    {rule_id, severity, slide, shape, shape_id, paragraph?, run?,
     property, expected, actual, message}

* ``severity`` is ``error`` / ``warn`` / ``info``; findings are ordered
  by severity, then slide, then rule_id (stable within).
* Messages are DISTRIBUTION-STYLE wherever the profile carries a
  learned distribution ("body run is 13pt; house type scale is
  {11, 14, 20, 30}pt"), never bare pass/fail.
* ``paragraph`` / ``run`` refs are 1-based and present only where they
  apply.

The rule catalog lives in ``utils.lint_rules`` (registry) with
implementations split across ``lint_rules`` (text/style rules) and
``lint_rules_geometry`` (geometry/media rules). Everything is strictly
READ-ONLY on the deck.
"""

from functools import cached_property
from pathlib import Path
from typing import Any, Dict, Iterator, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .profile_archetypes import FOOTER_FRACTION, classify_slide
from .profile_grid import POINTS_PER_INCH
from .resolve_analysis import build_resolved_analysis
from .style_apply import is_house_profile, leaf_points, leaf_value
from .style_roles import ph_type_label_name, placeholder_role

#: Internal resolver cap -- the lint engine needs the FULL deck.
_RESOLVE_MAX_BYTES = 64_000_000

SEVERITY_RANK = {"error": 0, "warn": 1, "info": 2}

#: Numeric comparison tolerances (points / unitless multiple).
SIZE_TOLERANCE_PT = 0.25
SPACING_TOLERANCE_PT = 0.25
LINE_SPACING_TOLERANCE = 0.011
BORDER_TOLERANCE_PT = 0.05

#: Body size used when the profile carries no body size at all.
_FALLBACK_BODY_SIZE_PT = 18.0


def make_finding(rule_id: str, severity: str, slide: int,
                 shape: Optional[str], prop: str, expected: Any,
                 actual: Any, message: str,
                 shape_id: Optional[int] = None,
                 paragraph: Optional[int] = None,
                 run: Optional[int] = None) -> Dict:
    """One lint finding; ``paragraph``/``run`` keys only when given."""
    if severity not in SEVERITY_RANK:
        raise ValueError(f"unknown severity {severity!r}")
    finding: Dict[str, Any] = {
        "rule_id": rule_id, "severity": severity, "slide": slide,
        "shape": shape, "shape_id": shape_id, "property": prop,
        "expected": expected, "actual": actual, "message": message,
    }
    if paragraph is not None:
        finding["paragraph"] = paragraph
    if run is not None:
        finding["run"] = run
    return finding


def format_scale(values: Sequence, unit: str = "") -> str:
    """``[11.0, 14.0]`` -> ``"{11, 14}pt"`` (distribution message)."""
    def fmt(value):
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)

    return "{" + ", ".join(fmt(v) for v in values) + "}" + unit


# ---------------------------------------------------------------------------
# Deck facts (resolver output + live shape references)
# ---------------------------------------------------------------------------

def _is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except NotImplementedError:  # exotic shapes report no type
        return False


def collect_deck_facts(deck_path: str) -> Dict:
    """One deck -> resolved slide facts with live shape references.

    Returns ``{"slides": [...], "themes": [...], "width_pt",
    "height_pt"}``. Each slide record matches the structure the Step 3
    profile layers consume (deck / slide_number / width_pt / height_pt /
    shapes) so ``classify_slide`` and the grid helpers run unchanged;
    each shape record additionally carries ``is_picture``, ``shape_id``
    and the non-serializable ``_shape`` (the python-pptx shape, for
    raw-XML probes). Strictly read-only.
    """
    path = Path(deck_path)
    if not path.is_file():
        raise FileNotFoundError(f"deck not found: {deck_path}")
    prs = Presentation(str(path))
    report = build_resolved_analysis(prs, detail="full",
                                     max_bytes=_RESOLVE_MAX_BYTES)
    if report.get("truncated"):
        raise ValueError(
            f"resolved analysis of {path.name} was truncated -- the "
            "deck is too large to lint in one pass"
        )
    width_pt, height_pt = prs.slide_width.pt, prs.slide_height.pt
    slides: List[Dict] = []
    for slide_record, slide in zip(report["slides"], prs.slides):
        shapes = []
        for record, shape in zip(slide_record["shapes"], slide.shapes):
            merged = dict(record)  # never mutate the resolver's output
            merged["is_picture"] = _is_picture(shape)
            merged["shape_id"] = shape.shape_id
            merged["_shape"] = shape
            shapes.append(merged)
        slides.append({
            "deck": path.name,
            "slide_number": slide_record["slide_number"],
            "width_pt": width_pt,
            "height_pt": height_pt,
            "shapes": shapes,
        })
    return {"slides": slides, "themes": report.get("themes", []),
            "width_pt": width_pt, "height_pt": height_pt}


# ---------------------------------------------------------------------------
# Lint context (profile sections pre-digested for the rules)
# ---------------------------------------------------------------------------

def _norm_hex(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    return str(value).lstrip("#").upper()


class LintContext:
    """Everything a rule needs: deck facts + pre-digested profile."""

    def __init__(self, facts: Dict, profile: Dict):
        self.slides: List[Dict] = facts["slides"]
        self.width_pt: float = facts["width_pt"]
        self.height_pt: float = facts["height_pt"]
        self.profile = profile

        palette = profile.get("palette", {})
        #: scheme token -> RRGGBB, in the profile's canonical order.
        self.scheme: Dict[str, str] = {
            token: _norm_hex(leaf_value(leaf))
            for token, leaf in palette.get("scheme", {}).items()
        }
        self.scheme_hexes = set(self.scheme.values())
        self.usage_hexes = {
            _norm_hex(entry["color"])
            for entry in palette.get("usage", []) if entry.get("color")
        }

        typography = profile.get("typography", {})
        self.house_fonts = sorted({
            leaf_value(spec["font"]) for spec in typography.values()
            if "font" in spec
        })
        body_size = typography.get("body", {}).get("size")
        self.body_size_pt = (leaf_points(body_size) if body_size
                             else _FALLBACK_BODY_SIZE_PT)

        self.paragraph_rules: Dict = profile.get("paragraph", {})
        self.border_rule: Optional[Dict] = (
            profile.get("shape_defaults", {}).get("border"))

        grid = profile.get("grid", {})
        self.grid_edges: Dict[str, List[float]] = {
            family: list(edges)
            for family, edges in grid.get("edges", {}).items()
        }
        tolerance = grid.get("tolerance")
        self.grid_tolerance_in: Optional[float] = (
            float(leaf_value(tolerance)) if tolerance else None)

        self.archetypes: Dict = profile.get("archetypes", {})
        distributions = profile.get("distributions", {})
        self.font_size_scale: List[float] = [
            float(v) for v in distributions.get(
                "font_sizes_pt", {}).get("values", [])
        ]
        self.space_after_scale: List[float] = [
            float(v) for v in distributions.get(
                "space_after_pt", {}).get("values", [])
        ]

    # -- shared derived views ------------------------------------------------

    def token_for_hex(self, hex_value: str) -> Optional[str]:
        """First scheme token (canonical order) carrying ``hex_value``."""
        wanted = _norm_hex(hex_value)
        for token, value in self.scheme.items():
            if value == wanted:
                return token
        return None

    def footer_cutoff_pt(self, slide: Dict) -> float:
        return FOOTER_FRACTION * slide["height_pt"]

    def in_footer_zone(self, slide: Dict, shape: Dict) -> bool:
        top = shape["geometry"].get("top_pt")
        return top is not None and top >= self.footer_cutoff_pt(slide)

    def shape_role(self, slide: Dict, shape: Dict) -> Optional[str]:
        """title/body/footer role, mirroring the profile's learn side."""
        role = placeholder_role(ph_type_label_name(shape.get("ph_type")))
        if role == "title":
            return "title"
        if self.in_footer_zone(slide, shape):
            return "footer"
        return role

    @cached_property
    def archetype_by_slide(self) -> Dict[int, str]:
        """slide_number -> closed-set archetype (computed once)."""
        return {
            slide["slide_number"]: classify_slide(slide, self.body_size_pt)
            for slide in self.slides
        }

    def all_grid_lines_in(self) -> List[float]:
        """Every learned grid line, all families merged (inches)."""
        return sorted(
            line for edges in self.grid_edges.values() for line in edges)


# ---------------------------------------------------------------------------
# Shared iteration helpers for rules
# ---------------------------------------------------------------------------

def iter_shapes(ctx: LintContext) -> Iterator[Tuple[Dict, Dict]]:
    for slide in ctx.slides:
        for shape in slide["shapes"]:
            yield slide, shape


def iter_runs(ctx: LintContext) -> Iterator[Tuple[Dict, Dict, int, int,
                                                  Dict]]:
    """Yield ``(slide, shape, paragraph_1based, run_1based, run)``."""
    for slide, shape in iter_shapes(ctx):
        for p_idx, paragraph in enumerate(shape.get("paragraphs", []),
                                          start=1):
            for r_idx, run in enumerate(paragraph.get("runs", []),
                                        start=1):
                yield slide, shape, p_idx, r_idx, run


def shape_has_text(shape: Dict) -> bool:
    return any(
        run.get("text", "").strip()
        for paragraph in shape.get("paragraphs", [])
        for run in paragraph.get("runs", [])
    )


def rect_of(shape: Dict) -> Optional[Tuple[float, float, float, float]]:
    """(left, top, width, height) in points, or None."""
    geometry = shape["geometry"]
    values = tuple(geometry.get(key) for key in
                   ("left_pt", "top_pt", "width_pt", "height_pt"))
    if any(value is None for value in values):
        return None
    return values  # type: ignore[return-value]


def pt_to_in(value: float) -> float:
    return value / POINTS_PER_INCH


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _strip_internal(finding: Dict) -> Dict:
    return {key: value for key, value in finding.items()
            if not key.startswith("_")}


def lint_against_profile(deck_path: str, profile: Dict) -> List[Dict]:
    """Lint a deck against a ``house-profile/1`` dict.

    Returns the ordered findings list (see module docstring). Raises
    ``FileNotFoundError`` for a missing deck and ``ValueError`` for a
    profile that is not a house-profile/1 dict.
    """
    if not is_house_profile(profile):
        detail = (f"schema_version {profile.get('schema_version')!r}"
                  if isinstance(profile, dict)
                  else f"a {type(profile).__name__}")
        raise ValueError(
            "lint_against_profile requires a house-profile/1 dict "
            f"(schema_version 'house-profile/1'); got {detail}"
        )
    from .lint_rules import RULES  # late import: rules import this module

    facts = collect_deck_facts(deck_path)
    ctx = LintContext(facts, profile)
    findings: List[Dict] = []
    for rule in RULES:
        findings.extend(rule.fn(ctx))
    findings.sort(key=lambda f: (SEVERITY_RANK[f["severity"]],
                                 f["slide"], f["rule_id"],
                                 f.get("shape") or "",
                                 f.get("paragraph") or 0,
                                 f.get("run") or 0))
    return [_strip_internal(finding) for finding in findings]
