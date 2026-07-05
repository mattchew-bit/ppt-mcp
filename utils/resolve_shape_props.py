"""Effective shape properties: line, fill, geometry, ``p:style`` refs.

Resolves what PowerPoint actually draws for a shape when python-pptx
returns ``None`` -- covering the two inheritance mechanisms shapes have:

* ``p:style`` references (``a:lnRef``/``a:fillRef``/``a:fontRef``):
  1-based indexes into the governing theme's ``a:fmtScheme`` style lists,
  with the reference's own color substituted wherever the theme entry
  says ``a:schemeClr val="phClr"`` (ECMA-376 20.1.4.2.10). ``fillRef``
  additionally addresses the background-fill list with the +1000 offset
  convention (idx 1001 -> ``a:bgFillStyleLst`` entry 1).
* explicit ``p:spPr`` children, which override the referenced style
  attribute-by-attribute (an explicit ``a:ln`` with only ``noFill``
  still inherits width/dash from the ``lnRef`` entry).

Geometry: placeholders without their own ``a:xfrm`` inherit position and
size from the matched layout placeholder, then the master placeholder
(same SCPShapeTree matching the text cascade uses). A missing
``a:prstGeom`` resolves to the ECMA default ``rect``. Adjustment values
are returned as fractions (``a:gd fmla="val 32000"`` -> 0.32, matching
what the COM object model reports); presets whose ``a:avLst`` is empty
fall back to the preset's documented default adjustments
(``PRESET_DEFAULT_ADJUSTMENTS`` -- deliberately small, fixture-arbitrated
table, unknown presets return no adjustments).

Ported semantics (re-implemented, not copied): ShapeCrawler ``ShapeColor``
(MIT) for reading the ``fontRef`` color; Apache POI ``XSLFSimpleShape``
(Apache-2.0) for the style-matrix indexing and phClr substitution.

Strictly read-only: nothing here mutates the presentation.
"""

from typing import Dict, List, Optional, Tuple

from lxml import etree
from pptx.oxml.ns import qn

from .resolve_colors import (
    PERCENT_DENOMINATOR,
    apply_color_transforms,
    find_color_child,
)
from .resolve_core import match_placeholder_in_tree, placeholder_of
from .resolve_theme import ThemeContext

#: English Metric Units per point.
EMU_PER_POINT = 12700.0

#: 60000ths of a degree per degree (xfrm/@rot).
ROTATION_DENOMINATOR = 60000.0

#: ECMA-376 default line width when neither spPr nor the theme says.
DEFAULT_LINE_WIDTH_PT = 0.75

#: fillRef idx >= 1000 addresses a:bgFillStyleLst (idx - 1000, 1-based).
BACKGROUND_FILL_IDX_OFFSET = 1000

#: Fill child local-name -> resolver fill type token.
_FILL_TYPES: Dict[str, str] = {
    "noFill": "none",
    "solidFill": "solid",
    "gradFill": "gradient",
    "blipFill": "picture",
    "pattFill": "pattern",
    "grpFill": "group",
}

#: Default adjustment fractions for presets whose a:avLst is empty.
#: Values are the ECMA-376 preset-definition defaults (val/100000);
#: fixture-arbitrated (chevron: PowerPoint COM reports [0.5]). Presets
#: not listed return [] -- documented limitation, extend as fixtures
#: demand.
PRESET_DEFAULT_ADJUSTMENTS: Dict[str, List[float]] = {
    "roundRect": [0.16667],
    "chevron": [0.5],
    "homePlate": [0.5],
    "snip1Rect": [0.16667],
    "snip2SameRect": [0.16667, 0.0],
    "round2SameRect": [0.16667, 0.0],
}

#: ``a:fontRef/@idx`` -> theme font token (see resolve_theme).
_FONT_REF_TOKENS = {"major": "+mj-lt", "minor": "+mn-lt"}


def _require(value, name: str) -> None:
    if value is None:
        raise ValueError(f"{name} must not be None")


def _local_name(element: etree._Element) -> str:
    tag = element.tag
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _int_attr(element: etree._Element, attr: str) -> int:
    """Required integer attribute; explicit ValueError when absent.

    Never trust file content: a malformed deck may drop a mandatory
    attribute (e.g. ``a:off`` without ``x``), which must surface as a
    clear ``ValueError``, not a bare ``TypeError`` from ``int(None)``.
    """
    raw = element.get(attr)
    if raw is None:
        raise ValueError(
            f"<{_local_name(element)}> element is missing required "
            f"attribute {attr!r}"
        )
    return int(raw)


def _sp_pr(shape) -> Optional[etree._Element]:
    return shape._element.find(qn("p:spPr"))


# ---------------------------------------------------------------------------
# p:style references
# ---------------------------------------------------------------------------

def style_ref(shape, tag: str) -> Optional[etree._Element]:
    """The shape's ``p:style`` child ``tag`` (e.g. ``a:lnRef``), if any."""
    _require(shape, "shape")
    style = shape._element.find(qn("p:style"))
    if style is None:
        return None
    return style.find(qn(tag))


def _ref_idx(ref: Optional[etree._Element]) -> int:
    if ref is None:
        return 0
    raw = ref.get("idx")
    return 0 if raw is None else int(raw)


def _ref_color_hex(
    ref: Optional[etree._Element], theme: ThemeContext
) -> Optional[str]:
    """The reference's own color child resolved to hex (the phClr value)."""
    if ref is None:
        return None
    color_child = find_color_child(ref)
    if color_child is None:
        return None
    return theme.resolve_color(color_child)


def resolve_color_with_ph(
    color_element: etree._Element,
    theme: ThemeContext,
    ph_hex: Optional[str],
) -> str:
    """Resolve a color element, substituting ``phClr`` with ``ph_hex``."""
    _require(color_element, "color_element")
    if (
        _local_name(color_element) == "schemeClr"
        and color_element.get("val") == "phClr"
    ):
        if ph_hex is None:
            raise ValueError(
                "schemeClr val='phClr' encountered with no style-reference "
                "color to substitute"
            )
        return apply_color_transforms(ph_hex, list(color_element))
    return theme.resolve_color(color_element)


def font_ref_color_hex(shape, theme: ThemeContext) -> Optional[str]:
    """``p:style/a:fontRef`` color -> hex, or ``None`` (ShapeCrawler
    ``ShapeColor`` semantics)."""
    return _ref_color_hex(style_ref(shape, "a:fontRef"), theme)


def font_ref_typeface(shape, theme: ThemeContext) -> Optional[str]:
    """``p:style/a:fontRef/@idx`` -> concrete theme typeface, or ``None``."""
    ref = style_ref(shape, "a:fontRef")
    if ref is None:
        return None
    token = _FONT_REF_TOKENS.get(ref.get("idx", ""))
    return None if token is None else theme.theme_font(token)


# ---------------------------------------------------------------------------
# Fill
# ---------------------------------------------------------------------------

def _find_fill_child(parent: Optional[etree._Element]) -> Optional[etree._Element]:
    """First ECMA fill child (noFill/solidFill/gradFill/...), if any."""
    if parent is None:
        return None
    for child in parent:
        if _local_name(child) in _FILL_TYPES:
            return child
    return None


def _fill_info(
    fill_element: etree._Element,
    theme: ThemeContext,
    ph_hex: Optional[str] = None,
) -> Dict:
    """A fill child element -> ``{"type": ..., "color_hex": ...}``."""
    fill_type = _FILL_TYPES[_local_name(fill_element)]
    info: Dict = {"type": fill_type, "color_hex": None}
    if fill_type == "solid":
        color_child = find_color_child(fill_element)
        if color_child is not None:
            info["color_hex"] = resolve_color_with_ph(
                color_child, theme, ph_hex
            )
    return info


def _fill_ref_entry(
    shape, theme: ThemeContext
) -> Tuple[Optional[etree._Element], Optional[str]]:
    """The theme fill element a ``fillRef`` points at, plus its phClr hex."""
    ref = style_ref(shape, "a:fillRef")
    idx = _ref_idx(ref)
    if idx <= 0:
        return None, None
    if idx > BACKGROUND_FILL_IDX_OFFSET:
        styles = theme.background_fill_styles()
        position = idx - BACKGROUND_FILL_IDX_OFFSET - 1
    else:
        styles = theme.fill_styles()
        position = idx - 1
    if position >= len(styles):
        raise ValueError(
            f"fillRef idx={idx} exceeds the theme's format-scheme fill list "
            f"({len(styles)} entries)"
        )
    return styles[position], _ref_color_hex(ref, theme)


def resolve_shape_fill(shape, theme: ThemeContext) -> Dict:
    """Effective shape fill: ``{"type": token, "color_hex": hex|None}``.

    Explicit ``p:spPr`` fill wins; otherwise the ``fillRef`` entry from
    the theme's ``a:fillStyleLst``/``a:bgFillStyleLst`` with phClr
    substituted; otherwise ``none`` (shapes with neither -- e.g. text
    placeholders -- draw no fill).
    """
    _require(theme, "theme")
    explicit = _find_fill_child(_sp_pr(shape))
    if explicit is not None:
        return _fill_info(explicit, theme)
    entry, ph_hex = _fill_ref_entry(shape, theme)
    if entry is not None:
        return _fill_info(entry, theme, ph_hex)
    return {"type": "none", "color_hex": None}


# ---------------------------------------------------------------------------
# Line
# ---------------------------------------------------------------------------

def _line_ref_entry(
    shape, theme: ThemeContext
) -> Tuple[Optional[etree._Element], Optional[str]]:
    """The theme ``a:ln`` a ``lnRef`` points at, plus its phClr hex."""
    ref = style_ref(shape, "a:lnRef")
    idx = _ref_idx(ref)
    if idx <= 0:
        return None, None
    styles = theme.line_styles()
    if idx > len(styles):
        raise ValueError(
            f"lnRef idx={idx} exceeds the theme's format-scheme line list "
            f"({len(styles)} entries)"
        )
    return styles[idx - 1], _ref_color_hex(ref, theme)


def _line_attribute(explicit_ln, theme_ln, getter):
    """Attribute-wise merge: explicit ``a:ln`` first, then the theme's."""
    for candidate in (explicit_ln, theme_ln):
        if candidate is None:
            continue
        value = getter(candidate)
        if value is not None:
            return value
    return None


def resolve_shape_line(shape, theme: ThemeContext) -> Dict:
    """Effective outline: visibility, weight, dash, color.

    ``{"visible": bool, "weight_pt": float|None, "dash": token|None,
    "color_hex": hex|None}``. The explicit ``p:spPr/a:ln`` overrides the
    ``lnRef`` theme entry attribute-by-attribute (ECMA style-matrix
    semantics): an explicit ``noFill`` hides a themed line, an explicit
    fill with no ``@w`` still takes the themed width.
    """
    _require(theme, "theme")
    sp_pr = _sp_pr(shape)
    explicit_ln = None if sp_pr is None else sp_pr.find(qn("a:ln"))
    theme_ln, ph_hex = _line_ref_entry(shape, theme)
    if explicit_ln is None and theme_ln is None:
        return {"visible": False, "weight_pt": None, "dash": None,
                "color_hex": None}

    fill_child = _line_attribute(explicit_ln, theme_ln, _find_fill_child)
    if fill_child is None or _FILL_TYPES[_local_name(fill_child)] == "none":
        return {"visible": False, "weight_pt": None, "dash": None,
                "color_hex": None}

    color_hex = None
    color_child = find_color_child(fill_child)
    if color_child is not None:
        color_hex = resolve_color_with_ph(color_child, theme, ph_hex)

    raw_width = _line_attribute(explicit_ln, theme_ln,
                                lambda ln: ln.get("w"))
    weight_pt = (DEFAULT_LINE_WIDTH_PT if raw_width is None
                 else int(raw_width) / EMU_PER_POINT)

    dash_element = _line_attribute(explicit_ln, theme_ln,
                                   lambda ln: ln.find(qn("a:prstDash")))
    dash = "solid" if dash_element is None else dash_element.get("val")

    return {"visible": True, "weight_pt": weight_pt, "dash": dash,
            "color_hex": color_hex}


# ---------------------------------------------------------------------------
# Geometry
# ---------------------------------------------------------------------------

#: Geometry-only placeholder-type aliases: when the strict SCPShapeTree
#: match (type+idx, then type-only) fails, retry with these equivalent
#: types. An absent type means "body" (ECMA-376 19.3.1.36) -- a slide /
#: layout content placeholder written as ``<p:ph idx="1"/>`` inherits its
#: position from the master's ``type="body" idx="1"`` shape
#: (fixture-verified against COM-reported geometry). The text cascade
#: deliberately does NOT alias (master txStyles covers those steps).
_GEOMETRY_PH_TYPE_ALIASES = {
    None: ("body",),
    "subTitle": ("body",),
    "ctrTitle": ("title",),
    "title": ("ctrTitle",),
}


def _match_placeholder_for_geometry(ph, tree) -> Optional[etree._Element]:
    """SCPShapeTree match, then the geometry-only type-alias retry."""
    matched = match_placeholder_in_tree(ph, tree)
    if matched is not None:
        return matched
    target_idx = ph.get("idx")
    for alias in _GEOMETRY_PH_TYPE_ALIASES.get(ph.get("type"), ()):
        for sp in tree.findall(qn("p:sp")):
            candidate = placeholder_of(sp)
            if (candidate is not None
                    and candidate.get("type") == alias
                    and candidate.get("idx") == target_idx):
                return sp
    return None


def _inherited_xfrm(shape, slide) -> Optional[etree._Element]:
    """The shape's ``a:xfrm``, walking layout/master for placeholders."""
    sp_pr = _sp_pr(shape)
    if sp_pr is not None:
        xfrm = sp_pr.find(qn("a:xfrm"))
        if xfrm is not None:
            return xfrm
    ph = placeholder_of(shape._element)
    if ph is None or slide is None:
        return None
    layout = slide.slide_layout
    lookup_ph = ph
    for tree in (layout.shapes._spTree, layout.slide_master.shapes._spTree):
        matched = _match_placeholder_for_geometry(lookup_ph, tree)
        if matched is None:
            continue
        matched_sp_pr = matched.find(qn("p:spPr"))
        if matched_sp_pr is not None:
            xfrm = matched_sp_pr.find(qn("a:xfrm"))
            if xfrm is not None:
                return xfrm
        # Continue the walk via the matched shape's own ph (a layout
        # placeholder may carry a type the slide's ph omits).
        matched_ph = placeholder_of(matched)
        if matched_ph is not None:
            lookup_ph = matched_ph
    return None


def _adjustments(geom: etree._Element, preset: str) -> List[float]:
    """Explicit ``a:gd`` fractions, else the preset's documented defaults."""
    av_lst = geom.find(qn("a:avLst"))
    values: List[float] = []
    if av_lst is not None:
        for gd in av_lst.findall(qn("a:gd")):
            formula = gd.get("fmla", "")
            parts = formula.split()
            if len(parts) == 2 and parts[0] == "val":
                values.append(int(parts[1]) / PERCENT_DENOMINATOR)
    if values:
        return values
    return list(PRESET_DEFAULT_ADJUSTMENTS.get(preset, []))


def resolve_shape_geometry(shape, slide=None) -> Dict:
    """Effective geometry: preset, position/size in points, rotation.

    ``slide`` enables placeholder xfrm inheritance (layout -> master);
    without it, non-inheriting shapes still resolve. Missing
    ``a:prstGeom`` resolves to ``rect`` (ECMA default); a ``a:custGeom``
    reports preset ``custom`` with no adjustments.
    """
    _require(shape, "shape")
    sp_pr = _sp_pr(shape)
    preset = "rect"
    adjustments: List[float] = []
    if sp_pr is not None:
        geom = sp_pr.find(qn("a:prstGeom"))
        if geom is not None:
            preset = geom.get("prst") or "rect"
            adjustments = _adjustments(geom, preset)
        elif sp_pr.find(qn("a:custGeom")) is not None:
            preset = "custom"

    result: Dict = {
        "preset": preset,
        "adjustments": adjustments,
        "rotation_deg": 0.0,
        "left_pt": None, "top_pt": None,
        "width_pt": None, "height_pt": None,
    }
    xfrm = _inherited_xfrm(shape, slide)
    if xfrm is None:
        return result

    offset = xfrm.find(qn("a:off"))
    extent = xfrm.find(qn("a:ext"))
    if offset is not None:
        result["left_pt"] = _int_attr(offset, "x") / EMU_PER_POINT
        result["top_pt"] = _int_attr(offset, "y") / EMU_PER_POINT
    if extent is not None:
        result["width_pt"] = _int_attr(extent, "cx") / EMU_PER_POINT
        result["height_pt"] = _int_attr(extent, "cy") / EMU_PER_POINT
    raw_rotation = xfrm.get("rot")
    if raw_rotation is not None:
        result["rotation_deg"] = int(raw_rotation) / ROTATION_DENOMINATOR
    return result
