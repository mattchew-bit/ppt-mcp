"""Multi-deck house-style profile builder (Step 3 of the plan).

``create_house_profile(paths)`` runs the Step 2 effective-style resolver
across a corpus of reference decks (via ``utils.resolve_analysis`` --
resolution is never reimplemented here) and aggregates the resolved
values into the compact, prescriptive ``house-profile/1`` rules JSON
documented in ``utils.profile_schema``: modal values + tolerances, not
raw dumps.

Pipeline:

    decks -> resolved slide facts -> aggregate sections
        typography / paragraph / palette / shape_defaults   (modal rules)
        grid                                (utils.profile_grid inference)
        archetypes                          (utils.profile_archetypes)
        images / distributions              (descriptive stats for lint)
    -> enforce the 8KB size budget -> profile dict

Grid semantics decision: footer-zone shapes (page numbers, source
notes -- anything whose top edge sits below ``FOOTER_FRACTION`` of the
slide height) are EXCLUDED from column-grid inference. Their x
positions repeat on nearly every slide, so feeding them to the
clustering mints phantom column edges out of what is really footer
furniture. The footer conventions themselves are not lost: the
typography ``footer`` role captures their text style and the archetype
zones capture their geometry, so no separate ``grid.aux`` key is
carried.

Strictly read-only on the input decks.
"""

import json
from collections import Counter
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .profile_archetypes import FOOTER_FRACTION, learn_archetypes
from .profile_grid import POINTS_PER_INCH, infer_grid
from .profile_schema import (
    SCHEMA_VERSION,
    color,
    dim,
    enforce_size_budget,
    rect_in,
    round2,
    token,
)
from .resolve_analysis import build_resolved_analysis
from .style_roles import ph_type_label_name, placeholder_role

#: Internal resolver cap -- the profile builder needs the FULL corpus.
_RESOLVE_MAX_BYTES = 64_000_000

#: Top-N caps (size budget levers; enforce_size_budget names them).
TOP_N_PALETTE_USAGE = 8
TOP_N_DISTRIBUTION = 8
TOP_N_IMAGE_BUCKETS = 3

#: Bullet levels the paragraph section learns (house scale: l1..l3).
MAX_BULLET_LEVELS = 3

#: Canonical clrScheme slot order for the palette section.
_SCHEME_ORDER = ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                 "accent3", "accent4", "accent5", "accent6",
                 "hlink", "folHlink")

#: Image-zone quantization step, in inches (buckets nudge noise away).
_ZONE_STEP_IN = 0.5

#: Fallback body size when a corpus carries no body placeholder at all.
_FALLBACK_BODY_SIZE_PT = 18.0


def _modal(counter: Counter):
    """Most frequent value; ties break lexically for determinism."""
    if not counter:
        return None
    ranked = sorted(counter.items(), key=lambda kv: (-kv[1], str(kv[0])))
    return ranked[0][0]


# ---------------------------------------------------------------------------
# Corpus collection (resolver-backed slide facts)
# ---------------------------------------------------------------------------

def _is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except NotImplementedError:  # exotic shapes report no type
        return False


def _deck_facts(path: Path) -> Tuple[List[Dict], List[Dict], Tuple]:
    """One deck -> (slide facts, theme summaries, slide size in inches)."""
    prs = Presentation(str(path))
    report = build_resolved_analysis(
        prs, detail="full", max_bytes=_RESOLVE_MAX_BYTES)
    if report.get("truncated"):
        raise ValueError(
            f"resolved analysis of {path.name} was truncated -- "
            "the deck is too large to profile in one pass"
        )
    width_pt, height_pt = prs.slide_width.pt, prs.slide_height.pt
    slides = []
    for slide_record, slide in zip(report["slides"], prs.slides):
        shapes = []
        for record, shape in zip(slide_record["shapes"], slide.shapes):
            merged = dict(record)  # never mutate the resolver's output
            merged["is_picture"] = _is_picture(shape)
            shapes.append(merged)
        slides.append({
            "deck": path.name,
            "slide_number": slide_record["slide_number"],
            "width_pt": width_pt,
            "height_pt": height_pt,
            "shapes": shapes,
        })
    size = (round2(width_pt / POINTS_PER_INCH),
            round2(height_pt / POINTS_PER_INCH))
    return slides, report.get("themes", []), size


def collect_corpus_facts(paths: Sequence[str]) -> Dict:
    """All decks -> merged slide facts + modal theme + slide size.

    Returns ``{"slides": [...], "theme": <summary>, "slide_size_in":
    (w, h)}`` where each slide-facts record carries the resolver's
    per-shape output plus ``is_picture``. Public so the archetype /
    grid layers (and their tests) can run on the same resolved facts
    the profile builder uses.
    """
    slides: List[Dict] = []
    theme_counter: Counter = Counter()
    themes_by_key: Dict[str, Dict] = {}
    sizes = set()
    for raw_path in paths:
        path = Path(raw_path)
        if not path.is_file():
            raise FileNotFoundError(f"reference deck not found: {raw_path}")
        deck_slides, themes, size = _deck_facts(path)
        slides.extend(deck_slides)
        sizes.add(size)
        for theme in themes:
            key = json.dumps(theme, sort_keys=True)
            theme_counter[key] += 1
            themes_by_key[key] = theme
    if not slides:
        raise ValueError("reference decks contain no slides to profile")
    if len(sizes) > 1:
        raise ValueError(
            f"reference decks disagree on slide size: {sorted(sizes)} -- "
            "profile one aspect ratio at a time"
        )
    modal_theme = themes_by_key[_modal(theme_counter)]
    return {"slides": slides, "theme": modal_theme,
            "slide_size_in": next(iter(sizes))}


# ---------------------------------------------------------------------------
# Role classification (typography scopes)
# ---------------------------------------------------------------------------

def _shape_role(shape: Dict, slide_height_pt: float) -> Optional[str]:
    """title / footer / body scope for typography aggregation.

    Placeholder types map through the shared exact-match table
    (``utils.style_roles``), so the learn side counts exactly the
    shapes the apply side will later write (SUBTITLE / VERTICAL_TITLE
    are neither). Non-placeholder shapes parked in the footer zone
    also count as footer -- the zone is a house convention regardless
    of the shape type carrying it.
    """
    role = placeholder_role(ph_type_label_name(shape.get("ph_type")))
    if role == "title":
        return "title"
    top = shape["geometry"].get("top_pt")
    if top is not None and top >= FOOTER_FRACTION * slide_height_pt:
        return "footer"
    return role


def _iter_role_runs(slides: Sequence[Dict]):
    """Yield ``(role, indent_level, run_font)`` for every resolved run."""
    for slide in slides:
        for shape in slide["shapes"]:
            role = _shape_role(shape, slide["height_pt"])
            for paragraph in shape.get("paragraphs", []):
                for run in paragraph.get("runs", []):
                    yield role, paragraph["indent_level"], run["font"]


def _typography_section(slides: Sequence[Dict]) -> Tuple[Dict, float]:
    """Modal font/size/bold/color per role; returns (section, body size)."""
    counters = {
        role: {"font": Counter(), "size": Counter(),
               "bold": Counter(), "color": Counter()}
        for role in ("title", "body", "footer")
    }
    all_sizes: Counter = Counter()
    for role, level, font in _iter_role_runs(slides):
        if font.get("size_pt") is not None:
            all_sizes[font["size_pt"]] += 1
        if role is None or (role == "body" and level != 1):
            continue
        bucket = counters[role]
        if font.get("name"):
            bucket["font"][font["name"]] += 1
        if font.get("size_pt") is not None:
            bucket["size"][font["size_pt"]] += 1
        bucket["bold"][bool(font.get("bold"))] += 1
        if font.get("color_hex"):
            bucket["color"][font["color_hex"]] += 1

    section: Dict[str, Dict] = {}
    for role, bucket in counters.items():
        spec: Dict[str, Dict] = {}
        if _modal(bucket["font"]) is not None:
            spec["font"] = token(_modal(bucket["font"]))
        if _modal(bucket["size"]) is not None:
            spec["size"] = dim(_modal(bucket["size"]), "pt")
        if bucket["bold"]:
            spec["bold"] = token(bool(_modal(bucket["bold"])))
        if _modal(bucket["color"]) is not None:
            spec["color"] = color(_modal(bucket["color"]))
        if spec:
            section[role] = spec
    body_size = _modal(counters["body"]["size"])
    if body_size is None:
        body_size = _modal(all_sizes) or _FALLBACK_BODY_SIZE_PT
    return section, float(body_size)


# ---------------------------------------------------------------------------
# Paragraph rules (body placeholder cascade, per level)
# ---------------------------------------------------------------------------

def _iter_body_paragraphs(slides: Sequence[Dict]):
    for slide in slides:
        for shape in slide["shapes"]:
            if _shape_role(shape, slide["height_pt"]) != "body":
                continue
            for paragraph in shape.get("paragraphs", []):
                if paragraph.get("runs"):
                    yield paragraph


def _bullet_rule(bullets: Dict[str, Counter]) -> Optional[Dict]:
    """Modal per-level bullet rule: {char, color?, size_pct?} ONLY.

    The pinned house-profile/1 contract carries no bullet ``font`` key
    (glyphs stay in the text font per the synthetic universal-font
    corpus); widening the contract requires a schema revision, not an
    extra key.
    """
    char = _modal(bullets["char"])
    if char is None:
        return None
    rule: Dict[str, Dict] = {"char": token(char)}
    size_pct = _modal(bullets["size_pct"])
    if size_pct is not None:
        rule["size_pct"] = dim(size_pct, None)
    bullet_color = _modal(bullets["color"])
    if bullet_color is not None:
        rule["color"] = color(bullet_color)
    return rule


def _count_bullet(paragraph: Dict, bullets: Dict[str, Counter]) -> None:
    bullet = paragraph.get("bullet") or {}
    if bullet.get("type") != "char" or not bullet.get("char"):
        return
    bullets["char"][bullet["char"]] += 1
    if not bullet.get("size_follows_text") and bullet.get("size_pct"):
        bullets["size_pct"][round2(bullet["size_pct"] * 100)] += 1
    if not bullet.get("color_follows_text") and bullet.get("color_hex"):
        bullets["color"][bullet["color_hex"]] += 1


def _paragraph_section(slides: Sequence[Dict]) -> Dict:
    """Modal spacing + per-level bullet rules from body paragraphs."""
    levels: Dict[int, Dict[str, Counter]] = {}
    for paragraph in _iter_body_paragraphs(slides):
        level = paragraph["indent_level"]
        if not 1 <= level <= MAX_BULLET_LEVELS:
            continue
        bucket = levels.setdefault(level, {
            "before": Counter(), "after": Counter(), "line": Counter(),
            "char": Counter(), "size_pct": Counter(), "color": Counter(),
        })
        if "points" in paragraph["space_before"]:
            bucket["before"][paragraph["space_before"]["points"]] += 1
        if "points" in paragraph["space_after"]:
            bucket["after"][paragraph["space_after"]["points"]] += 1
        if "multiple" in paragraph["line_spacing"]:
            bucket["line"][paragraph["line_spacing"]["multiple"]] += 1
        _count_bullet(paragraph, bucket)

    section: Dict = {}
    level_one = levels.get(1)
    if level_one is not None:
        if _modal(level_one["before"]) is not None:
            section["space_before"] = dim(_modal(level_one["before"]), "pt")
        if _modal(level_one["after"]) is not None:
            section["space_after"] = dim(_modal(level_one["after"]), "pt")
        if _modal(level_one["line"]) is not None:
            section["line_spacing"] = dim(_modal(level_one["line"]), None)
    rules = {
        f"l{level}": _bullet_rule(levels[level])
        for level in sorted(levels)
        if _bullet_rule(levels[level]) is not None
    }
    if rules:
        section["bullets"] = rules
    return section


# ---------------------------------------------------------------------------
# Palette (theme scheme + observed usage shares)
# ---------------------------------------------------------------------------

def _usage_counts(slides: Sequence[Dict]) -> Dict[str, Counter]:
    """Per-color total counts and per-color context counts."""
    totals: Counter = Counter()
    contexts: Dict[str, Counter] = {}

    def count(hex_value: Optional[str], context: str) -> None:
        if not hex_value:
            return
        totals[hex_value] += 1
        contexts.setdefault(hex_value, Counter())[context] += 1

    for slide in slides:
        for shape in slide["shapes"]:
            fill = shape.get("fill") or {}
            if fill.get("type") == "solid":
                count(fill.get("color_hex"), "fill")
            line = shape.get("line") or {}
            if line.get("visible"):
                count(line.get("color_hex"), "line")
            for paragraph in shape.get("paragraphs", []):
                for run in paragraph.get("runs", []):
                    count(run["font"].get("color_hex"), "text")
    return {"totals": totals, "contexts": contexts}


def _palette_section(slides: Sequence[Dict], theme: Dict) -> Dict:
    scheme_hex = theme.get("color_scheme", {})
    scheme = {
        name: color(scheme_hex[name])
        for name in _SCHEME_ORDER if name in scheme_hex
    }
    usage = _usage_counts(slides)
    total = sum(usage["totals"].values())
    ranked = sorted(usage["totals"].items(),
                    key=lambda kv: (-kv[1], kv[0]))[:TOP_N_PALETTE_USAGE]
    entries = [
        {
            "color": f"#{hex_value.upper()}",
            "share": round2(count / total),
            "role": _modal(usage["contexts"][hex_value]),
        }
        for hex_value, count in ranked
    ]
    return {"scheme": scheme, "usage": entries}


# ---------------------------------------------------------------------------
# Shape defaults (panel border / corner / fill conventions)
# ---------------------------------------------------------------------------

def _shape_defaults_section(slides: Sequence[Dict]) -> Dict:
    borders: Counter = Counter()
    corners: Counter = Counter()
    fills: Counter = Counter()
    for slide in slides:
        for shape in slide["shapes"]:
            if shape.get("is_placeholder") or shape.get("is_picture"):
                continue
            line = shape.get("line") or {}
            if not line.get("visible"):
                continue
            borders[(line.get("weight_pt"), line.get("color_hex"),
                     line.get("dash"))] += 1
            if (shape["geometry"].get("preset") == "roundRect"
                    and shape.get("adjustments")):
                corners[shape["adjustments"][0]] += 1
            fill = shape.get("fill") or {}
            if fill.get("type") == "solid" and fill.get("color_hex"):
                fills[fill["color_hex"]] += 1

    section: Dict = {}
    border = _modal(borders)
    if border is not None:
        weight_pt, color_hex, dash = border
        rule: Dict = {}
        if weight_pt is not None:
            rule["weight"] = dim(weight_pt, "pt")
        if color_hex is not None:
            rule["color"] = color(color_hex)
        if dash is not None:
            rule["dash"] = token(dash)
        section["border"] = rule
    if _modal(corners) is not None:
        section["corner_radius"] = dim(_modal(corners), None)
    if _modal(fills) is not None:
        section["fill"] = color(_modal(fills))
    return section


# ---------------------------------------------------------------------------
# Grid + archetypes (delegated inference, schema-shaped here)
# ---------------------------------------------------------------------------

def _footerless_slides(slides: Sequence[Dict]) -> List[Dict]:
    """New slide records without footer-zone shapes (grid input only).

    See the module docstring: footer furniture repeats at x positions
    that are footer conventions, not content columns, and would mint
    phantom column edges. Shapes with no resolvable top keep their say
    (they cannot be proven to be footer furniture). Input records are
    never mutated.
    """
    filtered = []
    for slide in slides:
        cutoff = FOOTER_FRACTION * slide["height_pt"]
        shapes = [
            shape for shape in slide["shapes"]
            if shape["geometry"].get("top_pt") is None
            or shape["geometry"]["top_pt"] < cutoff
        ]
        filtered.append({**slide, "shapes": shapes})
    return filtered


def _grid_section(slides: Sequence[Dict]) -> Dict:
    grid = infer_grid(_footerless_slides(slides))
    return {
        "edges": {
            family: [round2(edge) for edge in edges]
            for family, edges in grid["edges"].items()
        },
        "unit": "in",
        "tolerance": dim(grid["tolerance_in"], "in"),
    }


def _archetypes_section(slides: Sequence[Dict],
                        body_size_pt: float) -> Dict:
    learned = learn_archetypes(slides, body_size_pt)
    section: Dict = {}
    for name, spec in learned.items():
        entry: Dict = {}
        for box in ("title_band", "body_region"):
            rect = spec[box]
            if rect is not None:
                entry[box] = rect_in(rect["x"], rect["y"],
                                     rect["w"], rect["h"])
        entry["count"] = spec["count"]
        section[name] = entry
    return section


# ---------------------------------------------------------------------------
# Images (count / size distribution / placement zones)
# ---------------------------------------------------------------------------

def _image_rects_in(slides: Sequence[Dict]) -> List[Tuple]:
    rects = []
    for slide in slides:
        for shape in slide["shapes"]:
            if not shape.get("is_picture"):
                continue
            geometry = shape["geometry"]
            if geometry.get("left_pt") is None:
                continue
            rects.append(tuple(
                geometry[key] / POINTS_PER_INCH
                for key in ("left_pt", "top_pt", "width_pt", "height_pt")
            ))
    return rects


def _zone_buckets(rects: Sequence[Tuple]) -> List[Dict]:
    """Quantize image rects into placement zones with usage shares."""
    buckets: Dict[Tuple, List[Tuple]] = {}
    for rect in rects:
        key = tuple(round(value / _ZONE_STEP_IN) for value in rect)
        buckets.setdefault(key, []).append(rect)
    zones = []
    for members in buckets.values():
        mean = [sum(values) / len(members) for values in zip(*members)]
        zones.append({"rect": mean, "share": len(members) / len(rects)})
    zones.sort(key=lambda zone: (-zone["share"], zone["rect"]))
    return zones[:TOP_N_IMAGE_BUCKETS]


def _images_section(slides: Sequence[Dict]) -> Dict:
    rects = _image_rects_in(slides)
    per_slide = [
        sum(1 for shape in slide["shapes"] if shape.get("is_picture"))
        for slide in slides
    ]
    section: Dict = {
        "count_per_slide": {
            "mean": round2(sum(per_slide) / len(per_slide)),
            "max": max(per_slide),
        },
    }
    if not rects:
        section["size_distribution"] = []
        section["zones"] = []
        return section
    sizes: Counter = Counter(
        (round2(rect[2]), round2(rect[3])) for rect in rects)
    ranked = sorted(sizes.items(),
                    key=lambda kv: (-kv[1], kv[0]))[:TOP_N_IMAGE_BUCKETS]
    section["size_distribution"] = [
        {"width": dim(width, "in"), "height": dim(height, "in"),
         "share": round2(count / len(rects))}
        for (width, height), count in ranked
    ]
    section["zones"] = [
        {**rect_in(*zone["rect"]), "share": round2(zone["share"])}
        for zone in _zone_buckets(rects)
    ]
    return section


# ---------------------------------------------------------------------------
# Distributions (house scales for Step 5 lint messages)
# ---------------------------------------------------------------------------

def _distribution(counter: Counter, by_value: bool = True) -> Dict:
    """Counter -> ``{"values": [...], "shares": [...]}`` (capped)."""
    total = sum(counter.values())
    if total == 0:
        return {"values": [], "shares": []}
    top = sorted(counter.items(),
                 key=lambda kv: (-kv[1], str(kv[0])))[:TOP_N_DISTRIBUTION]
    if by_value:
        top = sorted(top, key=lambda kv: kv[0])
    return {
        "values": [round2(v) if isinstance(v, float) else v
                   for v, _ in top],
        "shares": [round2(count / total) for _, count in top],
    }


def _distributions_section(slides: Sequence[Dict]) -> Dict:
    sizes: Counter = Counter()
    for _, _, font in _iter_role_runs(slides):
        if font.get("size_pt") is not None:
            sizes[float(font["size_pt"])] += 1
    space_after: Counter = Counter()
    for slide in slides:
        for shape in slide["shapes"]:
            for paragraph in shape.get("paragraphs", []):
                if (paragraph.get("runs")
                        and "points" in paragraph["space_after"]):
                    space_after[float(
                        paragraph["space_after"]["points"])] += 1
    palette: Counter = Counter()
    totals = _usage_counts(slides)["totals"]
    for hex_value, count in totals.items():
        palette[f"#{hex_value.upper()}"] += count
    return {
        "font_sizes_pt": _distribution(sizes),
        "space_after_pt": _distribution(space_after),
        "palette_shares": _distribution(palette, by_value=False),
    }


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def create_house_profile(paths: Sequence[str],
                         name: str = "house") -> Dict:
    """Build a ``house-profile/1`` rules dict from reference decks.

    ``paths`` are .pptx files (5-10 decks recommended). The result is
    deterministic (same input -> byte-identical serialized profile),
    JSON-safe, and guaranteed to fit the 8KB budget (``ValueError``
    otherwise, naming the sections to trim). See ``utils.profile_schema``
    for the pinned schema and the consumption split.
    """
    if not paths or not all(isinstance(p, str) and p for p in paths):
        raise ValueError(
            "paths must be a non-empty list of .pptx file paths"
        )
    if not name or not isinstance(name, str):
        raise ValueError(f"profile name must be a non-empty string: {name!r}")
    corpus = collect_corpus_facts(paths)
    slides = corpus["slides"]
    typography, body_size_pt = _typography_section(slides)
    width_in, height_in = corpus["slide_size_in"]
    profile: Dict = {
        "schema_version": SCHEMA_VERSION,
        "name": name,
        "source_decks": [Path(p).name for p in paths],
        "slide_size": {"width": dim(width_in, "in"),
                       "height": dim(height_in, "in")},
        "typography": typography,
        "paragraph": _paragraph_section(slides),
        "palette": _palette_section(slides, corpus["theme"]),
        "shape_defaults": _shape_defaults_section(slides),
        "grid": _grid_section(slides),
        "archetypes": _archetypes_section(slides, body_size_pt),
        "images": _images_section(slides),
        "distributions": _distributions_section(slides),
    }
    enforce_size_budget(profile)
    return profile
