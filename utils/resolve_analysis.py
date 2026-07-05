"""Resolved-mode deck analysis for ``analyze_presentation_style``.

Walks a presentation with the Stage 2 effective-style resolver
(``utils.resolve_utils``) and produces a JSON-safe report of what
PowerPoint actually displays -- inheritance-resolved fonts, colors,
paragraph metrics, bullets and shape line/fill -- with the output
controls the plan mandates (a resolved full dump of a 30-slide deck
would blow the model's context):

    * ``slide_range``    -- 1-based selection string, e.g. ``"1-3,5"``
    * ``detail``         -- ``"summary"`` (deck-level rollup, default)
      or ``"full"`` (per-run dump)
    * ``groups``         -- property-group filter, any subset of
      ``{"fonts", "colors", "paragraphs", "shapes", "theme"}``
    * ``max_bytes``      -- hard response cap; when the serialized
      report would exceed it, trailing slides are dropped and an
      explicit ``{"truncated": true, "hint": ...}`` marker is added

Read-only: the presentation is never mutated.
"""

import json
from typing import Any, Dict, List, Optional, Sequence, Set

from pptx.oxml.ns import qn

from .resolve_theme import ThemeContext
from .resolve_utils import (
    TextContext,
    resolve_paragraph,
    resolve_run_font,
    resolve_shape,
)

#: Default hard cap for the serialized resolved report (~40KB).
MAX_RESPONSE_BYTES = 40_000

VALID_DETAILS = ("summary", "full")
VALID_GROUPS = ("fonts", "colors", "paragraphs", "shapes", "theme")


# ---------------------------------------------------------------------------
# Input validation / parsing
# ---------------------------------------------------------------------------

def parse_slide_range(spec: Optional[str], slide_count: int) -> List[int]:
    """``"1-3,5"`` -> 0-based slide indices (validated, de-duplicated).

    ``None``/empty selects every slide. Indices are 1-based in the spec,
    matching how users talk about slides.
    """
    if not spec:
        return list(range(slide_count))
    indices: List[int] = []
    seen: Set[int] = set()
    for chunk in str(spec).split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            start_raw, _, end_raw = chunk.partition("-")
            start, end = int(start_raw), int(end_raw)
        else:
            start = end = int(chunk)
        if start < 1 or end > slide_count or start > end:
            raise ValueError(
                f"slide_range chunk {chunk!r} outside 1..{slide_count}"
            )
        for number in range(start, end + 1):
            if number - 1 not in seen:
                seen.add(number - 1)
                indices.append(number - 1)
    if not indices:
        raise ValueError(f"slide_range {spec!r} selects no slides")
    return indices


def _validate_controls(detail: str, groups) -> Set[str]:
    if detail not in VALID_DETAILS:
        raise ValueError(
            f"detail must be one of {VALID_DETAILS}, got {detail!r}"
        )
    if groups is None:
        return set(VALID_GROUPS)
    selected = {str(g) for g in groups}
    unknown = selected - set(VALID_GROUPS)
    if unknown:
        raise ValueError(
            f"unknown property group(s) {sorted(unknown)}; "
            f"valid: {VALID_GROUPS}"
        )
    return selected


# ---------------------------------------------------------------------------
# Per-slide resolution
# ---------------------------------------------------------------------------

def _theme_for_master(master,
                      cache: Dict[str, ThemeContext]) -> ThemeContext:
    """ThemeContext for a master, parsing each theme part exactly once.

    ``cache`` is keyed by the master's part name; every slide governed
    by the same master shares one ``ThemeContext`` (multi-master decks
    get one per master).
    """
    key = str(master.part.partname)
    context = cache.get(key)
    if context is None:
        context = ThemeContext.for_master(master)
        cache[key] = context
    return context


def _paragraph_record(context: TextContext, groups: Set[str]) -> Dict:
    record: Dict[str, Any] = {}
    if "paragraphs" in groups:
        record.update(resolve_paragraph(context))
    if "fonts" in groups or "colors" in groups:
        runs = []
        for run in context.paragraph.runs:
            font = resolve_run_font(TextContext(
                context.slide, context.shape, context.paragraph, run,
                context.theme,
            ))
            runs.append({"text": run.text, "font": font})
        record["runs"] = runs
    return record


def _has_text_body(shape) -> bool:
    """True when the shape's XML actually carries a ``p:txBody``.

    Deliberately NOT ``shape.has_text_frame`` + ``shape.text_frame``:
    ``has_text_frame`` is True for every ``p:sp`` even when its XML
    legally omits ``p:txBody``, and ``shape.text_frame`` then calls
    ``get_or_add_txBody()``, CREATING a txBody with a fabricated empty
    paragraph -- a mutation the resolver's read-only guarantee forbids.
    """
    return shape._element.find(qn("p:txBody")) is not None


def _shape_record(slide, shape, theme: ThemeContext,
                  groups: Set[str]) -> Dict:
    record: Dict[str, Any] = {
        "name": shape.name,
        "is_placeholder": shape.is_placeholder,
        "ph_type": (str(shape.placeholder_format.type)
                    if shape.is_placeholder else None),
    }
    if "shapes" in groups:
        record.update(resolve_shape(shape, slide, theme))
    if groups & {"fonts", "colors", "paragraphs"} and _has_text_body(shape):
        record["paragraphs"] = [
            _paragraph_record(
                TextContext(slide, shape, paragraph, None, theme), groups)
            for paragraph in shape.text_frame.paragraphs
        ]
    return record


def _slide_record(prs, slide_index: int, groups: Set[str],
                  theme_cache: Dict[str, ThemeContext]) -> Dict:
    slide = prs.slides[slide_index]
    theme = _theme_for_master(slide.slide_layout.slide_master, theme_cache)
    return {
        "slide_number": slide_index + 1,
        "layout": slide.slide_layout.name,
        "shapes": [
            _shape_record(slide, shape, theme, groups)
            for shape in slide.shapes
        ],
    }


def _theme_summaries(prs,
                     theme_cache: Dict[str, ThemeContext]) -> List[Dict]:
    summaries = []
    for master in prs.slide_masters:
        context = _theme_for_master(master, theme_cache)
        summaries.append({
            "theme_name": context.theme_name,
            "color_scheme": context.color_scheme,
            "fonts": {
                "major_latin": context.major_latin_font,
                "minor_latin": context.minor_latin_font,
            },
        })
    return summaries


# ---------------------------------------------------------------------------
# Summary rollup
# ---------------------------------------------------------------------------

def _count(counter: Dict, key) -> None:
    if key is not None:
        counter[key] = counter.get(key, 0) + 1


#: Counter names accumulated by ``_count_slide_styles``.
_COUNTER_NAMES = ("fonts", "sizes", "colors", "bullets", "spacing",
                  "line_colors", "fill_colors")


def _count_paragraph_styles(paragraph: Dict,
                            counters: Dict[str, Dict]) -> None:
    """Accumulate one paragraph record into the summary counters."""
    bullet = paragraph.get("bullet")
    if bullet is not None and bullet["type"] != "none":
        _count(counters["bullets"],
               f"L{paragraph['indent_level']}:"
               f"{bullet.get('char') or bullet['type']}")
    if "space_after" in paragraph:
        _count(counters["spacing"], json.dumps({
            "before": paragraph["space_before"],
            "after": paragraph["space_after"],
            "line": paragraph["line_spacing"],
        }, sort_keys=True))
    for run in paragraph.get("runs", []):
        _count(counters["fonts"], run["font"]["name"])
        _count(counters["sizes"], run["font"]["size_pt"])
        _count(counters["colors"], run["font"]["color_hex"])


def _count_slide_styles(slide_records: List[Dict]) -> Dict[str, Dict]:
    """Frequency counters over every shape/paragraph/run record."""
    counters: Dict[str, Dict] = {name: {} for name in _COUNTER_NAMES}
    for slide_record in slide_records:
        for shape in slide_record["shapes"]:
            fill = shape.get("fill") or {}
            line = shape.get("line") or {}
            _count(counters["fill_colors"], fill.get("color_hex"))
            _count(counters["line_colors"], line.get("color_hex"))
            for paragraph in shape.get("paragraphs", []):
                _count_paragraph_styles(paragraph, counters)
    return counters


def _top(counter: Dict, limit: int = 12) -> List:
    """The ``limit`` most frequent counter entries, ranked."""
    ranked = sorted(counter.items(), key=lambda kv: -kv[1])[:limit]
    return [{"value": key, "count": count} for key, count in ranked]


def _summarize_slides(slide_records: List[Dict],
                      groups: Set[str]) -> Dict:
    counters = _count_slide_styles(slide_records)
    summary: Dict[str, Any] = {}
    if "fonts" in groups:
        summary["fonts"] = _top(counters["fonts"])
        summary["sizes_pt"] = _top(counters["sizes"])
    if "colors" in groups:
        summary["text_colors"] = _top(counters["colors"])
        summary["fill_colors"] = _top(counters["fill_colors"])
        summary["line_colors"] = _top(counters["line_colors"])
    if "paragraphs" in groups:
        summary["bullets"] = _top(counters["bullets"])
        summary["paragraph_spacing"] = [
            {"value": json.loads(key), "count": count}
            for key, count in sorted(counters["spacing"].items(),
                                     key=lambda kv: -kv[1])[:8]
        ]
    return summary


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _fit_to_cap(report: Dict, slide_records: List[Dict],
                max_bytes: int) -> None:
    """Drop trailing slides until the report fits; mark truncation."""
    kept = list(slide_records)
    report["slides"] = kept
    while kept and len(json.dumps(report)) > max_bytes:
        kept.pop()
        report["slides"] = kept
        report["truncated"] = True
        report["hint"] = (
            f"response capped at {max_bytes} bytes: {len(kept)} of "
            f"{len(slide_records)} selected slides returned -- re-run "
            "with slide_range for the rest, or detail='summary', or a "
            "groups filter"
        )
    if not kept and slide_records:
        report["hint"] = (
            f"response capped at {max_bytes} bytes: even one slide "
            "exceeds the cap -- use detail='summary' or a groups filter"
        )


def build_resolved_analysis(
    prs,
    slide_range: Optional[str] = None,
    detail: str = "summary",
    groups: Optional[Sequence[str]] = None,
    max_bytes: int = MAX_RESPONSE_BYTES,
) -> Dict:
    """Inheritance-resolved style report for an open presentation.

    ``prs`` is a python-pptx ``Presentation``. See the module docstring
    for the output controls. Always JSON-serializable and no larger
    than ``max_bytes`` when serialized (explicit truncation markers
    otherwise).
    """
    if prs is None:
        raise ValueError("prs must not be None")
    if not isinstance(max_bytes, int) or max_bytes <= 0:
        raise ValueError(f"max_bytes must be a positive int: {max_bytes!r}")
    selected_groups = _validate_controls(detail, groups)
    indices = parse_slide_range(slide_range, len(prs.slides))

    theme_cache: Dict[str, ThemeContext] = {}
    slide_records = [
        _slide_record(prs, index, selected_groups, theme_cache)
        for index in indices
    ]
    report: Dict[str, Any] = {
        "resolved": True,
        "detail": detail,
        "slide_numbers": [index + 1 for index in indices],
        "groups": sorted(selected_groups),
        "truncated": False,
    }
    if "theme" in selected_groups:
        report["themes"] = _theme_summaries(prs, theme_cache)

    if detail == "summary":
        report["summary"] = _summarize_slides(slide_records,
                                              selected_groups)
        if len(json.dumps(report)) > max_bytes:
            report.pop("summary", None)
            report["truncated"] = True
            report["hint"] = (
                f"summary exceeded the {max_bytes}-byte cap; narrow with "
                "slide_range or a groups filter"
            )
    else:
        _fit_to_cap(report, slide_records, max_bytes)
    return report
